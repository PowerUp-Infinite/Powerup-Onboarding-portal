"""
m2_engine.py - PowerUp Infinite Strategy Deck Generator.

Port of M2/app.py - logic is byte-for-byte identical to the local automation.
The only differences are:
  * Data loaded from Google Sheets via the `sheets` module
  * Base deck / Risk-Reward deck / Categorization / rating images downloaded
    from Google Drive (cached in a tempdir for the session)
  * generate_deck() returns (BytesIO, filename) instead of saving to disk
"""

import os
import sys
import re

if sys.platform == 'win32':
    try:
        sys.stdout.reconfigure(encoding='utf-8', errors='replace')
        sys.stderr.reconfigure(encoding='utf-8', errors='replace')
    except Exception:
        pass

from copy import deepcopy
from io import BytesIO

import pandas as pd
import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.dates
import matplotlib.ticker

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from lxml import etree

import sheets

# ----------------------------------------------------------------
# CONFIGURATION - all assets downloaded from Google Drive
# ----------------------------------------------------------------

import tempfile
from config import (
    M2_TEMPLATE_ID, M2_RISK_REWARD_TEMPLATE_ID,
    M2_CATEGORIZATION_FILE_ID,
    M2_IMG_INFORM_ID, M2_IMG_ONTRACK_ID,
    M2_IMG_OFFTRACK_ID, M2_IMG_OUTOFFORM_ID,
)

# Temporary directory for cached Drive downloads (persists across calls within one session)
_CACHE_DIR = tempfile.mkdtemp(prefix='m2_assets_')

def _download_drive_file(file_id, filename, export_mime=None):
    """Download a file from Google Drive into _CACHE_DIR. Returns local path.
    Uses export_mime for Google Workspace files (Slides->PPTX), direct download for binary files.
    Caches: only downloads once per session."""
    local_path = os.path.join(_CACHE_DIR, filename)
    if os.path.exists(local_path):
        return local_path

    from google_auth import get_drive_service
    from googleapiclient.http import MediaIoBaseDownload
    drive = get_drive_service()

    buf = BytesIO()
    if export_mime:
        request = drive.files().export_media(fileId=file_id, mimeType=export_mime)
    else:
        request = drive.files().get_media(fileId=file_id)
    downloader = MediaIoBaseDownload(buf, request)
    done = False
    while not done:
        _, done = downloader.next_chunk()
    buf.seek(0)

    with open(local_path, 'wb') as f:
        f.write(buf.read())
    return local_path


def _get_base_deck_path():
    """Download M2 Base Deck from Drive (Google Slides -> PPTX)."""
    return _download_drive_file(
        M2_TEMPLATE_ID, 'Base_Deck.pptx',
        export_mime='application/vnd.openxmlformats-officedocument.presentationml.presentation',
    )

def _get_rr_deck_path():
    """Download Risk Reward Slides from Drive (Google Slides -> PPTX)."""
    return _download_drive_file(
        M2_RISK_REWARD_TEMPLATE_ID, 'Risk_Reward.pptx',
        export_mime='application/vnd.openxmlformats-officedocument.presentationml.presentation',
    )

def _get_categorization_path():
    """Download Scheme Categorization Excel from Drive."""
    return _download_drive_file(
        M2_CATEGORIZATION_FILE_ID, 'Scheme_Category_Catgorization.xlsx',
    )

def _get_rating_image(key):
    """Download a rating image from Drive. key: IN_FORM / ON_TRACK / OUT_OF_FORM / OFF_TRACK."""
    file_map = {
        'IN_FORM':     (M2_IMG_INFORM_ID,    'InForm.png'),
        'ON_TRACK':    (M2_IMG_ONTRACK_ID,   'OnTrack.png'),
        'OUT_OF_FORM': (M2_IMG_OUTOFFORM_ID, 'OutofForm.png'),
        'OFF_TRACK':   (M2_IMG_OFFTRACK_ID,  'OffTrack.png'),
    }
    fid, fname = file_map[key]
    return _download_drive_file(fid, fname)

# Set of valid POWERRATING values — used as a membership check before
# attempting to fetch the rating image from Drive via _get_rating_image().
# Kept as a set (not a dict like in M2-automation/app.py) because the portal
# resolves image paths lazily through Drive instead of reading from disk.
RATING_IMAGES = {'IN_FORM', 'ON_TRACK', 'OUT_OF_FORM', 'OFF_TRACK'}

# Pie-chart colours — taken directly from base deck legend shapes
CHART_COLORS = {
    '1) Aggressive':   '#2E8AE5',
    '2) Balanced':     '#4E9EED',
    '3) Conservative': '#6DB0F2',
    'Hybrid':          '#FFE2BF',
    'Debt Like':       '#EBF2F2',  # light mint/seafoam — from PowerUp_Base_Deck legend
    'Gold & Silver':   '#F7CB88',  # golden amber      — from PowerUp_Base_Deck legend
    'Global':          '#FFC7B4',  # salmon/peach-pink  — from PowerUp_Base_Deck legend
    'Solution':        '#CABAF3',  # lavender           — from PowerUp_Base_Deck legend
}
CHART_LABELS = {
    '1) Aggressive':   'Aggressive',
    '2) Balanced':     'Balanced',
    '3) Conservative': 'Conservative',
    'Hybrid':          'Hybrid',
    'Gold & Silver':   'Gold & Silver',
    'Debt Like':       'Debt',
    'Solution':        'Solution',
    'Global':          'Global',
}

# Risk profile scale (index 0..4)
RISK_SCALE = ['Very Conservative', 'Conservative', 'Balanced', 'Aggressive', 'Very Aggressive']

HORIZON_DISPLAY = {
    # Most-specific keys first — dict iteration order is insertion order.
    'more than 8':'8+ Years','more than 7':'8+ Years',
    '5-8':'5-8 Years','medium to long':'5-8 Years',
    '3-5':'3-5 Years','3–5':'3-5 Years','medium-term':'3-5 Years',
    'less than 3':'Less than 3 Years','short':'Less than 3 Years',
    'long-term':'8+ Years','long':'8+ Years',
}

# XML namespaces
NS_R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
NS_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'

# Scheme-slide table X positions (EMU)
X_NAME       = 572501
X_RATING     = 2707092
X_VALUE      = 4271150
X_XIRR       = 5813826
X_MISSED     = 7635275
X_RATING_IMG = 3334942
RATING_IMG_SZ = 255600

# Base-deck slide-4 pie chart image dimensions (shape ;194)
PIE_IMG_LEFT   = 2667449
PIE_IMG_TOP    = 1043625
PIE_IMG_WIDTH  = 5159751
PIE_IMG_HEIGHT = 3190450

# ──────────────────────────────────────────────────────────────
# DATA LOADING
# ──────────────────────────────────────────────────────────────

def load_data(pf_id: str | None = None):
    """Load all data from Google Sheets (+ categorization file from Drive).
    If pf_id is given, the large per-client tabs (Lines, Invested_Value_Line)
    are filtered to just that client — drastically reduces peak memory on
    Streamlit Cloud (free tier has only ~1 GB, full Lines can be >1 M rows)."""
    data = {}
    data['pf_level']       = sheets.read_pf_level()
    data['riskgroup']      = sheets.read_riskgroup_level()
    data['scheme']         = sheets.read_scheme_level()
    data['categorization'] = pd.read_excel(_get_categorization_path())
    data['questionnaire']  = sheets.read_questionnaire()
    data['lines']          = sheets.read_lines(pf_id=pf_id) if pf_id else sheets.read_lines()
    data['results']        = sheets.read_results()
    data['invested']       = sheets.read_invested_value_line(pf_id=pf_id) if pf_id else sheets.read_invested_value_line()
    # is_demat: SOA/Demat split per PF_ID. Cloud flow reads the 'Is_demat' tab
    # on the MAIN spreadsheet. Desktop flow overrides this with the uploaded
    # Excel's 'Is demat' sheet (handled in desktop/workers/m2_worker.py).
    try:
        data['is_demat']   = sheets.read_is_demat()
    except Exception as e:
        print(f"  WARN: could not load Is_demat tab ({e}) — falling back to "
              f"PF_level SOA%/Demat% columns.")
        data['is_demat']   = pd.DataFrame()
    # Convert numeric columns that come back as strings from Sheets API
    _text_cols = {
        'PF_ID', 'ISIN', 'NAME', 'FUND_NAME', 'FUND_STANDARD_NAME',
        'FUND_LEGAL_NAME', 'TYPE', 'POWERRATING', 'DISTRIBUTION_STATUS',
        'RISK_GROUP_L0', 'UPDATED_SUBCATEGORY', 'UPDATED_BROAD_CATEGORY_GROUP',
        'BROAD_CATEGORY_GROUP', 'DERIVED_CATEGORY', 'Purchase Mode',
        'BM', 'DIR_ISIN', 'ALT_ISIN_J', 'DATE',
    }
    for key in ('pf_level', 'riskgroup', 'scheme', 'results', 'lines', 'invested'):
        df = data[key]
        for col in df.columns:
            if col in _text_cols:
                continue
            df[col] = pd.to_numeric(df[col], errors='coerce')
        data[key] = df
    for name, df in data.items():
        print(f"  Loaded {name}: {len(df)} rows, cols={list(df.columns)[:6]}...")
    return data


# ──────────────────────────────────────────────────────────────
# FORMATTING HELPERS
# ──────────────────────────────────────────────────────────────

def fmt_inr_rupee(value, prefix='₹'):
    """Rs.94.9L / Rs.1.2Cr / Rs.25K"""
    if pd.isna(value) or value == 0:
        return f'{prefix}0'
    av = abs(value); s = '' if value >= 0 else '-'
    if av >= 1e7:
        cr = av / 1e7
        return f'{s}{prefix}{cr:.1f}Cr' if cr < 10 else f'{s}{prefix}{cr:.0f}Cr'
    if av >= 1e5:
        l = av / 1e5
        return f'{s}{prefix}{l:.1f}L' if l < 10 else f'{s}{prefix}{l:.0f}L'
    if av >= 1e3:
        k = av / 1e3
        return f'{s}{prefix}{k:.1f}K' if k < 10 else f'{s}{prefix}{k:.0f}K'
    return f'{s}{prefix}{av:.0f}'

# Keep backward compat alias
fmt_inr = fmt_inr_rupee

def _fmt_inr_2dp(value, prefix=''):
    """Like fmt_inr_rupee but always 2 decimal places — used in appendix tables."""
    if pd.isna(value) or value == 0:
        return f'{prefix}0'
    av = abs(value); s = '' if value >= 0 else '-'
    if av >= 1e7:
        return f'{s}{prefix}{av/1e7:.2f}Cr'
    if av >= 1e5:
        return f'{s}{prefix}{av/1e5:.2f}L'
    if av >= 1e3:
        return f'{s}{prefix}{av/1e3:.2f}K'
    return f'{s}{prefix}{av:.0f}'

def fmt_inr_display(value):
    """INR 50L / INR 1 Cr — for slide display"""
    if pd.isna(value) or value == 0:
        return None
    av = abs(value)
    if av >= 1e7:
        cr = av / 1e7
        return f'INR {cr:.0f} Cr' if cr == int(cr) else f'INR {cr:.1f} Cr'
    if av >= 1e5:
        l = av / 1e5
        return f'INR {l:.0f}L' if l == int(l) else f'INR {l:.1f}L'
    if av >= 1e3:
        k = av / 1e3
        return f'INR {k:.0f}K' if k == int(k) else f'INR {k:.1f}K'
    return f'INR {av:.0f}'

def _safe_inr(val):
    """Format as INR if numeric; return string as-is; NaN -> '-'."""
    if val is None: return '-'
    try:
        if pd.isna(val): return '-'
    except (TypeError, ValueError):
        pass
    # Handle int/float/numpy numeric types uniformly
    try:
        fv = float(val)
        return fmt_inr_display(fv) or '-'
    except (ValueError, TypeError):
        s = str(val).strip()
        return s if s else '-'

def _safe_pct(val):
    """Format as percentage. NaN -> '-'. Handles decimals (0.05 -> 5%)."""
    if val is None: return '-'
    try:
        if pd.isna(val): return '-'
    except (TypeError, ValueError):
        pass
    s = str(val).strip()
    if '%' in s: return s
    try:
        fv = float(s)
        # Values stored as decimal fraction (0.05 = 5%) → multiply by 100
        if 0 < abs(fv) < 1:
            fv *= 100
        return f'{fv:.0f}%'
    except Exception:
        return s if s else '-'

_EXCEL_EPOCH = pd.Timestamp('1899-12-30')

def _parse_dates(series):
    """Parse a date series that may contain date strings OR Excel serial numbers."""
    try:
        # Try normal string parsing first
        return pd.to_datetime(series, dayfirst=True)
    except Exception:
        pass
    # Fall back: treat numeric values as Excel serial-day offsets from 1899-12-30
    def _conv(v):
        try:
            n = float(v)
            if 20000 < n < 100000:   # plausible Excel serial (1954–2173)
                return _EXCEL_EPOCH + pd.Timedelta(days=int(n))
        except (ValueError, TypeError):
            pass
        # ISO format strings ("2024-09-01" or "2024-09-01 00:00:00") MUST be
        # parsed as YYYY-MM-DD, not dayfirst=True (which would give Jan 9).
        # Try ISO first, fall back to dayfirst for "01/09/2024" style strings.
        s = str(v).strip()
        if not s:
            return pd.NaT
        # ISO-ish: starts with 4 digits then '-'
        if len(s) >= 10 and s[:4].isdigit() and s[4] == '-':
            try:
                return pd.to_datetime(s, format='ISO8601', errors='raise')
            except Exception:
                try:
                    return pd.to_datetime(s, errors='raise')
                except Exception:
                    return pd.NaT
        try:
            return pd.to_datetime(v, dayfirst=True)
        except Exception:
            return pd.NaT
    return series.apply(_conv)


def _safe_str(val):
    """Return string value; NaN/None -> '-'. Whole-number floats (e.g. 2028.0) -> '2028'."""
    if val is None: return '-'
    if isinstance(val, float):
        if pd.isna(val): return '-'
        if val == int(val):
            return str(int(val))
    s = str(val).strip()
    return s if s else '-'

def fmt_scheme_val(cv, pf_pct):
    v = _fmt_inr_2dp(cv)
    return f'{v} ({pf_pct * 100:.1f}%)'

def fmt_xirr_pair(x, bx):
    def f(v): return '-' if pd.isna(v) else f'{v * 100:.1f}%'
    return f'{f(x)} | {f(bx)}'

def fmt_missed(mg):
    if pd.isna(mg) or mg == 0: return '-'
    return _fmt_inr_2dp(mg)

def _match(text, mapping):
    if pd.isna(text): return None
    t = str(text).lower()
    for k, v in mapping.items():
        if k in t: return v
    return None

# ──────────────────────────────────────────────────────────────
# RISK PROFILE CALCULATION
# ──────────────────────────────────────────────────────────────

def calc_risk_profile(q):
    """
    4-step risk logic:
    1. Base from Portfolio Preference return % (15% VeryAgg, 12% Agg, 9% Bal, 6% Con)
    2. Horizon adjustment: short/medium/medium-to-long -> -1; long-term -> 0
    3. Fall Reaction: invest more -> +1; stay invested -> 0; exit* -> -1
    4. Liability management: 'Yes - comfortably' -> 0; 'Just about' -> -1; other struggling -> -1
    """
    # Step 1: Base index from Portfolio Preference
    pref = str(q.get('Portfolio Preference', '')).lower()
    if '15%' in pref:
        idx = 4   # Very Aggressive
    elif '12%' in pref:
        idx = 3   # Aggressive
    elif '9%' in pref:
        idx = 2   # Balanced
    elif '6%' in pref:
        idx = 1   # Conservative
    else:
        idx = 2   # Balanced default
    base = RISK_SCALE[idx]

    # Step 2: Horizon adjustment
    # Only truly long-term (8+ years) gets h_adj=0 (no change).
    # "Medium to long-term" must NOT match as long-term — check 'medium' not in horizon first.
    horizon = str(q.get('Investment Horizon', '')).lower()
    long_kws = ['more than 7', 'more than 8', 'long-term', 'long term', '8+']
    is_long = any(k in horizon for k in long_kws) and 'medium' not in horizon
    h_adj = 0 if is_long else -1
    idx = max(0, min(4, idx + h_adj))

    # Step 3: Fall Reaction adjustment
    # "Invest more" and "Stay invested" both leave risk unchanged; only an
    # exit response (partial or full) pulls the profile down by one.
    fall = str(q.get('Fall Reaction', '')).lower()
    if 'invest more' in fall or 'stay invested' in fall or 'stay' in fall:
        f_adj = 0
    else:  # exit all / exit partially / anything else
        f_adj = -1
    idx = max(0, min(4, idx + f_adj))

    # Step 4: Liability management adjustment
    #
    # The three real answers currently in the questionnaire form are:
    #   "Yes - comfortably"  → has liabilities, managing fine        → 0
    #   "No I don't foresee" → no liabilities at all                 → 0
    #   "Just about"         → has liabilities and barely managing   → -1
    #
    # Only "Just about" (or equivalent stress wording) downgrades. Matching
    # by 'yes' / 'comfort' alone misclassified "No I don't foresee" as stress
    # and incorrectly subtracted 1 from the risk index.
    liab = str(q.get('Liability Followup Answer', '')).lower()
    stress_kws = ('just about', 'barely', 'struggl', 'difficult', 'tight', 'hardly')
    if any(k in liab for k in stress_kws):
        l_adj = -1
    else:
        # empty / "yes comfortably" / "no I don't foresee" / etc. → no change
        l_adj = 0
    idx = max(0, min(4, idx + l_adj))

    profile = RISK_SCALE[idx]
    print(f"  Risk: base={base} h_adj={h_adj} f_adj={f_adj} l_adj={l_adj} -> {profile}")
    return profile

def get_horizon(text):
    return _match(text, HORIZON_DISPLAY) or str(text)

def parse_goals(text):
    if pd.isna(text): return []
    return [g.strip() for g in str(text).split(',') if g.strip()]

# ──────────────────────────────────────────────────────────────
# SLIDE MANIPULATION (low-level)
# ──────────────────────────────────────────────────────────────

def delete_slide(prs, idx):
    sldIdLst = prs.slides._sldIdLst
    rId = sldIdLst[idx].get(f'{{{NS_R}}}id')
    prs.part.drop_rel(rId)
    sldIdLst.remove(sldIdLst[idx])

def clone_slide(prs, tpl_idx):
    """Deep-copy a slide (with images) and append at end. Returns new Slide."""
    tpl = prs.slides[tpl_idx]
    new = prs.slides.add_slide(tpl.slide_layout)

    img_map = {}
    for key, rel in tpl.part.rels.items():
        if 'image' in rel.reltype:
            img_map[key] = new.part.relate_to(rel.target_part, rel.reltype)

    sp_tree = new.shapes._spTree
    for ch in list(sp_tree):
        tag = etree.QName(ch.tag).localname
        if tag in ('sp', 'grpSp', 'pic', 'graphicFrame', 'cxnSp'):
            sp_tree.remove(ch)

    for ch in tpl.shapes._spTree:
        tag = etree.QName(ch.tag).localname
        if tag in ('sp', 'grpSp', 'pic', 'graphicFrame', 'cxnSp'):
            el = deepcopy(ch)
            for blip in el.iter(f'{{{NS_A}}}blip'):
                old = blip.get(f'{{{NS_R}}}embed')
                if old and old in img_map:
                    blip.set(f'{{{NS_R}}}embed', img_map[old])
            sp_tree.append(el)

    bg = tpl._element.find(f'{{{NS_P}}}bg')
    if bg is not None:
        old_bg = new._element.find(f'{{{NS_P}}}bg')
        if old_bg is not None:
            new._element.remove(old_bg)
        new._element.insert(0, deepcopy(bg))

    return new

def move_slide(prs, src, dst):
    lst = prs.slides._sldIdLst
    entries = list(lst)
    el = entries[src]
    lst.remove(el)
    entries = list(lst)
    if dst >= len(entries):
        lst.append(el)
    else:
        entries[dst].addprevious(el)

def replace_text(shape, new_text):
    """Set shape text preserving first-run formatting."""
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        if para.runs:
            para.runs[0].text = str(new_text)
            for r in para.runs[1:]:
                r.text = ''
            return
    shape.text_frame.paragraphs[0].text = str(new_text)

def set_table_cell(cell, text):
    para = cell.text_frame.paragraphs[0]
    if para.runs:
        para.runs[0].text = str(text)
        for r in para.runs[1:]:
            r.text = ''
    else:
        cell.text = str(text)


def _set_cell_transparent(cell):
    """Remove any background fill from a table cell (make it transparent)."""
    tc = cell._tc
    tcPr = tc.find(f'{{{NS_A}}}tcPr')
    if tcPr is None:
        tcPr = etree.SubElement(tc, f'{{{NS_A}}}tcPr')
    for child in list(tcPr):
        lname = etree.QName(child.tag).localname
        if lname in ('noFill', 'solidFill', 'gradFill', 'blipFill', 'pattFill', 'grpFill'):
            tcPr.remove(child)
    etree.SubElement(tcPr, f'{{{NS_A}}}noFill')

def remove_shape(slide, shape):
    el = shape._element
    el.getparent().remove(el)

# ──────────────────────────────────────────────────────────────
# SLIDE 1 — Title
# ──────────────────────────────────────────────────────────────

def do_slide1(prs, full_name):
    slide = prs.slides[0]
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            txt = ''.join(r.text for r in para.runs)
            if 'with' in txt.lower() and len(txt) < 80:
                para.runs[0].text = f'with {full_name}'
                para.runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                for r in para.runs[1:]:
                    r.text = ''
                print(f"  Slide 1: title -> 'with {full_name}'")
                return
    print("  Slide 1: WARNING - name placeholder not found")

# ──────────────────────────────────────────────────────────────
# SLIDE 2 — Welcome / Agenda
# ──────────────────────────────────────────────────────────────

def do_slide2(prs, first_name):
    slide = prs.slides[1]
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        paras = shape.text_frame.paragraphs
        if len(paras) >= 2 and 'Welcome' in paras[0].text:
            p1 = paras[1]
            if p1.runs:
                p1.runs[0].text = first_name
                for r in p1.runs[1:]:
                    r.text = ''
            else:
                p1.text = first_name
            print(f"  Slide 2: welcome name -> '{first_name}'")
            return
    print("  Slide 2: WARNING - name placeholder not found")

# ──────────────────────────────────────────────────────────────
# SLIDE 3 — You at a Glance
# ──────────────────────────────────────────────────────────────

def do_slide3(prs, q_row, risk_profile, pf_row=None):
    slide   = prs.slides[2]
    goals   = parse_goals(q_row.get('Goals', ''))
    horizon = get_horizon(q_row.get('Investment Horizon', ''))
    # Age preference order:
    #   1. PF_level.Age (authoritative — set by the data team)
    #   2. Questionnaire 'Age' field (fallback for old workflow)
    age = ''
    if pf_row is not None and not (hasattr(pf_row, 'empty') and pf_row.empty):
        v = pf_row.get('Age')
        if v is not None and not (isinstance(v, float) and pd.isna(v)) and str(v).strip():
            age = str(v).strip()
    if not age:
        age = q_row.get('Age', '') or ''

    lump_val  = q_row.get('Lumpsum Amount (with Infinite)', 0)
    sip_val   = q_row.get('Monthly SIP Amount (with Infinite)', 0)
    stepup_raw = q_row.get('Ret: YoY Investment Increase %', 0)
    lump_str  = fmt_inr_display(lump_val)  if not pd.isna(lump_val)   and lump_val   else 'INR 0'
    sip_str   = fmt_inr_display(sip_val)   if not pd.isna(sip_val)    and sip_val    else 'INR 0'
    try:
        step_up = float(stepup_raw) if not pd.isna(stepup_raw) else 0.0
    except (TypeError, ValueError):
        # Could be "10%" string
        step_up = float(str(stepup_raw).rstrip('%')) if str(stepup_raw).rstrip('%').replace('.','').isdigit() else 0.0
    # Data stored as fraction (0.2 = 20%) — normalize to percent for display
    if 0 < step_up < 1:
        step_up *= 100
    has_stepup = step_up > 0

    shapes_to_remove = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        if shape.name == 'Google Shape;129;p18':
            _set_goals_text(shape, goals)

        elif shape.name == 'Google Shape;125;p18':
            replace_text(shape, horizon)
            print(f"  Slide 3: horizon -> '{horizon}'")

        elif shape.name == 'Google Shape;123;p18':
            _set_investment_text(shape, lump_str, sip_str, has_stepup)

        elif shape.name == 'Google Shape;127;p18':
            replace_text(shape, risk_profile)
            print(f"  Slide 3: risk profile -> '{risk_profile}'")

        elif shape.name == 'Google Shape;130;p18':
            replace_text(shape, f'{risk_profile} Investor')

        elif shape.name in ('Google Shape;132;p18', 'Google Shape;133;p18'):
            # Shape 132 = age, shape 133 = SIP step-up. (Older templates
            # had both as ;132;p18 — the text-content check handles both.)
            cur = shape.text_frame.text.strip()
            if 'SIP' in cur or 'Step' in cur or 'step' in cur:
                if has_stepup:
                    replace_text(shape, f'SIP Step-Up every year: {step_up:.0f}%')
                    print(f"  Slide 3: SIP step-up -> {step_up:.0f}%")
                else:
                    shapes_to_remove.append(shape)
                    print("  Slide 3: SIP step-up = 0, removing text box")
            else:
                if age:
                    replace_text(shape, f'Current Age: {age}')
                    print(f"  Slide 3: age -> '{age}'")

    # Remove shapes after iteration to avoid modifying the collection mid-loop
    sp_tree = slide.shapes._spTree
    for shape in shapes_to_remove:
        sp_tree.remove(shape._element)


def _set_goals_text(shape, goals):
    """Primary goal in large text (template size); secondary goals smaller below."""
    if not goals:
        replace_text(shape, 'Wealth Creation')
        print("  Slide 3: goals -> 'Wealth Creation'")
        return

    primary   = goals[0]
    secondary = goals[1:] if len(goals) > 1 else []

    tf   = shape.text_frame
    txBody = tf._txBody
    a_ns   = f'{{{NS_A}}}'

    # Set first paragraph (primary goal) — preserve existing run formatting
    p0 = tf.paragraphs[0]
    if p0.runs:
        p0.runs[0].text = primary
        for r in p0.runs[1:]:
            r.text = ''
    else:
        p0.text = primary

    # Remove any extra paragraphs from a previous run
    existing = txBody.findall(f'{a_ns}p')
    for ep in existing[1:]:
        txBody.remove(ep)

    # Add secondary goals as a smaller second paragraph
    if secondary:
        p1  = tf.add_paragraph()
        run = p1.add_run()
        run.text           = ', '.join(secondary)
        run.font.name      = 'IBM Plex Sans'
        run.font.size      = Pt(11)
        run.font.color.rgb = RGBColor(0x23, 0x23, 0x23)

    print(f"  Slide 3: goals -> primary='{primary}'"
          + (f" secondary={secondary}" if secondary else ""))


def _set_investment_text(shape, lump_str, sip_str, has_stepup):
    """
    Reproduce the 7-run formatting of base deck shape ;123:
      [lump/dark]  [ ]  [with/gray]  [ ]  [sip/dark]  [ ]  [monthly SIP/gray]
    """
    sip_label = 'monthly SIP*' if has_stepup else 'monthly SIP'
    tf   = shape.text_frame
    para = tf.paragraphs[0]
    runs = para.runs

    if len(runs) >= 7:
        # Template runs intact — update text; colours already set in template
        runs[0].text = lump_str
        runs[1].text = ' '
        runs[2].text = 'with'
        runs[3].text = ' '
        runs[4].text = sip_str
        runs[5].text = ' '
        runs[6].text = sip_label
        # Ensure correct run colours
        for ri, clr in [(0, RGBColor(0x23,0x23,0x23)),
                        (2, RGBColor(0x8E,0x93,0x93)),
                        (4, RGBColor(0x23,0x23,0x23)),
                        (6, RGBColor(0x8E,0x93,0x93))]:
            runs[ri].font.color.rgb = clr
    else:
        # Rebuild from scratch
        for r in runs:
            r.text = ''
        if runs:
            runs[0].text = lump_str
            runs[0].font.size      = Pt(18)
            runs[0].font.color.rgb = RGBColor(0x23, 0x23, 0x23)
        parts = [
            (' ',       Pt(18), RGBColor(0x23,0x23,0x23)),
            ('with',    Pt(15), RGBColor(0x8E,0x93,0x93)),
            (' ',       Pt(18), RGBColor(0x23,0x23,0x23)),
            (sip_str,   Pt(18), RGBColor(0x23,0x23,0x23)),
            (' ',       Pt(18), RGBColor(0x23,0x23,0x23)),
            (sip_label, Pt(15), RGBColor(0x8E,0x93,0x93)),
        ]
        for text, size, color in parts:
            r = para.add_run()
            r.text           = text
            r.font.name      = 'IBM Plex Sans'
            r.font.size      = size
            r.font.color.rgb = color

    print(f"  Slide 3: investment -> '{lump_str} with {sip_str} {sip_label}'")

# ──────────────────────────────────────────────────────────────
# SLIDE 4 — Portfolio Snapshot  (metrics + pie chart)
# ──────────────────────────────────────────────────────────────

def _portfolio_risk(sm):
    """Derive portfolio risk from Small + Mid allocation percentage."""
    if sm < 15:  return 'Very Conservative'
    if sm < 20:  return 'Conservative'
    if sm < 40:  return 'Balanced'
    if sm < 45:  return 'Aggressive'
    return 'Very Aggressive'


def _shape_max_font_emu(shape):
    """Return the largest font size (EMU) across a shape's runs, or 0."""
    if not shape.has_text_frame:
        return 0
    mx = 0
    for para in shape.text_frame.paragraphs:
        for r in para.runs:
            if r.font.size is None:
                continue
            try:
                sz = int(r.font.size)
                if sz > mx:
                    mx = sz
            except Exception:
                pass
    return mx


def do_slide4(prs, pf, rg_agg, risk_profile, is_demat_df=None):
    """
    Slide 4 uses text-based shape lookup (not hardcoded Google Shape IDs)
    because every edit of the Drive template shifts shape IDs by +1. Labels
    are identified by their keyword text + small font (sz < 200000 EMU);
    values are the large-font shape (sz >= 200000) sitting directly below
    each label in the same column.
    """
    slide = prs.slides[3]
    cv   = pf['PF_CURRENT_VALUE']
    iv   = pf['INVESTED_VALUE']
    xirr = pf['PF_XIRR']
    bxir = pf['BM_XIRR']
    pg   = pf['PF_GAINS']
    bg   = pf.get('BM_CURRENT_VALUE', iv) - iv
    sm   = (pf.get('SMALL', 0) + pf.get('MID', 0)) * pf.get('EQUITY', 0) * 100
    pf_risk = _portfolio_risk(sm)
    matches = pf_risk == risk_profile

    # Classify top-level text shapes as labels (small font) vs values (large font)
    text_shapes = [s for s in slide.shapes
                   if s.has_text_frame and s.text_frame.text.strip()]
    labels = [s for s in text_shapes if 90000 <= _shape_max_font_emu(s) < 200000]
    values = [s for s in text_shapes if _shape_max_font_emu(s) >= 200000]

    def _value_below(label_shape):
        """Find the large-font shape directly below a label (same column)."""
        best, best_dy = None, 10**12
        lL, lT, lW = label_shape.left, label_shape.top, label_shape.width
        for v in values:
            if v.top <= lT:
                continue
            # horizontal overlap with label
            if v.left > lL + lW or v.left + v.width < lL:
                continue
            dy = v.top - lT
            if dy < best_dy and dy < 600000:   # within ~0.6M EMU ≈ 0.6 in
                best_dy, best = dy, v
        return best

    def _find_label(*keywords):
        for lbl in labels:
            t = lbl.text_frame.text.strip().lower()
            if all(k in t for k in keywords):
                return lbl
        return None

    for name, kws, new_val in [
        ('Current Value',   ('current', 'value'),   fmt_inr_rupee(cv)),
        ('Invested Amount', ('invested',),          fmt_inr_rupee(iv)),
        ('Current XIRR',    ('current', 'xirr'),    f'{xirr * 100:.1f}%'),
        ('Benchmark XIRR',  ('benchmark', 'xirr'),  f'{bxir * 100:.1f}%'),
    ]:
        lbl = _find_label(*kws)
        if lbl is None:
            print(f"  Slide 4: WARN label {name!r} not found"); continue
        val_sh = _value_below(lbl)
        if val_sh is None:
            print(f"  Slide 4: WARN value shape for {name!r} not found"); continue
        replace_text(val_sh, new_val)
        print(f"  Slide 4: {name} -> {new_val}")

    # Portfolio risk — find "Portfolio Risk" label, write risk below it
    risk_lbl = (_find_label('portfolio', 'risk')
                or _find_label('risk', 'profile')
                or _find_label('risk'))
    if risk_lbl is not None:
        risk_val_sh = _value_below(risk_lbl)
        if risk_val_sh is not None:
            replace_text(risk_val_sh, pf_risk)
            print(f"  Slide 4: portfolio risk (S+M={sm:.0f}%) -> '{pf_risk}'")

    # SOA % and Demat % — these labels live in a smaller-font row at the
    # bottom-left of the slide. The value shapes underneath them are also
    # smaller font (≈101350 EMU) so _value_below() (which only looks at
    # values >= 200000 EMU) won't find them. Search by exact label text +
    # nearest "%" shape directly below in the same column.
    def _pct_value_below(label_shape):
        """Find any text shape containing '%' below `label_shape` in the
        same column. Used for SOA / Demat values which are smaller-font
        than the main metric values."""
        if label_shape is None:
            return None
        lL, lT, lW = label_shape.left, label_shape.top, label_shape.width
        best, best_dy = None, 1 << 30
        for s in text_shapes:
            if s is label_shape:
                continue
            if not s.text_frame.text.strip().endswith('%'):
                continue
            if s.top <= lT:
                continue
            if s.left > lL + lW or s.left + s.width < lL:
                continue
            dy = s.top - lT
            if dy < best_dy and dy < 400000:
                best_dy, best = dy, s
        return best

    def _fmt_pct_for_slide(val):
        """0.96 -> '96%'  ;  96 -> '96%'  ;  1 -> '1%'  ;  None/NaN -> '-'.

        Disambiguation: STRICTLY less than 1 means fraction (0.85 = 85%).
        1 or above means it's already a percentage integer (1 = 1%, 99 =
        99%). Treating 1 as 100% would mis-classify rows like Demat=1."""
        if val is None or (isinstance(val, float) and pd.isna(val)):
            return '-'
        try:
            f = float(val)
        except (TypeError, ValueError):
            s = str(val).strip()
            return s if s else '-'
        if 0 < f < 1:
            f *= 100
        return f'{int(round(f))}%'

    # SOA% / Demat% — prefer the dedicated 'Is demat' sheet (rows split by
    # IS_DEMAT bool, PCT_OF_USER sums to 1). Fall back to legacy columns on
    # PF_level for back-compat with older Excels / the cloud flow.
    soa_val, demat_val = None, None
    if is_demat_df is not None and not is_demat_df.empty \
            and 'IS_DEMAT' in is_demat_df.columns \
            and 'PCT_OF_USER' in is_demat_df.columns:
        def _is_true(v):
            if isinstance(v, bool):
                return v
            if isinstance(v, (int, float)) and not pd.isna(v):
                return float(v) != 0.0
            return str(v).strip().lower() in ('true', 'yes', '1', 't', 'y')
        # Each row specifies ONE side; the other is the complement.
        #   IS_DEMAT=True, PCT=0.05  → Demat=5%,  SOA=95%
        #   IS_DEMAT=False, PCT=0.05 → SOA=5%,    Demat=95%
        # If both True and False rows are present, use each side explicitly.
        mask_true = is_demat_df['IS_DEMAT'].apply(_is_true)
        true_sum = float(
            pd.to_numeric(is_demat_df.loc[mask_true, 'PCT_OF_USER'],
                          errors='coerce').fillna(0).sum()
        )
        false_sum = float(
            pd.to_numeric(is_demat_df.loc[~mask_true, 'PCT_OF_USER'],
                          errors='coerce').fillna(0).sum()
        )
        if true_sum > 0 and false_sum > 0:
            demat_frac, soa_frac = true_sum, false_sum
        elif true_sum > 0:
            demat_frac = true_sum
            soa_frac   = max(0.0, 1.0 - demat_frac)
        elif false_sum > 0:
            soa_frac   = false_sum
            demat_frac = max(0.0, 1.0 - soa_frac)
        else:
            demat_frac = soa_frac = 0.0

        # Convert fractions (0..1) to integer percent here so the downstream
        # _fmt_pct_for_slide doesn't mis-classify a true 100% (fraction=1.0)
        # as "1%". The "fraction vs already-percent" boundary check inside
        # _fmt_pct_for_slide is only correct for ambiguous PF_level columns.
        demat_val = int(round(demat_frac * 100))
        soa_val   = int(round(soa_frac   * 100))
        print(f"  Slide 4: is_demat sheet -> SOA={soa_val}%, "
              f"Demat={demat_val}%")
    else:
        soa_val   = pf.get('SOA%') if hasattr(pf, 'get') else None
        demat_val = pf.get('Demat%') if hasattr(pf, 'get') else None

    for label_text, raw, kind in (('SOA', soa_val, 'soa'),
                                  ('Demat', demat_val, 'demat')):
        if raw is None or (isinstance(raw, float) and pd.isna(raw)):
            continue   # column missing/empty — leave template default
        # Iterate ALL text shapes (not just `labels`) since "SOA"/"Demat"
        # labels are sz=133350 and live in our `labels` set anyway.
        lbl = None
        for s in text_shapes:
            if s.text_frame.text.strip() == label_text:
                lbl = s; break
        if lbl is None:
            print(f"  Slide 4: WARN {label_text} label not found"); continue
        val_sh = _pct_value_below(lbl)
        if val_sh is None:
            print(f"  Slide 4: WARN {label_text} value shape not found"); continue
        new_text = _fmt_pct_for_slide(raw)
        replace_text(val_sh, new_text)
        print(f"  Slide 4: {label_text} -> {new_text}")

    # Gains + S+M + match/no-match — match by existing text content
    match_text   = "Matches your risk profile"
    nomatch_text = "Doesn't match your risk profile"
    match_color   = RGBColor(0x2A, 0x9C, 0x4A)   # green
    nomatch_color = RGBColor(0xCC, 0x00, 0x00)   # red
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        txt = shape.text_frame.text.strip()
        tl  = txt.lower()
        if 'portfolio gains' in tl:
            replace_text(shape, f'Portfolio gains: {fmt_inr_rupee(pg)}')
            print(f"  Slide 4: PF gains -> {fmt_inr_rupee(pg)}")
        elif 'benchmark gains' in tl:
            replace_text(shape, f'Benchmark gains: {fmt_inr_rupee(bg)}')
            print(f"  Slide 4: BM gains -> {fmt_inr_rupee(bg)}")
        elif 'small' in tl and 'mid' in tl and 'allocation' in tl:
            replace_text(shape, f'Small + Mid Allocation: {sm:.0f}%')
            print(f"  Slide 4: S+M -> {sm:.0f}%")
        elif 'match' in tl and 'risk profile' in tl:
            new_text  = match_text if matches else nomatch_text
            new_color = match_color if matches else nomatch_color
            replace_text(shape, new_text)
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        run.font.color.rgb = new_color
                    except Exception:
                        pass
            print(f"  Slide 4: match/no-match -> '{new_text}' "
                  f"(pf_risk='{pf_risk}', q_risk='{risk_profile}')")

    _make_pie(slide, rg_agg)


def _make_pie(slide, rg):
    """
    Regenerate the portfolio allocation donut.
    - Generates ONLY the donut ring (no text, no legend) as transparent PNG
    - Same dimensions as original image ;194 so shape ;155 centre-text aligns
    - Inserts behind shape ;155 in z-order
    - Updates existing legend group percentages (;166 / ;185)
    """
    if rg.empty:
        print("  Slide 4: no riskgroup data — skipping chart"); return

    parts = []
    for _, row in rg.iterrows():
        g = row['RISK_GROUP_L0']
        p = row['% of PF']
        if pd.isna(p) or p <= 0: continue
        # Skip rows with missing or unknown risk group (NaN / blank). These
        # are dust rows in Riskgroup_level that would otherwise crash the
        # legend writer with "Argument must be bytes or unicode, got float".
        if g is None or (isinstance(g, float) and pd.isna(g)) or str(g).strip() == '':
            continue
        parts.append((CHART_LABELS.get(g, g), p * 100, CHART_COLORS.get(g, '#808080')))
    if not parts:
        print("  Slide 4: empty riskgroup — skipping chart"); return

    eq_order    = ['Aggressive', 'Balanced', 'Conservative']
    # Legend order from reference deck: Equity → Hybrid → Debt Like → Gold & Silver → Global → Solution
    other_order = ['Hybrid', 'Debt', 'Gold & Silver', 'Global', 'Solution']

    def sk(x):
        if x[0] in eq_order:    return (0, eq_order.index(x[0]))
        if x[0] in other_order: return (1, other_order.index(x[0]))
        return (2, 0)

    parts.sort(key=sk)
    sizes  = [p[1] for p in parts]
    colors = [p[2] for p in parts]
    eq_total = sum(s for l, s, _ in parts if l in eq_order)

    # Enforce hard 1% minimum so any non-zero slice is always visible
    sizes = [max(s, 1.0) if s > 0 else s for s in sizes]

    # ── Generate donut ring: same aspect ratio as original image ─────────────
    # Original image is 5159751 × 3190450 EMU  (ratio ≈ 1.617)
    # Shape ;155 centre is at exactly 50% w / 50% h of the image
    # So we size our figure to match, put the donut at centre, rest transparent
    ratio = PIE_IMG_WIDTH / PIE_IMG_HEIGHT   # 1.6172
    fig_h = 5.0
    fig_w = fig_h * ratio
    ypad  = 1.15
    xpad  = ypad * ratio                     # keeps equal aspect visually

    fig, ax = plt.subplots(figsize=(fig_w, fig_h))
    fig.patch.set_alpha(0)
    fig.subplots_adjust(0, 0, 1, 1)          # remove all white margins

    ax.set_aspect('equal')
    ax.pie(
        sizes, colors=colors, radius=1.0, startangle=90, counterclock=False,
        wedgeprops=dict(width=0.35, edgecolor='white', linewidth=2.5),
    )
    ax.axis('off')
    ax.set_xlim(-xpad, xpad)
    ax.set_ylim(-ypad, ypad)

    buf = BytesIO()
    fig.savefig(buf, format='png', dpi=200,
                transparent=True, pad_inches=0, bbox_inches=None)
    plt.close(fig)
    buf.seek(0)

    # ── Locate the centre-text shape by content ("Portfolio / Allocation") ──
    shape_155_el = None
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        tl = shape.text_frame.text.strip().lower()
        if 'allocation' in tl and 'portfolio' in tl and len(tl) < 40:
            shape_155_el = shape._element
            break

    # ── Remove old pie-chart image (large, height > 1 M EMU) ─────────────────
    for shape in list(slide.shapes):
        if (shape.shape_type == 13
                and shape.width  > 3_000_000
                and shape.height > 1_000_000
                and shape.top    < 2_000_000):
            try:
                _ = shape.image
                remove_shape(slide, shape)
                print("  Slide 4: removed old pie image")
                break
            except Exception:
                pass

    # ── Add new donut ring at the exact original image position ──────────────
    new_pic = slide.shapes.add_picture(
        buf,
        Emu(PIE_IMG_LEFT), Emu(PIE_IMG_TOP),
        Emu(PIE_IMG_WIDTH), Emu(PIE_IMG_HEIGHT),
    )

    # Move new picture behind centre-text shape so the centre text stays on top
    if shape_155_el is not None:
        sp_tree = slide.shapes._spTree
        new_el  = new_pic._element
        sp_tree.remove(new_el)
        shape_155_el.addprevious(new_el)
        print("  Slide 4: donut placed behind centre text")

    # ── Update existing legend group text ─────────────────────────────────────
    _update_legend_groups(slide, parts, eq_total)
    print(f"  Slide 4: pie chart done ({len(parts)} segments, eq={eq_total:.0f}%)")


def _update_legend_groups(slide, parts, eq_total):
    """
    Update percentage (and optionally label) text in the existing base-deck
    legend groups. Labels are matched by TEXT CONTENT (Conservative / Balanced
    / Aggressive / Equity / Hybrid) rather than hardcoded Google Shape IDs, so
    the code survives the shape-ID reshuffling that Drive performs on every
    edit of the template. Each label is paired with its nearest % sibling in
    the same group (closest `top` coordinate).
    Extra non-equity categories (Debt, Gold & Silver, etc.) get cloned rows.
    """
    pct_map = {lb: pc for lb, pc, _ in parts}
    col_map = {lb: col for lb, _, col in parts}

    KNOWN_LABELS = {'Conservative', 'Balanced', 'Aggressive', 'Equity', 'Hybrid'}

    def _set_pct(pct_sh, cat):
        val = eq_total if cat == 'Equity' else pct_map.get(cat, 0)
        if val <= 0:
            new_text = ''
        elif val < 1:
            new_text = '<1%'
        else:
            new_text = f'{val:.0f}%'
        para = pct_sh.text_frame.paragraphs[0]
        if para.runs:
            para.runs[0].text = new_text
            for r in para.runs[1:]:
                r.text = ''
        else:
            para.text = new_text

    def _hide_label(lbl_sh):
        para = lbl_sh.text_frame.paragraphs[0]
        if para.runs:
            para.runs[0].text = ''

    for shape in slide.shapes:
        if shape.shape_type != 6:
            continue
        try:
            children = list(shape.shapes)
        except Exception:
            continue

        # Identify label shapes by exact text match and % shapes by '%' presence
        label_children = []   # list of (cat, shape)
        pct_children   = []
        for ch in children:
            if not ch.has_text_frame:
                continue
            t = ch.text_frame.text.strip()
            if t in KNOWN_LABELS:
                label_children.append((t, ch))
            elif '%' in t:
                pct_children.append(ch)

        # Pair each label with its closest (same-row) % sibling
        used = set()
        for cat, lbl_sh in label_children:
            best, best_dt = None, 1 << 30
            for ps in pct_children:
                if id(ps) in used:
                    continue
                dt = abs(ps.top - lbl_sh.top)
                if dt < best_dt:
                    best_dt, best = dt, ps
            if best is None:
                continue
            used.add(id(best))
            _set_pct(best, cat)
            if cat != 'Equity' and pct_map.get(cat, 0) == 0:
                _hide_label(lbl_sh)

    # ── Clone the "Hybrid" group for extra non-equity categories ─────────────
    eq_labels = {'Aggressive', 'Balanced', 'Conservative'}
    extra_cats = [(lb, pc, col_map.get(lb, '#808080'))
                  for lb, pc, _ in parts
                  if lb not in eq_labels and lb != 'Hybrid' and pc > 0]
    if not extra_cats:
        return

    # Find the Hybrid group: a GROUP that contains a child whose text == "Hybrid"
    hybrid_grp = None
    for shape in slide.shapes:
        if shape.shape_type != 6:
            continue
        try:
            for ch in shape.shapes:
                if ch.has_text_frame and ch.text_frame.text.strip() == 'Hybrid':
                    hybrid_grp = shape
                    break
        except Exception:
            continue
        if hybrid_grp is not None:
            break
    if hybrid_grp is None:
        return

    next_top = hybrid_grp.top + hybrid_grp.height + 80000  # 80000 EMU gap

    for label, pct, color in extra_cats:
        clone_el = deepcopy(hybrid_grp._element)

        # Reposition the cloned group — only the TOP-LEVEL group's grpSpPr
        # xfrm off.y should be touched. Using .find('.//xfrm') returns the
        # FIRST xfrm anywhere (which is the group's own), but to be explicit
        # we walk into grpSpPr directly so we never accidentally move a
        # child's xfrm.
        grpSpPr = clone_el.find(f'{{{NS_P}}}grpSpPr')
        if grpSpPr is not None:
            xfrm = grpSpPr.find(f'{{{NS_A}}}xfrm')
            if xfrm is not None:
                off = xfrm.find(f'{{{NS_A}}}off')
                if off is not None:
                    off.set('y', str(next_top))

        # Update children: identify the HYBRID LABEL shape (text == "Hybrid")
        # and the PERCENTAGE shape (text contains "%") by CURRENT text content.
        # Empty shapes (the background rect, the coloured dot) must be left
        # alone — they are NOT percentage fields.
        label_sp   = None
        pct_sp     = None
        dot_fill   = None
        for child_el in clone_el.iter(f'{{{NS_P}}}sp'):
            txBody = child_el.find(f'{{{NS_P}}}txBody')
            raw = ''
            if txBody is not None:
                raw = ''.join(
                    (t.text or '') for t in txBody.findall(f'.//{{{NS_A}}}t')
                ).strip()
            if raw == 'Hybrid':
                label_sp = child_el
            elif '%' in raw:
                pct_sp = child_el
            elif raw == '':
                # Candidate for color dot — pick the SMALLEST solid-filled
                # shape (the dot, not the background rect).
                spPr = child_el.find(f'{{{NS_P}}}spPr')
                if spPr is not None:
                    sf = spPr.find(f'.//{{{NS_A}}}solidFill')
                    if sf is not None:
                        clr = sf.find(f'{{{NS_A}}}srgbClr')
                        if clr is not None:
                            # Track the one with the smallest ext (the dot).
                            xf = spPr.find(f'{{{NS_A}}}xfrm')
                            if xf is not None:
                                ext = xf.find(f'{{{NS_A}}}ext')
                                if ext is not None:
                                    try:
                                        area = int(ext.get('cx', 0)) * int(ext.get('cy', 0))
                                    except (TypeError, ValueError):
                                        area = 0
                                    if dot_fill is None or area < dot_fill[0]:
                                        dot_fill = (area, clr)

        # Write the label text
        if label_sp is not None:
            txBody = label_sp.find(f'{{{NS_P}}}txBody')
            t_els = txBody.findall(f'.//{{{NS_A}}}t') if txBody is not None else []
            if t_els:
                t_els[0].text = label
                for t in t_els[1:]:
                    t.text = ''

        # Write the percentage text
        if pct_sp is not None:
            txBody = pct_sp.find(f'{{{NS_P}}}txBody')
            t_els = txBody.findall(f'.//{{{NS_A}}}t') if txBody is not None else []
            if t_els:
                pct_text = '<1%' if pct < 1 else f'{pct:.0f}%'
                t_els[0].text = pct_text
                for t in t_els[1:]:
                    t.text = ''

        # Recolor the dot
        if dot_fill is not None:
            dot_fill[1].set('val', color.lstrip('#'))

        slide.shapes._spTree.append(clone_el)
        next_top += hybrid_grp.height + 80000

    # ── Move 'Small + Mid Allocation' text below last legend row ─────────────
    # next_top now points past the last clone; last clone bottom = next_top - 80000
    last_clone_bottom = next_top - 80000
    for sh in slide.shapes:
        if (sh.has_text_frame and
                'Small' in sh.text_frame.text and
                'Mid' in sh.text_frame.text):
            sh.top = last_clone_bottom + 45000   # ~0.05 cm gap
            break

# ──────────────────────────────────────────────────────────────
# APPENDIX — Scheme Slides
# ──────────────────────────────────────────────────────────────

def do_appendix(prs, pf_id, data):
    """Build scheme-data slides and insert after Appendix divider (index 22)."""
    sch = data['scheme'][data['scheme']['PF_ID'] == pf_id].copy()
    cat = data['categorization']

    if sch.empty:
        print("  Appendix: no schemes — deleting template slides")
        for i in [25, 24, 23, 22]:
            delete_slide(prs, i)
        return 0

    c_map = dict(zip(cat['Proposed Sub-Category'], cat['Powerup Broad Category']))
    s_map = dict(zip(cat['Proposed Sub-Category'], cat['Final Serialing']))
    n_map = dict(zip(cat['Proposed Sub-Category'], cat['Names']))

    # Arbitrage funds belong under Debt, not Hybrid
    c_map['ARBITRAGE_FUND']         = 'Debt'
    c_map['FUND_OF_FUNDS_ARBITRAGE'] = 'Debt'

    sch['_cat']  = sch['UPDATED_SUBCATEGORY'].map(c_map).fillna(
                       sch['UPDATED_SUBCATEGORY'].str.replace('_', ' ').str.title())
    sch['_sort'] = sch['UPDATED_SUBCATEGORY'].map(s_map).fillna(999)
    sch['_disp'] = sch['UPDATED_SUBCATEGORY'].map(n_map).fillna(
                       sch['UPDATED_SUBCATEGORY'].str.replace('_', ' ').str.title())

    # Drop rows with no subcategory (can't be grouped or displayed)
    sch = sch.dropna(subset=['UPDATED_SUBCATEGORY'])

    # Ordered unique subcategories by sort number
    seen = {}
    for _, r in sch.sort_values('_sort').iterrows():
        sc = r['UPDATED_SUBCATEGORY']
        if sc not in seen:
            seen[sc] = r['_sort']

    groups = []
    for subcat in seen:
        grp = sch[sch['UPDATED_SUBCATEGORY'] == subcat].sort_values(
            'CURRENT_VALUE', ascending=False)
        if grp.empty:
            continue
        groups.append(dict(
            cat  = grp['_cat'].iloc[0],
            disp = grp['_disp'].iloc[0],
            sort = grp['_sort'].iloc[0],
            rows = list(grp.iterrows()),
        ))

    # Build slide specs.
    # Packed specs ('tpl'='packed') contain two independent sections on one slide:
    #   - sections[0]: first subcategory (1 row)
    #   - sections[1]: second subcategory (1-2 rows)
    specs = []
    gi = 0
    while gi < len(groups):
        g = groups[gi]
        rows = list(g['rows'])

        # Try to pack: if this group has exactly 1 row AND next group has 1-2 rows
        # Only pack within the same broad category (never mix Equity with Hybrid, etc.)
        if len(rows) == 1 and gi + 1 < len(groups):
            g_next = groups[gi + 1]
            n_next = min(2, len(g_next['rows']))
            if g['cat'] == g_next['cat'] and 1 <= n_next <= 2:
                specs.append(dict(
                    tpl='packed',
                    sec2_n=n_next,
                    sections=[
                        dict(cat=g['cat'],      disp=g['disp'],
                             rows=rows),
                        dict(cat=g_next['cat'], disp=g_next['disp'],
                             rows=list(g_next['rows'][:n_next])),
                    ],
                ))
                remaining = list(g_next['rows'][n_next:])
                if remaining:
                    groups.insert(gi + 2, dict(cat=g_next['cat'], disp=g_next['disp'],
                                               sort=g_next['sort'], rows=remaining))
                gi += 2
                continue

        while rows:
            n = min(4, len(rows))
            specs.append(dict(tpl=n, cat=g['cat'], disp=g['disp'], rows=rows[:n]))
            rows = rows[n:]
        gi += 1

    print(f"  Appendix: {len(sch)} schemes -> {len(specs)} slides")

    tpl_idx = {4: 22, 3: 23, 2: 24, 1: 25}

    # Keep references to the template slides BEFORE any cloning shifts their indices.
    tpl_slides_ref = {n: prs.slides[tpl_idx[n]] for n in [1, 2, 3, 4]}

    new_slides = []
    for sp in specs:
        if sp['tpl'] == 'packed':
            # Base slide: always use the 1-row template (section 1)
            ns = clone_slide(prs, tpl_idx[1])
            # Append section 2 shapes (shifted downward) from the appropriate template
            sec2_src = tpl_slides_ref[sp['sec2_n']]
            _clone_section_onto_slide(sec2_src, ns, _SEC2_Y_OFFSET)
        else:
            ns = clone_slide(prs, tpl_idx[sp['tpl']])
        new_slides.append((ns, sp))

    for i in [25, 24, 23, 22]:
        delete_slide(prs, i)

    n = len(specs)
    first_new = len(prs.slides) - n
    for i in range(n):
        move_slide(prs, first_new + i, 22 + i)

    for ns, sp in new_slides:
        if sp['tpl'] == 'packed':
            sec1, sec2 = sp['sections']
            _fill_scheme_slide(ns, sec1, y_min=0,             y_max=_SEC_Y_SPLIT)
            _fill_scheme_slide(ns, sec2, y_min=_SEC_Y_SPLIT,  y_max=9_000_000)
        else:
            _fill_scheme_slide(ns, sp)

    print(f"  Appendix: {n} slides created & filled")
    return n


def do_hyperlinks(prs, n_appendix):
    """Wire internal hyperlinks: slide4 'see here' -> first appendix; appendix 'Go back' -> slide4."""
    if n_appendix == 0:
        return

    REL_SLIDE = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide'
    NS_R_ID   = f'{{{NS_R}}}id'

    slide4     = prs.slides[3]
    first_app  = prs.slides[22]   # first appendix slide is always at index 22

    # ── Slide 4 "see here" → first appendix slide ──────────────────
    for shape in slide4.shapes:
        if ';190;' not in shape.name:
            continue
        for hlink in shape._element.iter(f'{{{NS_A}}}hlinkClick'):
            rId = slide4.part.relate_to(first_app.part, REL_SLIDE)
            hlink.set(NS_R_ID, rId)
            hlink.set('action', 'ppaction://hlinksldjump')
            print(f"  Hyperlink: slide4 -> appendix[0]  rId={rId}")

    # ── Each appendix slide "Go back" → slide 4 ─────────────────────
    for i in range(n_appendix):
        app_slide = prs.slides[22 + i]
        rId = app_slide.part.relate_to(slide4.part, REL_SLIDE)
        for shape in app_slide.shapes:
            for hlink in shape._element.iter(f'{{{NS_A}}}hlinkClick'):
                hlink.set(NS_R_ID, rId)
                hlink.set('action', 'ppaction://hlinksldjump')
        print(f"  Hyperlink: appendix[{i}] -> slide4  rId={rId}")


def _clone_shape_at(slide, template_shape, new_y, new_text):
    """Clone a slide shape, reposition it vertically, and update its first text run."""
    clone_el = deepcopy(template_shape._element)
    # Give the clone a unique shape id
    existing_ids = set()
    for sh in slide.shapes:
        try:
            existing_ids.add(sh.shape_id)
        except Exception:
            pass
    new_id = max(existing_ids) + 1 if existing_ids else 9999
    for el in clone_el.iter():
        if el.get('id') is not None:
            el.set('id', str(new_id))
            break
    # Reposition: update the xfrm off y
    xfrm = clone_el.find(f'.//{{{NS_A}}}xfrm')
    if xfrm is not None:
        off = xfrm.find(f'{{{NS_A}}}off')
        if off is not None:
            off.set('y', str(int(new_y)))
    # Update text (first t element)
    for t_el in clone_el.iter(f'{{{NS_A}}}t'):
        t_el.text = new_text
        break
    slide.shapes._spTree.append(clone_el)


# y boundaries for section shapes on a scheme template slide
_SECTION_MIN_Y = 1_200_000
_SECTION_MAX_Y = 4_795_000

# tpl=1 bg bottom — section 1 ends here; section 2 starts after this + gap
_SEC1_BG_BOTTOM = 2_475_195
_SEC1_BG_TOP    = 1_315_158
_SEC_GAP        = 200_000   # EMU gap between the two sections
# y_offset: shift cloned section 2 shapes so section 2 starts right after section 1
_SEC2_Y_OFFSET  = (_SEC1_BG_BOTTOM + _SEC_GAP) - _SEC1_BG_TOP   # 1_360_037
_SEC_Y_SPLIT    = _SEC1_BG_BOTTOM + _SEC_GAP // 2               # midpoint of gap


def _shift_element_y(el, dy):
    """Shift the top-level y coordinate of a shape XML element by dy EMU."""
    tag = etree.QName(el.tag).localname
    if tag == 'graphicFrame':
        # p:xfrm > a:off
        xfrm = el.find(f'{{{NS_P}}}xfrm')
        if xfrm is not None:
            off = xfrm.find(f'{{{NS_A}}}off')
            if off is not None:
                off.set('y', str(int(off.get('y', '0')) + dy))
    else:
        # p:spPr / p:grpSpPr > a:xfrm > a:off
        for pr_name in (f'{{{NS_P}}}spPr', f'{{{NS_P}}}grpSpPr'):
            pr = el.find(pr_name)
            if pr is not None:
                xfrm = pr.find(f'{{{NS_A}}}xfrm')
                if xfrm is not None:
                    off = xfrm.find(f'{{{NS_A}}}off')
                    if off is not None:
                        off.set('y', str(int(off.get('y', '0')) + dy))
                break


def _clone_section_onto_slide(source_slide, target_slide, y_offset):
    """Clone section-level shapes from source_slide to target_slide, shifted by y_offset.

    Section shapes are those whose top is in [_SECTION_MIN_Y, _SECTION_MAX_Y).
    """
    existing_ids = set()
    for sh in target_slide.shapes:
        try:
            existing_ids.add(sh.shape_id)
        except Exception:
            pass
    next_id = max(existing_ids) + 1 if existing_ids else 9000

    # Build an image rId map: source rId -> target rId (only for images on source)
    img_rId_map = {}
    for src_rId, rel in source_slide.part.rels.items():
        if 'image' in rel.reltype:
            img_part = rel.target_part
            new_rId = target_slide.part.relate_to(img_part, rel.reltype)
            img_rId_map[src_rId] = new_rId

    for shape in source_slide.shapes:
        try:
            top = shape.top
        except Exception:
            continue
        if not (_SECTION_MIN_Y <= top < _SECTION_MAX_Y):
            continue
        # Skip standalone picture shapes — rating images are re-added by _fill_scheme_slide.
        # Groups may contain a separator-line image that must be preserved.
        if shape.shape_type == 13:
            continue

        clone_el = deepcopy(shape._element)

        # Remap any image rIds inside the clone (e.g. separator-line image inside GROUP)
        for blip in clone_el.iter(f'{{{NS_A}}}blip'):
            old_rId = blip.get(f'{{{NS_R}}}embed')
            if old_rId and old_rId in img_rId_map:
                blip.set(f'{{{NS_R}}}embed', img_rId_map[old_rId])

        # Assign fresh unique IDs to every element that carries an 'id' attribute
        for el in clone_el.iter():
            if el.get('id') is not None:
                el.set('id', str(next_id))
                next_id += 1

        # Shift the shape's y coordinate
        _shift_element_y(clone_el, y_offset)

        target_slide.shapes._spTree.append(clone_el)


def _fill_scheme_slide(slide, spec, y_min=0, y_max=9_000_000):
    cat_name  = spec['cat']
    sub_name  = spec['disp']
    rows_data = spec['rows']

    # Normalize "Equity - Index" (and variants) to "Equity"
    if 'equity' in cat_name.lower() and 'index' in cat_name.lower():
        cat_name = 'Equity'

    # Update category and sub-category text labels (left side, short text)
    # Only touch shapes within the y_min/y_max window so packed sections don't clobber each other.
    KNOWN_CATS = {
        'Equity', 'Hybrid', 'Debt', 'Gold & Silver', 'Fund of Funds',
        'Equity - Index', 'Global Funds', 'Solution Oriented',
        'Precious Metals', 'Debt - Index', 'Solution', 'Gold',
    }
    KNOWN_SUBCATS = {
        'Flexi Cap', 'Mid Cap', 'Small Cap', 'Large Cap', 'Index Large',
        'Value & Contra', 'ELSS', 'Focused Fund', 'Multi Cap',
        'Large & Mid', 'Dividend Yield', 'Short Duration', 'Liquid',
        'Hybrid Aggressive', 'Hybrid Conservative', 'Balanced Advantage',
        'Gold ETF', 'Gold Fund', 'Intl Fund', 'Arbitrage', 'Dynamic Bond',
        'Index Mid/Small', 'Index Fund', 'Thematic', 'Sectoral',
    }
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        try:
            top = shape.top
        except Exception:
            continue
        if not (y_min <= top < y_max):
            continue
        txt = shape.text_frame.text.strip()
        if shape.left < 2_000_000 and len(txt) < 45:
            if txt in KNOWN_CATS:
                replace_text(shape, cat_name)
            elif txt in KNOWN_SUBCATS:
                replace_text(shape, sub_name)

    # Collect table cells by Y-position within window
    tables_by_y = {}
    for shape in slide.shapes:
        if shape.has_table:
            try:
                top = shape.top
            except Exception:
                continue
            if y_min <= top < y_max:
                tables_by_y.setdefault(top, {})[shape.left] = shape
    row_ys = sorted(tables_by_y.keys())

    # Remove existing power-rating images within the y window (~256k square)
    for shape in list(slide.shapes):
        if shape.shape_type == 13:
            try:
                _ = shape.image
                if 200_000 < shape.width < 320_000 and 200_000 < shape.height < 320_000:
                    if y_min <= shape.top < y_max:
                        remove_shape(slide, shape)
            except Exception:
                pass

    # Fill each row
    for ri, y in enumerate(row_ys):
        cells = tables_by_y[y]
        if ri >= len(rows_data):
            for _, sh in cells.items():
                set_table_cell(sh.table.cell(0, 0), '')
            continue

        _, sr = rows_data[ri]

        for x, sh in sorted(cells.items()):
            cell = sh.table.cell(0, 0)
            if   abs(x - X_NAME)   < 60_000:
                set_table_cell(cell, sr.get('FUND_NAME', sr.get('FUND_STANDARD_NAME', '')))
            elif abs(x - X_RATING) < 60_000:
                set_table_cell(cell, '')
            elif abs(x - X_VALUE)  < 60_000:
                set_table_cell(cell, fmt_scheme_val(sr.get('CURRENT_VALUE', 0),
                                                    sr.get('% of PF', 0)))
            elif abs(x - X_XIRR)   < 60_000:
                set_table_cell(cell, fmt_xirr_pair(sr.get('XIRR_VALUE'),
                                                   sr.get('BM_XIRR')))
            elif abs(x - X_MISSED) < 60_000:
                set_table_cell(cell, fmt_missed(sr.get('MG_AS_ON_APP', 0)))

        # Add power-rating image (or centered '-' if none)
        rating = sr.get('POWERRATING')
        has_rating = pd.notna(rating) and str(rating) in RATING_IMAGES
        if not has_rating:
            for x, sh in sorted(cells.items()):
                if abs(x - X_RATING) < 60_000:
                    cell = sh.table.cell(0, 0)
                    set_table_cell(cell, '-')
                    # Normalize font size
                    para = cell.text_frame.paragraphs[0]
                    if para.runs:
                        para.runs[0].font.size = Pt(9)
                    # Horizontal center
                    pPr = para._p.find(f'{{{NS_A}}}pPr')
                    if pPr is None:
                        pPr = etree.SubElement(para._p, f'{{{NS_A}}}pPr')
                    pPr.set('algn', 'ctr')
                    # Vertical center
                    tc = cell._tc
                    tcPr = tc.find(f'{{{NS_A}}}tcPr')
                    if tcPr is None:
                        tcPr = etree.SubElement(tc, f'{{{NS_A}}}tcPr')
                    tcPr.set('anchor', 'ctr')
                    break
        if has_rating:
            img_path = _get_rating_image(str(rating))
            if os.path.exists(img_path):
                row_h = (row_ys[ri + 1] - y) if ri + 1 < len(row_ys) else 552450
                img_y = y + (row_h - RATING_IMG_SZ) // 2
                slide.shapes.add_picture(
                    img_path,
                    Emu(X_RATING_IMG), Emu(img_y),
                    Emu(RATING_IMG_SZ), Emu(RATING_IMG_SZ),
                )

# ──────────────────────────────────────────────────────────────
# SLIDE 6 — What's working well
# ──────────────────────────────────────────────────────────────

def do_slide6(prs, pf, risk_profile):
    """Fill slide 6 'What's working well' with two dynamic boxes.

    Box 01: Consistent investing discipline (always shown)
      Description: "SIPs & lump sum over {years} years building a corpus of {current_value}"

    Box 02 (conditional on PF XIRR vs BM XIRR):
      Variant A — PF XIRR > BM XIRR:
        Title: "Delivering competitive returns"
        Desc:  "Portfolio performance has edged past benchmark by {diff}%, you're on the right track"
      Variant B — otherwise:
        Title: "Aligned to your risk profile"
        Desc:  "Your portfolio reflects your preferred risk level: {risk_profile}*"
    """
    slide = prs.slides[5]   # 0-based index 5 = slide 6

    years_raw = pf.get('YEARS_SINCE_FIRST_TRANSACTION', 0)
    try:
        years_int = int(float(years_raw)) if not pd.isna(years_raw) else 0
    except Exception:
        years_int = 0

    cv      = pf.get('PF_CURRENT_VALUE', 0)
    cv_str  = fmt_inr_rupee(cv, prefix='')

    pf_xirr = pf.get('PF_XIRR', 0)
    bm_xirr = pf.get('BM_XIRR', 0)
    try:
        use_competitive = float(pf_xirr) > float(bm_xirr)
        diff = (float(pf_xirr) - float(bm_xirr)) * 100
    except Exception:
        use_competitive = False
        diff = 0.0

    # Template text markers — match by content, not by hardcoded shape IDs.
    # Box 01: heading "Consistent investing discipline" (left unchanged),
    #         body starts with "SIPs & lump sum".
    # Box 02: heading is one of three variants (Solid Absolute Gains / etc.);
    #         body starts with "Solid absolute gains of" or one of the variant
    #         body prefixes.
    box01_body = f'SIPs & lump sum over {years_int} years building a corpus of {cv_str}'
    if use_competitive:
        box02_heading = 'Delivering competitive returns'
        box02_body    = (f"Portfolio performance has edged past benchmark "
                         f"by {diff:.1f}%, you're on the right track")
    else:
        box02_heading = 'Aligned to your risk profile'
        box02_body    = f'Your portfolio reflects your preferred risk level: {risk_profile}*'

    BOX02_HEADINGS = {
        'Solid Absolute Gains',
        'Delivering competitive returns',
        'Aligned to your risk profile',
    }

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        txt = shape.text_frame.text.strip()
        tl  = txt.lower()
        if 'sips' in tl and 'lump sum' in tl:
            replace_text(shape, box01_body)
        elif txt in BOX02_HEADINGS:
            replace_text(shape, box02_heading)
        elif ('solid absolute gains of' in tl
              or 'portfolio performance has edged' in tl
              or 'portfolio reflects your preferred' in tl):
            replace_text(shape, box02_body)

    variant = 'Delivering competitive returns' if use_competitive else 'Aligned to your risk profile'
    print(f"  Slide 6: {years_int}y, {cv_str}, variant='{variant}'")


# ──────────────────────────────────────────────────────────────
# SLIDE 13 — Portfolio vs Infinite comparison chart
# ──────────────────────────────────────────────────────────────

# Maps risk profile → TYPE prefix used in Lines/Results CSVs
RISK_TYPE_PREFIX = {
    'Very Aggressive':  'VA',
    'Aggressive':       'A',
    'Balanced':         'B',
    'Conservative':     'C',
    'Very Conservative':'VC',
}


def _best_infinite_type(pf_id, prefix, results_df):
    """Pick variant-1 lumpsum type for the given risk prefix."""
    cust = results_df[results_df['PF_ID'] == pf_id]
    # Always prefer variant 1 lumpsum (e.g. "VA1 - lumpsum - 24M")
    v1_lump = f'{prefix}1 - lumpsum - 24M'
    if not cust[cust['TYPE'] == v1_lump].empty:
        return v1_lump
    # Fall back: any variant-1 type
    v1 = cust[cust['TYPE'].str.startswith(f'{prefix}1')]
    if not v1.empty:
        return v1.iloc[0]['TYPE']
    # Final fall back: any type with that prefix
    any_pref = cust[cust['TYPE'].str.startswith(prefix)]
    return any_pref.iloc[0]['TYPE'] if not any_pref.empty else None


def do_slide13(prs, pf_id, risk_profile, data):
    slide = prs.slides[12]  # index 12 = slide 13

    prefix = RISK_TYPE_PREFIX.get(risk_profile, 'B')
    inf_type = _best_infinite_type(pf_id, prefix, data['results'])
    if inf_type is None:
        print(f"  Slide 13: no Infinite type found for prefix '{prefix}' — skipping")
        return

    # ── Get chart lines ──────────────────────────────────────────────────────
    lines_df = data['lines']
    cust_lines = lines_df[lines_df['PF_ID'] == pf_id].copy()
    cust_lines['DATE'] = _parse_dates(cust_lines['DATE'])

    # Dedupe by DATE — Lines often has each date twice (one row stored as
    # "2024-09-01" and another as "2024-09-01 00:00:00"). Without this the
    # chart shows a sawtooth/zigzag pattern because matplotlib draws a line
    # between every consecutive pair, including the duplicate same-date pairs.
    pf_line  = (cust_lines[cust_lines['TYPE'] == 'pf']
                .sort_values('DATE')
                .drop_duplicates(subset='DATE', keep='last'))
    inf_line = (cust_lines[cust_lines['TYPE'] == inf_type]
                .sort_values('DATE')
                .drop_duplicates(subset='DATE', keep='last'))

    if pf_line.empty or inf_line.empty:
        print(f"  Slide 13: missing line data for '{pf_id}' — skipping")
        return

    # Clip invested value to the same date range as pf line
    date_min, date_max = pf_line['DATE'].min(), pf_line['DATE'].max()
    inv_df = data['invested'][data['invested']['PF_ID'] == pf_id].copy()
    inv_df['DATE'] = _parse_dates(inv_df['DATE'])
    inv_df = inv_df[(inv_df['DATE'] >= date_min) & (inv_df['DATE'] <= date_max)].sort_values('DATE')

    # Merge all lines on pf_line dates (forward-fill invested amount)
    dates = pf_line['DATE'].values
    pf_vals  = pf_line.set_index('DATE')['CURRENT_VALUE']
    inf_vals = inf_line.set_index('DATE')['CURRENT_VALUE']
    inv_vals = inv_df.set_index('DATE')['INVESTED_AMOUNT'] if not inv_df.empty else None

    # ── Results for XIRR table and end-value labels ───────────────────────────
    res = data['results'][data['results']['PF_ID'] == pf_id]
    pf_res  = res[res['TYPE'] == 'pf']
    inf_res = res[res['TYPE'] == inf_type]

    pf_xirr  = float(pf_res['XIRR'].iloc[0])  if not pf_res.empty  else 0.0
    inf_xirr = float(inf_res['XIRR'].iloc[0]) if not inf_res.empty else 0.0
    pf_final  = float(pf_res['CURRENT_VALUE'].iloc[0])  if not pf_res.empty  else 0.0
    inf_final = float(inf_res['CURRENT_VALUE'].iloc[0]) if not inf_res.empty else 0.0

    # ── Generate chart image ──────────────────────────────────────────────────
    IMG_W_EMU = 5110622
    IMG_H_EMU = 3160301
    ratio = IMG_W_EMU / IMG_H_EMU   # ≈ 1.617
    fig_h = 4.2
    fig_w = fig_h * ratio

    fig, ax = plt.subplots(figsize=(fig_w, fig_h))
    fig.patch.set_facecolor('white')
    ax.set_facecolor('white')
    fig.subplots_adjust(left=0.06, right=0.97, top=0.97, bottom=0.12)

    # Plot actual portfolio line — blue dotted
    ax.plot(pf_vals.index, pf_vals.values / 1e5,
            color='#4E9EED', linewidth=0.8, linestyle=(0, (3, 2)))

    # Plot Infinite strategy line — dark solid
    ax.plot(inf_vals.index, inf_vals.values / 1e5,
            color='#1A1A2E', linewidth=1.0)

    # X-axis: "01-01-YYYY" labels at each year, no y-axis labels, no grid
    ax.xaxis.set_major_formatter(matplotlib.dates.DateFormatter('%d-%m-%Y'))
    ax.xaxis.set_major_locator(matplotlib.dates.YearLocator())
    plt.setp(ax.get_xticklabels(), rotation=0, ha='center',
             fontsize=8, color='#555555')
    ax.yaxis.set_visible(False)
    ax.grid(visible=False)
    for spine in ax.spines.values():
        spine.set_visible(False)
    ax.tick_params(axis='x', length=0)

    buf13 = BytesIO()
    fig.savefig(buf13, format='png', dpi=200,
                facecolor='white', bbox_inches=None, pad_inches=0)
    plt.close(fig)
    buf13.seek(0)

    # ── Remove the old chart PICTURE (find by size/position, not shape ID) ──
    for shape in list(slide.shapes):
        if (shape.shape_type == 13
                and shape.width  > 3_000_000
                and shape.height > 2_000_000
                and shape.top    > 1_000_000):
            try:
                _ = shape.image
                remove_shape(slide, shape)
                print("  Slide 13: removed old chart image")
                break
            except Exception:
                pass

    new_pic = slide.shapes.add_picture(
        buf13,
        Emu(-152400), Emu(1666799),
        Emu(IMG_W_EMU), Emu(IMG_H_EMU),
    )
    # Move to same z-position (bottom of stack)
    sp_tree = slide.shapes._spTree
    new_el  = new_pic._element
    sp_tree.remove(new_el)
    sp_tree.insert(2, new_el)

    # ── Update text shapes by content-matching (resilient to shape-ID drift) ─
    import re as _re
    final_val_pat = _re.compile(r'^[₹]?\s*[\d.]+\s*L\s*$', _re.IGNORECASE)

    # Final-value labels (e.g. "41.5L", "49.7L"): sorted top→bottom,
    # upper = Infinite, lower = Actual.
    final_shapes = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        t = shape.text_frame.text.strip()
        if final_val_pat.match(t):
            final_shapes.append(shape)
    final_shapes.sort(key=lambda s: s.top)
    if len(final_shapes) >= 2:
        replace_text(final_shapes[0], fmt_inr_rupee(inf_final).replace('₹', ''))
        replace_text(final_shapes[1], fmt_inr_rupee(pf_final).replace('₹', ''))
    elif len(final_shapes) == 1:
        replace_text(final_shapes[0], fmt_inr_rupee(inf_final).replace('₹', ''))

    # Infinite label — starts with "Infinite" and contains a risk-profile word
    risk_words = ('Very Aggressive', 'Very Conservative',
                  'Aggressive', 'Conservative', 'Balanced')
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        t = shape.text_frame.text.strip()
        if t.startswith('Infinite') and any(w in t for w in risk_words):
            replace_text(shape, f'Infinite {risk_profile}')
            break

    # ── XIRR table: find the single TABLE on the slide ──────────────────────
    for shape in slide.shapes:
        if shape.has_table:
            tbl = shape.table
            set_table_cell(tbl.cell(1, 1), f'{inf_xirr * 100:.2f}%')
            set_table_cell(tbl.cell(1, 2), f'{pf_xirr * 100:.2f}%')
            break

    print(f"  Slide 13: chart generated ({risk_profile} → {inf_type})")
    print(f"  Slide 13: Actual XIRR={pf_xirr*100:.2f}%, Infinite XIRR={inf_xirr*100:.2f}%")
    print(f"  Slide 13: Actual final={fmt_inr_rupee(pf_final)}, Infinite={fmt_inr_rupee(inf_final)}")


# ──────────────────────────────────────────────────────────────
# SHAPE UTILITIES
# ──────────────────────────────────────────────────────────────

def _iter_shapes_recursive(shapes):
    """Yield every shape, descending into GROUP shapes."""
    for shape in shapes:
        yield shape
        if shape.shape_type == 6:   # MSO_SHAPE_TYPE.GROUP
            try:
                yield from _iter_shapes_recursive(shape.shapes)
            except Exception:
                pass


# ──────────────────────────────────────────────────────────────
# QUESTIONNAIRE — answer lookup
# ──────────────────────────────────────────────────────────────

def _get_answer(question_text, q_row, context=''):
    """
    Map a questionnaire slide question to the customer's Excel answer.
    Returns the formatted answer string, or None if not matched.
    context: 'vehicle' | 'home' | '' — slide-level context for ambiguous questions.
    """
    q = question_text.lower().strip()

    # Slide 1: Age, Employment Status, Income Source
    if 'your age' in q:
        return _safe_str(q_row.get('Age'))
    if 'employment status' in q:
        return _safe_str(q_row.get('Employment Status'))
    if 'source of income' in q:
        return _safe_str(q_row.get('Income Source'))

    # Slide 2: Goals, Liability Type, Liability Followup
    if 'reason for investing' in q or 'investing in mutual' in q:
        return _safe_str(q_row.get('Goals'))
    if 'types of liabilities' in q:
        return _safe_str(q_row.get('Liability Type'))
    if 'comfortably meet' in q:
        return _safe_str(q_row.get('Liability Followup Answer'))

    # Slide 3: Emergency Fund, Portfolio Preference, Investment Horizon
    if 'emergency fund' in q:
        return _safe_str(q_row.get('Emergency Fund'))
    if 'portfolio to grow' in q or ('prefer' in q and 'portfolio' in q):
        return _safe_str(q_row.get('Portfolio Preference'))
    if 'investment horizon' in q:
        return _safe_str(q_row.get('Investment Horizon'))

    # Slide 4: Fall Reaction, Lumpsum, Monthly SIP
    if 'investments fall' in q or 'fall by 20' in q:
        return _safe_str(q_row.get('Fall Reaction'))
    if 'lumpsum' in q:
        return _safe_inr(q_row.get('Lumpsum Amount (with Infinite)', 0))
    if 'monthly sip' in q and 'amount' in q:
        return _safe_inr(q_row.get('Monthly SIP Amount (with Infinite)', 0))

    # Post-Retirement Income Planning slide (context='postret')
    if context == 'postret':
        if 'discretionary' in q:
            return _safe_inr(q_row.get('PostRet: Discretionary Expenses', 0))
        if 'monthly income' in q and 'expense' in q:
            inc = q_row.get('PostRet: Passive+Pension Income', 0)
            exp = q_row.get('PostRet: Living Expenses', 0)
            return f'Income: {_safe_inr(inc)} ; Expenses: {_safe_inr(exp)}'
        if 'financial investments apart' in q or 'apart from mutual' in q:
            v = q_row.get('PostRet: Other Instruments')
            if v is not None and not (isinstance(v, float) and pd.isna(v)):
                return _safe_str(v)
            return _safe_str(q_row.get('Other Investments Value'))

    # Retirement slides 5-6
    if 'monthly income' in q and 'expense' in q:
        inc = q_row.get('Ret: Monthly Income', 0)
        exp = q_row.get('Ret: Monthly Expenses', 0)
        return f'Income: {_safe_inr(inc)} ; Expenses: {_safe_inr(exp)}'
    if 'change in expenses' in q:
        return _safe_pct(q_row.get('Ret: Expense Change %'))
    if 'current monthly investment' in q or ('monthly investment' in q and 'current' in q):
        return _safe_inr(q_row.get('Ret: Monthly Investment', 0))
    if 'year-on-year' in q or 'yoy' in q:
        return _safe_pct(q_row.get('Ret: YoY Investment Increase %'))
    if 'financial investments apart' in q or 'apart from mutual' in q:
        return _safe_inr(q_row.get('Other Investments Value'))
    if 'liabilities' in q and ('emi' in q or 'loan' in q):
        v = q_row.get('Ret: Liabilities Detail')
        return _safe_str(v) if not (isinstance(v, float) and pd.isna(v)) else 'None'

    # Home Purchase  (context='home' or unambiguous keywords)
    if context != 'vehicle':
        if 'when do you want to purchase' in q:
            return _safe_str(q_row.get('Home: Purchase Year'))
        if 'flexibility to shift' in q:
            v = q_row.get('Home: Flexibility Yrs')
            if isinstance(v, float) and pd.isna(v): return '-'
            try: return f'{int(v)} years'
            except Exception: return _safe_str(v)
        # Order matters: 'down payment' must come BEFORE 'debt financing'
        # because slide questions like "Will you take debt financing for the
        # down payment?" contain BOTH keywords. The first match wins, so
        # the down payment % field would otherwise be returned as 'Yes'
        # from Home: Loan Y/N.
        if 'down payment' in q:
            return _safe_pct(q_row.get('Home: Down Payment %'))
        if 'debt financing' in q or 'loan' in q:
            return _safe_str(q_row.get('Home: Loan Y/N'))
    if 'value of home' in q:
        return _safe_inr(q_row.get('Home: Value', 0))
    if 'monthly rent' in q:
        v = q_row.get('Home: Monthly Rent')
        if isinstance(v, float) and pd.isna(v): return 'None'
        return _safe_inr(v)

    # Children's Education
    def _slot_has_value(v):
        """True when a per-child column actually contains data.
        Treats None / NaN / empty string / lone dash as unpopulated.
        Needed because Google Sheets returns '' for blank cells, which
        would otherwise be counted as populated."""
        if v is None:
            return False
        if isinstance(v, float) and pd.isna(v):
            return False
        s = str(v).strip()
        return bool(s) and s != '-'

    if 'number of children' in q and 'education' in q:
        count = sum(1 for i in range(1, 5)
                    if _slot_has_value(q_row.get(f'Edu: Child {i} UG Year')))
        return str(max(count, 1))
    if 'undergraduate' in q and 'start year' in q:
        years = [_safe_str(q_row.get(f'Edu: Child {i} UG Year'))
                 for i in range(1, 5)
                 if _slot_has_value(q_row.get(f'Edu: Child {i} UG Year'))]
        return ' & '.join(years) if years else '-'
    if 'cost' in q and 'undergraduate' in q:
        costs = []
        for i in range(1, 5):
            c = q_row.get(f'Edu: Child {i} UG Cost')
            if c is not None and not (isinstance(c, float) and pd.isna(c)):
                try:
                    if float(c) > 0:
                        costs.append(_safe_inr(float(c)))
                except (ValueError, TypeError):
                    pass
        return ' & '.join(costs) if costs else '-'
    if 'postgraduate' in q and 'start year' in q:
        years = [_safe_str(q_row.get(f'Edu: Child {i} PG Year'))
                 for i in range(1, 5)
                 if _slot_has_value(q_row.get(f'Edu: Child {i} PG Year'))]
        return ' & '.join(years) if years else '-'
    if 'cost' in q and 'postgraduate' in q:
        costs = []
        for i in range(1, 5):
            c = q_row.get(f'Edu: Child {i} PG Cost')
            if c is not None and not (isinstance(c, float) and pd.isna(c)):
                try:
                    if float(c) > 0:
                        costs.append(_safe_inr(float(c)))
                except (ValueError, TypeError):
                    pass
        return ' & '.join(costs) if costs else '-'

    # Children's Marriage
    if 'number of children' in q and 'marriage' in q:
        names = []
        for i in range(1, 5):
            n = q_row.get(f'Marriage: Child {i} Name', '')
            if not (isinstance(n, float) and pd.isna(n)) and str(n).strip():
                names.append(str(n).strip())
        count = max(len(names), 1)
        if names:
            return f'{count} - {", ".join(names)}'
        return str(count)
    if 'timeframe' in q and 'marriage' in q:
        times = []
        for i in range(1, 5):
            t = q_row.get(f'Marriage: Child {i} Timeframe')
            # Treat None / NaN / empty string / lone dash as "no child in this slot"
            if t is None or (isinstance(t, float) and pd.isna(t)):
                continue
            s = str(t).strip()
            if not s or s == '-':
                continue
            # Normalize: bare number → "N years";  "6-8years" → "6-8 years"
            if re.fullmatch(r'\d+', s):
                s = f'{s} years'
            else:
                s = re.sub(r'(\d)\s*years?', r'\1 years', s)
            times.append(s)
        return ' & '.join(times) if times else '-'
    if 'budget for marriage' in q or ('budget' in q and 'marriage' in q):
        budgets = []
        for i in range(1, 5):
            b = q_row.get(f'Marriage: Child {i} Budget')
            if b is not None and not (isinstance(b, float) and pd.isna(b)):
                try:
                    fv = float(b)
                    if fv > 0:
                        budgets.append(_safe_inr(fv))
                except (ValueError, TypeError):
                    if str(b).strip():
                        budgets.append(str(b).strip())
        return ' & '.join(budgets) if budgets else '-'

    # Vehicle Purchase  (explicit 'vehicle' in question OR context='vehicle')
    if 'vehicle' in q or context == 'vehicle':
        # Order matters: most-specific checks first
        if 'flexibility' in q or 'shift' in q:
            v = q_row.get('Vehicle: Flexibility Yrs')
            if isinstance(v, float) and pd.isna(v): return '-'
            try: return f'{int(v)} years'
            except Exception: return _safe_str(v)
        if 'down payment' in q:
            return _safe_pct(q_row.get('Vehicle: Down Payment %'))
        if 'when do you want' in q or 'purchase year' in q:
            return _safe_str(q_row.get('Vehicle: Purchase Year'))
        if 'value of vehicle' in q or ('value' in q and context == 'vehicle'):
            return _safe_inr(q_row.get('Vehicle: Value', 0))
        if 'debt financing' in q or 'financing' in q:
            return _safe_str(q_row.get('Vehicle: Loan Y/N'))

    return None  # no match — leave template text unchanged


# ── Complete answer → subcaption lookup (sourced from questionnaire form screenshots) ──
ANSWER_SUBCAPTIONS = {
    # Employment status
    'actively working':                     'Engaged in a full-time or part-time job, business, or self-employed with regular active income.',
    'soon to be retiring (within 5 yrs)':   'Planning to retire within the next few years',
    'soon to be retiring':                  'Planning to retire within the next few years',
    'retired early':                        'Not currently working by choice, but financially independent and living off savings or investments.',
    'retired':                              'No longer in active employment; primarily dependent on pension, savings, or investment income for expenses.',
    # Income source
    'active income only':                   'Regular earnings from salary, freelancing, or business.',
    'active + passive income':              'A mix of regular job/business income and recurring passive streams.',
    'active + passive':                     'A mix of regular job/business income and recurring passive streams.',
    'passive income only':                  'Recurring income from house rental, dividends, interest etc.',
    'pension income only':                  'Monthly pension received after retirement.',
    'passive + pension':                    'Combination of investment-based income (rent, dividends, interest) and pension inflows.',
    'no regular source':                    'No income inflow.',
    # Liability type
    'none':                                 'No loans or dependents',
    'financial liabilities only':           'Loans/EMIs but no dependents.',
    'dependent liabilities only':           'People depend on your income, no loans.',
    'both financial & dependent':           'Loans/EMIs and dependents rely on your income.',
    'both financial and dependent':         'Loans/EMIs and dependents rely on your income.',
    # Meet liabilities
    'yes - comfortably':                    'I have enough surplus, no stress.',
    'just about':                           "I manage, but it's tight some months.",
    'no - struggling':                      'I often find it difficult to meet liabilities.',
    # Investment horizon
    'short-term goals':                     'Less than 3 years',
    'medium-term goals':                    '3–5 years',
    'medium to long-term goals':            '5–8 years',
    'long-term wealth creation':            'More than 8 years',
    # Fall reaction
    'exit all investments immediately':     'To prevent further loss',
    'exit partially':                       'Shift to safer options',
    'stay invested':                        "I'm comfortable with market fluctuations",
    'invest more':                          'I will average my cost',
}


def _set_answer(shape, text):
    """Overwrite the answer text (para 0) and set the correct subcaption in para 1
    using the hardcoded ANSWER_SUBCAPTIONS lookup.  If no subcaption exists for this
    answer, the italic para 1 is cleared so the wrong default never shows."""
    if not shape.has_text_frame:
        return
    tf  = shape.text_frame
    p0  = tf.paragraphs[0]

    # Write the answer into para 0
    if p0.runs:
        p0.runs[0].text      = str(text)
        p0.runs[0].font.name = 'IBM Plex Sans'
        for r in p0.runs[1:]:
            r.text = ''
    else:
        p0.text = str(text)

    if len(tf.paragraphs) <= 1:
        return  # no subcaption row in this shape — nothing more to do

    # Look up the correct subcaption for this answer
    subcap = ANSWER_SUBCAPTIONS.get(str(text).strip().lower(), '')

    # Write (or clear) every run in para 1+
    for para in tf.paragraphs[1:]:
        first = True
        for r in para.runs:
            if first:
                r.text = subcap
                first  = False
            else:
                r.text = ''
        # If para had no runs but we have a subcaption, nothing we can do without
        # rebuilding the run — the existing italic endParaRPr will keep styling.


def _parse_portfolio_pref(text):
    """
    Parse: 'Grow well with some ups and downs — Moderate risk, ~12% p.a., worst -7% / best +20%'
    Returns: (description, risk_level, avg_return_pct, worst_pct, best_pct)
    """
    parts = re.split(r'\s+[—–-]\s+', str(text), maxsplit=1)
    description = parts[0].strip()
    rest = parts[1] if len(parts) > 1 else ''

    risk_level = avg_return = worst = best = ''
    if rest:
        rm = re.match(r'^([^,]+)', rest)
        if rm:
            risk_level = rm.group(1).strip()
        ret_m = re.search(r'~?([\d.]+)%\s*p', rest)
        if ret_m:
            avg_return = ret_m.group(1)
        w_m = re.search(r'worst\s+([+-]?[\d.]+)%', rest)
        b_m = re.search(r'best\s+([+-]?[\d.]+)%', rest)
        if w_m:
            worst = w_m.group(1)
            if not worst.startswith('-') and not worst.startswith('+'):
                worst = f'-{worst}'
        if b_m:
            best = b_m.group(1)
            if not best.startswith('+') and not best.startswith('-'):
                best = f'+{best}'
    return description, risk_level, avg_return, worst, best


def _set_portfolio_pref(shape, text):
    """
    Format portfolio preference answer as 4 styled lines:
      Line 1: description               (black)
      Line 2: risk level                (gray #8E9393)
      Line 3: Average returns: X% p.a.  (black)
      Line 4: Worst & best case: Y to Z (black, Y=red, Z=green)
    """
    desc, risk, ret, worst, best = _parse_portfolio_pref(text)
    if not desc:
        _set_answer(shape, text)
        return

    tf   = shape.text_frame
    txBody = tf._txBody
    a_ns  = f'{{{NS_A}}}'

    # Helper to build a fresh <a:p> with one or more runs
    def _new_para(runs_spec):
        """runs_spec = list of (text, bold, color_hex_or_None, size_pt_or_None)"""
        p_el = etree.SubElement(txBody, f'{a_ns}p')
        pPr  = etree.SubElement(p_el, f'{a_ns}pPr')
        pPr.set('indent', '0')
        for txt, bold, color, size in runs_spec:
            r_el = etree.SubElement(p_el, f'{a_ns}r')
            rPr  = etree.SubElement(r_el, f'{a_ns}rPr', lang='en-IN')
            rPr.set('dirty', '0')
            if size:
                rPr.set('sz', str(int(size * 100)))
            if bold:
                rPr.set('b', '1')
            rPr.set('spc', '-200')
            # Font
            latin = etree.SubElement(rPr, f'{a_ns}latin')
            latin.set('typeface', 'IBM Plex Sans')
            # Color
            if color:
                solidFill = etree.SubElement(rPr, f'{a_ns}solidFill')
                srgb = etree.SubElement(solidFill, f'{a_ns}srgbClr')
                srgb.set('val', color.lstrip('#'))
            t_el = etree.SubElement(r_el, f'{a_ns}t')
            t_el.text = txt
        return p_el

    # Remove all existing paragraphs
    for p in txBody.findall(f'{a_ns}p'):
        txBody.remove(p)

    font_sz = 10  # pt
    _new_para([(desc, False, '232323', font_sz)])
    _new_para([(risk, False, '8E9393', font_sz)])
    if ret:
        _new_para([(f'Average returns: {ret}% p.a.', False, '232323', font_sz)])
    if worst and best:
        p_el = etree.SubElement(txBody, f'{a_ns}p')
        pPr  = etree.SubElement(p_el, f'{a_ns}pPr')
        pPr.set('indent', '0')

        def _run(t, color):
            r_el = etree.SubElement(p_el, f'{a_ns}r')
            rPr  = etree.SubElement(r_el, f'{a_ns}rPr', lang='en-IN')
            rPr.set('dirty', '0')
            rPr.set('sz', str(int(font_sz * 100)))
            rPr.set('spc', '-200')
            latin = etree.SubElement(rPr, f'{a_ns}latin')
            latin.set('typeface', 'IBM Plex Sans')
            solidFill = etree.SubElement(rPr, f'{a_ns}solidFill')
            srgb = etree.SubElement(solidFill, f'{a_ns}srgbClr')
            srgb.set('val', color.lstrip('#'))
            t_el = etree.SubElement(r_el, f'{a_ns}t')
            t_el.text = t
        _run('Worst & best case: ', '232323')
        _run(f'{worst}%', 'C0392B')   # red
        _run(' to ', '232323')
        _run(f'{best}%', '27AE60')    # green


def populate_questionnaire_slide(slide, q_row):
    """
    For each GROUP on a questionnaire slide, find the question child
    (larger font ~228600) and answer child (smaller font ~133350),
    then fill in the customer's answer.
    """
    # Detect slide context by scanning all shapes (including inside groups)
    slide_context = ''
    for shape in _iter_shapes_recursive(slide.shapes):
        if not shape.has_text_frame:
            continue
        t = shape.text_frame.text
        tl = t.lower()
        if 'Vehicle Purchase' in t or ('Vehicle' in t and 'Purchase' in t):
            slide_context = 'vehicle'
            break
        if 'Home Purchase' in t or ('Home' in t and 'Purchase' in t):
            slide_context = 'home'
            break
        # Post-Retirement Income Planning slide has the unique
        # "Annual discretionary expenses" question
        if 'discretionary' in tl:
            slide_context = 'postret'
            break

    for shape in slide.shapes:
        if shape.shape_type != 6:
            continue
        try:
            children = list(shape.shapes)
        except Exception:
            continue

        q_shape = None
        a_shape = None

        for ch in children:
            if not ch.has_text_frame:
                continue
            for para in ch.text_frame.paragraphs:
                for run in para.runs:
                    sz = run.font.size
                    if sz is None:
                        continue
                    if sz >= 200_000:   # ~228600 = question
                        if q_shape is None:
                            q_shape = ch
                    elif sz <= 150_000: # ~133350 = answer
                        if a_shape is None:
                            a_shape = ch
                    break
                break

        # Fallback: first text child = question, second = answer
        if q_shape is None or a_shape is None:
            tc = [c for c in children if c.has_text_frame]
            if len(tc) >= 2:
                q_shape = tc[0]
                a_shape = tc[1]
            else:
                continue

        q_text = q_shape.text_frame.text.strip()
        answer = _get_answer(q_text, q_row, context=slide_context)
        if answer is not None:
            q_lo = q_text.lower()
            if 'portfolio to grow' in q_lo or ('prefer' in q_lo and 'portfolio' in q_lo):
                _set_portfolio_pref(a_shape, answer)
            else:
                _set_answer(a_shape, answer)
            print(f"    Q: '{q_text[:55]}' -> '{answer[:40]}'")

# ──────────────────────────────────────────────────────────────
# RISK REWARD SLIDES — insert 4 slides from risk_reward deck at slide 15
# ──────────────────────────────────────────────────────────────

# Risk Reward deck is fetched from Drive lazily via _get_rr_deck_path()

# Risk profile → 0-based start index in the risk-reward deck (groups of 4)
RISK_REWARD_IDX = {
    'Very Aggressive':  0,
    'Aggressive':       4,
    'Balanced':         8,
    'Conservative':     12,
    'Very Conservative': 12,   # same as Conservative
}


def _cross_deck_clone(src_slide, dst_prs):
    """Clone a slide from src_slide (another Presentation) into dst_prs.
    Images are copied as raw bytes so there are no part-name collisions."""
    layout = dst_prs.slide_layouts[0]
    ns = dst_prs.slides.add_slide(layout)

    # Copy each image from source to destination package and remap rIds
    img_map = {}
    for rId, rel in src_slide.part.rels.items():
        if 'image' not in rel.reltype:
            continue
        img_bytes = rel.target_part.blob
        img_part  = dst_prs.part.package.get_or_add_image_part(BytesIO(img_bytes))
        new_rId   = ns.part.relate_to(img_part, rel.reltype)
        img_map[rId] = new_rId

    # Replace spTree content
    sp_tree = ns.shapes._spTree
    for ch in list(sp_tree):
        tag = etree.QName(ch.tag).localname
        if tag in ('sp', 'grpSp', 'pic', 'graphicFrame', 'cxnSp'):
            sp_tree.remove(ch)
    for ch in src_slide.shapes._spTree:
        tag = etree.QName(ch.tag).localname
        if tag in ('sp', 'grpSp', 'pic', 'graphicFrame', 'cxnSp'):
            el = deepcopy(ch)
            for blip in el.iter(f'{{{NS_A}}}blip'):
                old = blip.get(f'{{{NS_R}}}embed')
                if old and old in img_map:
                    blip.set(f'{{{NS_R}}}embed', img_map[old])
            sp_tree.append(el)

    # Copy background
    bg = src_slide._element.find(f'{{{NS_P}}}bg')
    if bg is not None:
        old_bg = ns._element.find(f'{{{NS_P}}}bg')
        if old_bg is not None:
            ns._element.remove(old_bg)
        ns._element.insert(0, deepcopy(bg))

    return ns


def _replace_slide_content(dst_slide, src_slide, dst_prs):
    """Replace the content (shapes + background) of dst_slide with src_slide,
    copying images from src_slide's presentation into dst_prs in-place.
    No slide add/delete — avoids part-name collisions."""
    # Copy image bytes from source and relate to destination slide
    img_map = {}
    for rId, rel in src_slide.part.rels.items():
        if 'image' not in rel.reltype:
            continue
        img_bytes = rel.target_part.blob
        img_part  = dst_prs.part.package.get_or_add_image_part(BytesIO(img_bytes))
        new_rId   = dst_slide.part.relate_to(img_part, rel.reltype)
        img_map[rId] = new_rId

    # Clear existing shapes from destination spTree
    sp_tree = dst_slide.shapes._spTree
    for ch in list(sp_tree):
        tag = etree.QName(ch.tag).localname
        if tag in ('sp', 'grpSp', 'pic', 'graphicFrame', 'cxnSp'):
            sp_tree.remove(ch)

    # Copy shapes from source, remapping image rIds
    for ch in src_slide.shapes._spTree:
        tag = etree.QName(ch.tag).localname
        if tag in ('sp', 'grpSp', 'pic', 'graphicFrame', 'cxnSp'):
            el = deepcopy(ch)
            for blip in el.iter(f'{{{NS_A}}}blip'):
                old = blip.get(f'{{{NS_R}}}embed')
                if old and old in img_map:
                    blip.set(f'{{{NS_R}}}embed', img_map[old])
            sp_tree.append(el)

    # Replace background
    bg = src_slide._element.find(f'{{{NS_P}}}bg')
    if bg is not None:
        old_bg = dst_slide._element.find(f'{{{NS_P}}}bg')
        if old_bg is not None:
            dst_slide._element.remove(old_bg)
        dst_slide._element.insert(0, deepcopy(bg))


def _fill_rr_goals(slide, goals):
    """Replace goal placeholders on a risk-reward slide.
    Handles all known patterns across all risk profiles:
      - {{main_goal}} / {{secondary_goal}} (Very Aggressive L1, Conservative L1/S1)
      - 'Wealth Growth'                    (Very Aggressive S1)
      - 'Financial Freedom\\n...'           (Aggressive L1/S1, Balanced L1/S1)
    Searches recursively through GROUP shapes."""
    primary   = goals[0] if goals else 'Wealth Creation'
    secondary = ', '.join(goals[1:]) if len(goals) > 1 else ''

    # Known hardcoded multi-paragraph goal texts (first line is sufficient to identify)
    HARDCODED_FIRST_LINES = {'Financial Freedom'}

    for shape in _iter_shapes_recursive(slide.shapes):
        if not shape.has_text_frame:
            continue
        txt = shape.text_frame.text
        first_line = shape.text_frame.paragraphs[0].text.strip() if shape.text_frame.paragraphs else ''

        if '{{main_goal}}' in txt or '{{secondary_goal}}' in txt:
            para = shape.text_frame.paragraphs[0]
            runs = para.runs
            if len(runs) >= 2:
                runs[0].text = primary
                runs[1].text = secondary
            elif len(runs) == 1:
                runs[0].text = primary
            print(f"  Risk Reward: goals placeholder -> '{primary}' / '{secondary}'")

        elif first_line in HARDCODED_FIRST_LINES:
            # Multi-paragraph hardcoded goal text: replace para[0]=primary, para[1]=secondary, clear rest
            paras = shape.text_frame.paragraphs
            # Para 0 → primary
            if paras[0].runs:
                paras[0].runs[0].text = primary
                for r in paras[0].runs[1:]:
                    r.text = ''
            else:
                paras[0].text = primary
            # Para 1 → secondary (if exists)
            if len(paras) > 1:
                if paras[1].runs:
                    paras[1].runs[0].text = secondary
                    for r in paras[1].runs[1:]:
                        r.text = ''
                else:
                    paras[1].text = secondary
            # Clear remaining paragraphs
            for para in paras[2:]:
                for r in para.runs:
                    r.text = ''
            print(f"  Risk Reward: hardcoded goals -> '{primary}' / '{secondary}'")

        elif txt.strip() == 'Wealth Growth':
            # Hardcoded goal text on S1 variant slides — replace with all customer goals
            # Replicate the L1 structure: primary run, <a:br/>, secondary run
            para = shape.text_frame.paragraphs[0]
            p_el = para._p   # lxml element for the paragraph

            # Set primary in the first run
            if para.runs:
                run0_el = para.runs[0]._r
                run0_el.find(f'{{{NS_A}}}t').text = primary
                # Remove any existing extra runs and <a:br/> elements
                for child in list(p_el):
                    tag = etree.QName(child.tag).localname
                    if tag in ('br', 'r') and child is not run0_el:
                        p_el.remove(child)
            else:
                r_el = etree.SubElement(p_el, f'{{{NS_A}}}r')
                t_el = etree.SubElement(r_el, f'{{{NS_A}}}t')
                t_el.text = primary
                run0_el = r_el

            # Append line break + secondary run if there are secondary goals
            if secondary:
                rPr = run0_el.find(f'{{{NS_A}}}rPr')
                br_el = etree.SubElement(p_el, f'{{{NS_A}}}br')
                if rPr is not None:
                    br_el.append(deepcopy(rPr))

                r2_el = etree.SubElement(p_el, f'{{{NS_A}}}r')
                # Build rPr for secondary run: copy primary rPr then shrink the font
                rPr2 = deepcopy(rPr) if rPr is not None else etree.SubElement(r2_el, f'{{{NS_A}}}rPr')
                rPr2.set('sz', '1000')   # 10pt (127000 EMU) — smaller than primary 12pt
                r2_el.append(rPr2)
                t2_el = etree.SubElement(r2_el, f'{{{NS_A}}}t')
                t2_el.text = secondary

            print(f"  Risk Reward: 'Wealth Growth' -> '{primary}' / '{secondary}'")


def do_risk_reward_slides(prs, risk_profile, goals=None):
    """
    Replace slides 15-18 (indices 14-17) in prs in-place with the 4 slides
    from the risk-reward deck that correspond to risk_profile.
    Uses in-place content replacement to avoid XML part-name collisions.
    """
    try:
        _rr_path = _get_rr_deck_path()
    except Exception as e:
        print(f"  Risk Reward: Drive fetch failed - skipping ({e})"); return
    rr_prs = Presentation(_rr_path)
    start  = RISK_REWARD_IDX.get(risk_profile, 8)

    count = 0
    for offset in range(4):
        dst_idx = 14 + offset
        src_idx = start + offset
        if dst_idx >= len(prs.slides) or src_idx >= len(rr_prs.slides):
            break
        try:
            _replace_slide_content(
                prs.slides[dst_idx], rr_prs.slides[src_idx], prs
            )
            if goals:
                _fill_rr_goals(prs.slides[dst_idx], goals)
            count += 1
        except Exception as e:
            print(f"  Risk Reward: WARNING slide {src_idx+1}: {e}")

    print(f"  Risk Reward: replaced {count} slides for '{risk_profile}' "
          f"(source idx {start}-{start+count-1})")


# ──────────────────────────────────────────────────────────────
# QUESTIONNAIRE SLIDES — populate, filter, renumber
# ──────────────────────────────────────────────────────────────

GOAL_KW_SLIDES = {
    'Home Purchase':        'Home Purchase',
    "Children's Education": "Children's Education",
    "Children's Marriage":  "Children's Marriage",
    'Vehicle Purchase':     'Vehicle Purchase',
    'Vehicle':              'Vehicle Purchase',      # fallback keyword
}


def do_questionnaire(prs, goals, q_row):
    """
    1. Populate all questionnaire slides with the customer's answers.
    2. Remove slides for goals the customer did not select.
    3. Renumber the remaining slides (X/total).
    """
    norm = set()
    for g in goals:
        gl = g.lower()
        # Distinguish "Post-Retirement Income Planning" from plain "Retirement Planning"
        if 'post' in gl and 'retirement' in gl:
            norm.add('Post-Retirement Income Planning')
        elif 'retirement' in gl:
            norm.add('Retirement Planning')
        if 'home'       in gl:  norm.add('Home Purchase')
        if 'education'  in gl:  norm.add("Children's Education")
        if 'marriage'   in gl:  norm.add("Children's Marriage")
        if 'vehicle'    in gl:  norm.add('Vehicle Purchase')

    # Find all questionnaire slides
    q_indices = []
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame and 'Infinite Questionnaire' in shape.text_frame.text:
                q_indices.append(i); break

    if not q_indices:
        print("  Questionnaire: no slides found"); return
    print(f"  Questionnaire: found {len(q_indices)} slides, goals={norm}")

    # ── Step 1: Populate answers ──────────────────────────────────────────────
    if not q_row.empty:
        for idx in q_indices:
            populate_questionnaire_slide(prs.slides[idx], q_row)
    else:
        print("  Questionnaire: no q_row — skipping population")

    # ── Step 2: Remove slides for unselected goals ────────────────────────────
    def _norm_ap(s): return s.replace('\u2019', "'").replace('\u2018', "'")

    to_del = []
    for idx in q_indices:
        slide = prs.slides[idx]
        title = ''
        for shape in slide.shapes:
            if shape.has_text_frame and 'Infinite Questionnaire' in shape.text_frame.text:
                title = shape.text_frame.text.strip(); break

        title_norm = _norm_ap(title)
        for kw, goal_name in GOAL_KW_SLIDES.items():
            if _norm_ap(kw) in title_norm and goal_name not in norm:
                to_del.append(idx)
                print(f"    Remove: '{title}'")
                break

    # Remove PG (postgraduate) slide if PG cost is 0 or NaN
    pg_cost = q_row.get('Edu: Child 1 PG Cost', 0) if not q_row.empty else 0
    try:
        pg_zero = pd.isna(pg_cost) or float(pg_cost) == 0
    except (TypeError, ValueError):
        pg_zero = True
    if pg_zero:
        for idx in q_indices:
            slide = prs.slides[idx]
            slide_has_pg = False
            for shape in slide.shapes:
                if shape.shape_type == 6:  # group
                    try:
                        for ch in shape.shapes:
                            if ch.has_text_frame and 'postgraduate' in ch.text_frame.text.lower():
                                slide_has_pg = True; break
                    except Exception:
                        pass
                if slide_has_pg: break
            if slide_has_pg and idx not in to_del:
                to_del.append(idx)
                print(f"    Remove: PG slide (PG cost is 0 or missing)")

    # Remove Retirement / Post-Retirement Goals slides based on goal selection.
    # There are three "| Goals" slides in the base deck:
    #   - Retirement Planning (5/6):        has "Expected change in expenses post-retirement"
    #   - Retirement Planning (6/6):        has "Expected year-on-year increase in investment"
    #         (a generic follow-up only relevant to ongoing Retirement Planning)
    #   - Post-Retirement Income Planning:  has "Annual discretionary expenses"
    # Rules:
    #   has Retirement Planning only       → keep 5/6 + 6/6,   delete post-ret
    #   has Post-Retirement Planning only  → keep post-ret,    delete 5/6 + 6/6
    #   has both                           → keep all three
    #   has neither                        → delete all three
    has_ret      = 'Retirement Planning' in norm
    has_post_ret = 'Post-Retirement Income Planning' in norm
    for idx in q_indices:
        slide = prs.slides[idx]
        title = ''
        is_retirement_variant  = False   # 5/6 "change in expenses post-retirement"
        is_postret_variant     = False   # new 5/6 "annual discretionary expenses"
        is_generic_goals       = False   # 6/6 "year-on-year increase"
        for shape in _iter_shapes_recursive(slide.shapes):
            if not shape.has_text_frame: continue
            t  = shape.text_frame.text
            tl = t.lower()
            if 'Infinite Questionnaire' in t and '| Goals' in t:
                title = t.strip()
            if 'change in expenses' in tl and 'post-retirement' in tl:
                is_retirement_variant = True
            if 'discretionary' in tl:
                is_postret_variant = True
            if 'year-on-year' in tl:
                is_generic_goals = True
        if not title:
            continue
        # Decide whether to keep this slide
        if is_postret_variant and not has_post_ret:
            if idx not in to_del:
                to_del.append(idx)
                print(f"    Remove: Post-Retirement Goals slide (goal not selected)")
        elif is_retirement_variant and not has_ret:
            if idx not in to_del:
                to_del.append(idx)
                print(f"    Remove: Retirement Goals slide (goal not selected)")
        elif is_generic_goals and not has_ret:
            # Generic 6/6 goals slide belongs to Retirement Planning — delete it if
            # regular Retirement Planning is not selected (even if post-retirement is).
            if idx not in to_del:
                to_del.append(idx)
                print(f"    Remove: Generic Goals (6/6) slide (Retirement Planning not selected)")

    for idx in sorted(set(to_del), reverse=True):
        delete_slide(prs, idx)

    # ── Step 3: Renumber ──────────────────────────────────────────────────────
    q_slides = []
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame and 'Infinite Questionnaire' in shape.text_frame.text:
                q_slides.append((i, shape)); break

    total = len(q_slides)
    for seq, (_, title_shape) in enumerate(q_slides, 1):
        old = title_shape.text_frame.text.strip()
        new = re.sub(r'\(\d+/\d+\)', f'({seq}/{total})', old)
        if new != old:
            replace_text(title_shape, new)

    print(f"  Questionnaire: {len(to_del)} removed, {total} remaining (renumbered)")

# ──────────────────────────────────────────────────────────────
# MAIN
# ──────────────────────────────────────────────────────────────

def generate_deck(pf_id, customer_name, data=None, questionnaire_name=None,
                  override_risk_profile=None):
    """Generate a personalised M2 strategy deck.

    Args:
        pf_id: Client PF_ID
        customer_name: Client display name
        data: Pre-loaded data dict (optional - loaded from Google Sheets if None)
        questionnaire_name: If provided, match questionnaire row by this name
                            instead of using customer_name for matching.
        override_risk_profile: One of RISK_SCALE values
            ('Very Conservative' / 'Conservative' / 'Balanced' / 'Aggressive'
             / 'Very Aggressive'). When provided, overrides the calculated
            risk profile from the questionnaire — propagates to slide 3
            (display), slide 4 (match check vs portfolio risk), slide 6,
            slide 13 (Infinite line type), and Risk Reward slide selection.

    Returns:
        (pptx_bytes: BytesIO, filename: str)
    """
    print(f"\n{'='*60}")
    print(f"  Generating deck for: {customer_name}  ({pf_id})")
    print(f"{'='*60}\n")

    # Load data
    if data is None:
        print("[1/8] Loading data...")
        data = load_data()

    # Validate PF_ID
    pf_df  = data["pf_level"]
    pf_row = pf_df[pf_df["PF_ID"] == pf_id]
    if pf_row.empty:
        raise ValueError(f"PF_ID '{pf_id}' not found in PF_level sheet.")
    pf_row = pf_row.iloc[0]

    # Match questionnaire row - PF_ID first, then saved name, then exact/partial
    qdf   = data["questionnaire"]
    q_row = pd.Series(dtype=object)
    match_name = questionnaire_name or customer_name

    if "PF_ID" in qdf.columns:
        match = qdf[qdf["PF_ID"] == pf_id]
        if not match.empty:
            q_row = match.iloc[0]
            print(f"  Questionnaire: matched by PF_ID -> '{q_row.get('Name', '?')}'")

    if q_row.empty and questionnaire_name:
        name_col = next((c for c in qdf.columns if c.lower() == "name"), None)
        if name_col:
            match = qdf[qdf[name_col].astype(str).str.strip().str.lower() == questionnaire_name.strip().lower()]
            if not match.empty:
                q_row = match.iloc[0]
                print(f"  Questionnaire: matched by saved mapping -> '{q_row.get(name_col, '?')}'")

    if q_row.empty:
        name_col = next((c for c in qdf.columns if c.lower() == "name"), None)
        if name_col:
            for _, r in qdf.iterrows():
                if str(r.get(name_col, "")).lower().strip() == match_name.lower().strip():
                    q_row = r
                    print(f"  Questionnaire: matched by exact name '{r[name_col]}'")
                    break

    if q_row.empty and match_name:
        name_col = next((c for c in qdf.columns if c.lower() == "name"), None)
        if name_col:
            first = match_name.lower().split()[0]
            for _, r in qdf.iterrows():
                if first in str(r.get(name_col, "")).lower():
                    q_row = r
                    print(f"  Questionnaire: partial name match '{r[name_col]}'")
                    break

    # Fuzzy fallback — handles spelling variations (e.g. "Sunder Raj" vs "Sundar")
    if q_row.empty and match_name:
        from difflib import SequenceMatcher
        name_col = next((c for c in qdf.columns if c.lower() == "name"), None)
        if name_col:
            ml = match_name.lower().strip()
            best_row, best_score, best_name = None, 0.0, ''
            for _, r in qdf.iterrows():
                qn = str(r.get(name_col, '')).strip()
                if not qn or qn.lower() == 'nan':
                    continue
                s = SequenceMatcher(None, ml, qn.lower()).ratio()
                if s > best_score:
                    best_score, best_name, best_row = s, qn, r
            if best_row is not None and best_score >= 0.5:
                q_row = best_row
                print(f"  Questionnaire: fuzzy match '{best_name}' "
                      f"({int(round(best_score * 100))}% similarity)")

    if q_row.empty:
        print(f"  WARNING: no questionnaire row for '{customer_name}'")

    # Riskgroup aggregation
    rg     = data["riskgroup"][data["riskgroup"]["PF_ID"] == pf_id]
    rg_agg = (rg.groupby("RISK_GROUP_L0")
               .agg({"% of PF": "sum", "CURRENT_VALUE": "sum"})
               .reset_index())

    # Risk profile (calculated from questionnaire, then optionally overridden)
    print("\n[2/8] Risk profile...")
    calculated = calc_risk_profile(q_row) if not q_row.empty else "Balanced"
    if override_risk_profile and override_risk_profile in RISK_SCALE:
        risk_profile = override_risk_profile
        print(f"  Override applied: '{calculated}' (calculated) -> "
              f"'{risk_profile}' (override)")
    else:
        risk_profile = calculated

    # Download base deck from Google Drive and open
    print("\n[2b/8] Downloading base deck from Drive...")
    base_deck_path = _get_base_deck_path()
    prs = Presentation(base_deck_path)
    print(f"  Opened base deck template (from Drive)")

    first_name = customer_name.split()[0] if customer_name else "Client"

    # Process slides (same order as local M2/app.py)
    print("\n[3/9] Slide 1 - Title")
    do_slide1(prs, customer_name)

    print("[4/9] Slide 2 - Welcome")
    do_slide2(prs, first_name)

    print("[5/9] Slide 3 - You at a Glance")
    if not q_row.empty:
        do_slide3(prs, q_row, risk_profile, pf_row=pf_row)
    else:
        print("  SKIPPED (no questionnaire data)")

    print("[6/9] Slide 4 - Portfolio Snapshot")
    is_demat_df = data.get('is_demat', pd.DataFrame())
    if not is_demat_df.empty and 'PF_ID' in is_demat_df.columns:
        is_demat_df = is_demat_df[is_demat_df['PF_ID'].astype(str).str.strip()
                                  == str(pf_id).strip()]
    do_slide4(prs, pf_row, rg_agg, risk_profile, is_demat_df=is_demat_df)

    print("[6b/9] Slide 6 - What's working well")
    do_slide6(prs, pf_row, risk_profile)

    print("[7/9] Slide 13 - Portfolio vs Infinite")
    do_slide13(prs, pf_id, risk_profile, data)

    print("[8/9] Appendix - Scheme Slides")
    n_appendix = do_appendix(prs, pf_id, data) or 0

    print("[8b/9] Hyperlinks")
    do_hyperlinks(prs, n_appendix)

    print("[9/9] Risk Reward Slides (15-18)")
    rr_goals = parse_goals(q_row.get("Goals", "")) if not q_row.empty else []
    do_risk_reward_slides(prs, risk_profile, goals=rr_goals)

    print("[10/9] Questionnaire Slides")
    goals = parse_goals(q_row.get("Goals", "")) if not q_row.empty else []
    do_questionnaire(prs, goals, q_row)

    # Save to BytesIO
    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)

    safe = re.sub(r"[^\w\s-]", "", customer_name).strip().replace(" ", "_")
    filename = f"{safe}_{pf_id[:12]}_deck.pptx"

    print(f"\n{'='*60}")
    print(f"  DONE  ->  {filename}")
    print(f"{'='*60}\n")
    return buf, filename
