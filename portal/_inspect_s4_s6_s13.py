import sys, io, os
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
from dotenv import load_dotenv
load_dotenv('.env')
from googleapiclient.discovery import build
from google.oauth2 import service_account
creds = service_account.Credentials.from_service_account_file(
    os.getenv('GOOGLE_SERVICE_ACCOUNT_JSON'),
    scopes=['https://www.googleapis.com/auth/drive']
)
drive = build('drive', 'v3', credentials=creds)
from config import M2_TEMPLATE_ID
req = drive.files().export_media(
    fileId=M2_TEMPLATE_ID,
    mimeType='application/vnd.openxmlformats-officedocument.presentationml.presentation'
)
from googleapiclient.http import MediaIoBaseDownload
buf = io.BytesIO()
downloader = MediaIoBaseDownload(buf, req)
done = False
while not done:
    _, done = downloader.next_chunk()
buf.seek(0)
with open('Base_Deck_Latest.pptx', 'wb') as f:
    f.write(buf.read())

from pptx import Presentation
p = Presentation('Base_Deck_Latest.pptx')

def fullwalk(shape, indent=0):
    pad = '  '*indent
    name = shape.name
    print(f'{pad}[{name}] type={shape.shape_type} L={shape.left} T={shape.top} W={shape.width} H={shape.height}')
    if shape.has_text_frame:
        tf = shape.text_frame
        for pi, para in enumerate(tf.paragraphs):
            runs_info = []
            for run in para.runs:
                sz = run.font.size
                runs_info.append(f'sz={sz} text={run.text!r}')
            full = para.text
            print(f'{pad}  para[{pi}] full={full!r} runs={runs_info}')
    if shape.shape_type == 6:
        try:
            for c in shape.shapes:
                fullwalk(c, indent+1)
        except Exception:
            pass

print(f'\n========== SLIDE 6 ==========')
for s in p.slides[5].shapes:
    fullwalk(s)

print(f'\n========== SLIDE 13 ==========')
for s in p.slides[12].shapes:
    fullwalk(s)
