"""
m3_engine.py — PowerUp Infinite M3 Transition Deck Generator.

Ported from M3-automation/ to read reference data from Google Sheets,
download template from Google Drive, and output PPTX as BytesIO.

Combined from: reference_data.py + excel_reader.py + deck_writer.py
"""

import copy
import os
import re
import tempfile
from collections import OrderedDict
from datetime import date
from io import BytesIO

import openpyxl
import pandas as pd
from pptx import Presentation
from pptx.oxml.ns import qn

import sheets
from config import M3_TEMPLATE_ID


def _norm_name(name):
    """Normalize fund name for fuzzy matching.
    Strips 'Fund'/'Fd', lowercases, collapses spaces, and normalises
    common spelling variants (e.g. 'Smallcap' <-> 'Small Cap').
    """
    s = str(name).strip().lower()
    s = re.sub(r'\bfund\b', '', s)
    s = re.sub(r'\bfd\b', '', s)       # AUM file abbreviates 'Fund' as 'Fd'
    # Normalise compound-cap words so both spellings resolve to the same token
    s = s.replace('smallcap', 'small cap')
    s = s.replace('midcap',   'mid cap')
    s = s.replace('largecap', 'large cap')
    s = re.sub(r'\s+', ' ', s).strip()
    return s


def load_reference_data():
    """
    Load reference data from Google Sheets (M3_SPREADSHEET_ID).
    Returns (lookup, rr_category):
      lookup: dict keyed by ISIN with aum_cr, powerrank, upside, downside, five_year
      rr_category: dict keyed by UPDATED_SUBCATEGORY name -> 5Y category avg return
    """
    lookup = {}

    # --- AUM + fund-name fallback ---
    aum_df = sheets.read_m3_aum()
    name_to_isin = {}
    norm_name_to_isin = {}
    for _, row in aum_df.iterrows():
        isin = str(row.get('ISIN', '')).strip()
        if not isin or isin == 'nan':
            continue
        lookup.setdefault(isin, {})
        aum_val = row.get('AUM')
        try:
            aum_float = float(aum_val)
            lookup[isin]['aum_cr'] = aum_float / 10_000_000 if pd.notna(aum_float) else None
        except (TypeError, ValueError):
            lookup[isin]['aum_cr'] = None
        fn = row.get('FUND_NAME')
        if fn and str(fn).strip() and str(fn).strip().lower() != 'nan':
            fn_str = str(fn).strip()
            name_to_isin[fn_str] = isin
            norm_name_to_isin[_norm_name(fn_str)] = isin

    lookup['_name_to_isin'] = name_to_isin
    lookup['_norm_name_to_isin'] = norm_name_to_isin
    lookup['_isin_to_name'] = {isin: fn for fn, isin in name_to_isin.items()}

    # --- Power Rank ---
    pr_df = sheets.read_m3_powerranking()
    for _, row in pr_df.iterrows():
        isin = str(row.get('ISIN', '')).strip()
        if not isin or isin == 'nan':
            continue
        lookup.setdefault(isin, {})
        pr = row.get('POWERRANK')
        try:
            lookup[isin]['powerrank'] = int(float(pr)) if pr and str(pr).strip() else None
        except (TypeError, ValueError):
            lookup[isin]['powerrank'] = None

    # --- Upside / Downside ---
    ud_df = sheets.read_m3_upside_downside()
    isin_col = 'Scheme ISIN' if 'Scheme ISIN' in ud_df.columns else 'ISIN'
    down_col = 'Downside Capture Ratio' if 'Downside Capture Ratio' in ud_df.columns else None
    up_col = 'Upside Capture Ratio' if 'Upside Capture Ratio' in ud_df.columns else None
    for _, row in ud_df.iterrows():
        isin = str(row.get(isin_col, '')).strip()
        if not isin or isin == 'nan':
            continue
        lookup.setdefault(isin, {})
        if down_col:
            d = row.get(down_col)
            try:
                lookup[isin]['downside'] = float(d) if d and str(d).strip() not in ('', '--', 'nan') else None
            except (TypeError, ValueError):
                lookup[isin]['downside'] = None
        if up_col:
            u = row.get(up_col)
            try:
                lookup[isin]['upside'] = float(u) if u and str(u).strip() not in ('', '--', 'nan') else None
            except (TypeError, ValueError):
                lookup[isin]['upside'] = None

    # --- 5Y Rolling Returns ---
    rr_df = sheets.read_m3_rolling_returns()
    # Coerce ROLLING_PERIOD to numeric
    rr_df['ROLLING_PERIOD'] = pd.to_numeric(rr_df['ROLLING_PERIOD'], errors='coerce')
    df5 = rr_df[rr_df['ROLLING_PERIOD'] == 60]
    rr_scheme = {}
    rr_category = {}
    for _, row in df5.iterrows():
        eid = str(row.get('ENTITYID', '')).strip()
        try:
            val = float(row.get('RETURN_VALUE', 0))
        except (TypeError, ValueError):
            continue
        if eid.startswith('INF'):
            rr_scheme[eid] = val
        else:
            rr_category[eid] = val

    for isin, ret in rr_scheme.items():
        lookup.setdefault(isin, {})
        lookup[isin]['five_year'] = ret

    print(f"  Reference data loaded: {len(lookup)} ISINs, {len(rr_category)} categories with 5Y avg")
    return lookup, rr_category


# ---------------------------------------------------------------------------
# Formatting helpers
# ---------------------------------------------------------------------------

def fmt_aum(aum_cr):
    """Format AUM in Crores with Indian comma notation. e.g. 76060.9 -> '76,061 Cr'"""
    if aum_cr is None:
        return '-'
    n = round(aum_cr)
    s = str(n)
    if len(s) <= 3:
        return f"{s} Cr"
    result = s[-3:]
    s = s[:-3]
    while s:
        result = s[-2:] + ',' + result
        s = s[:-2]
    return result.lstrip(',') + ' Cr'


def fmt_upside_downside(upside, downside):
    """Format as 'Upside | Downside' rounded to nearest integer. e.g. '93 | 89'"""
    if upside is None and downside is None:
        return '-'
    up_str = str(round(upside)) if upside is not None else '-'
    dn_str = str(round(downside)) if downside is not None else '-'
    return f"{up_str} | {dn_str}"


def fmt_powerrank(powerrank):
    """Format power rank as integer string."""
    if powerrank is None:
        return '-'
    return str(int(powerrank))
# ──────────────────────────────────────────────────────────────
# EXCEL READER — reads PF_Curation / PF_MasterPlan from uploaded workbook
# ──────────────────────────────────────────────────────────────


def get_curation_sheet(wb):
    for name in wb.sheetnames:
        if name.startswith("PF_Curation"):
            return wb[name]
    raise ValueError("No PF_Curation* sheet found in workbook")


def col_map(header_row):
    """Return {column_name: 0-based_index} for non-None header values."""
    return {v: i for i, v in enumerate(header_row) if v is not None}


def detect_sections(ws):
    """
    Scan for section header rows by content.
    Returns dict with 's1','s2','s3','s4' as 1-based row numbers.

    How each section is identified:
      s1 — "Row Labels | Total Selected Value"  → Portfolio summary (AUM overview)
      s2 — FUND_NAME header row that contains FOLIO_NUMBER or TOTAL_UNITS
            → Current holdings / action master (per-folio sell/retain instructions)
      s3 — "Row Labels | Sum of TOTAL_VALUE" or "Sum of Current Value Amount"
            → Risk-group pivot summary
      s4 — FUND_NAME header row that contains "Allocation M1" or "Buy Value Amount"
            → Ideal portfolio (target allocation, buy/sell plan, milestones)
    If the smarter detection cannot tell s2 from s4, it falls back to:
      s2 = first FUND_NAME row, s4 = last FUND_NAME row.
    """
    sections = {}
    fund_name_rows = []   # (row_idx, header_tuple)
    _S3_B_VALUES = {"Sum of TOTAL_VALUE", "Sum of Current Value Amount"}

    all_rows = list(ws.iter_rows(values_only=True))
    for row_idx, row in enumerate(all_rows, start=1):
        a = row[0] if row else None
        b = row[1] if len(row) > 1 else None
        if a == "Row Labels" and b == "Total Selected Value":
            sections['s1'] = row_idx
        elif a == "Row Labels" and b in _S3_B_VALUES:
            sections['s3'] = row_idx
        elif a == "FUND_NAME":
            fund_name_rows.append((row_idx, row))

    # Identify s2 (action master) vs s4 (ideal portfolio) from FUND_NAME rows
    # s4 has columns like "Allocation M1", "Buy Value Amount", "SIP Allocation Amount"
    # s2 has columns like "FOLIO_NUMBER", "TOTAL_UNITS", "Sell Plan"
    _S4_MARKERS = {"Allocation M1", "Buy Value Amount", "SIP Allocation Amount"}
    _S2_MARKERS = {"FOLIO_NUMBER", "TOTAL_UNITS", "TOTAL_VALUE"}

    s2_candidates = []
    s4_candidates = []
    for row_idx, hdr in fund_name_rows:
        cols = set(c for c in hdr if c is not None)
        if cols & _S4_MARKERS:
            s4_candidates.append(row_idx)
        elif cols & _S2_MARKERS:
            s2_candidates.append(row_idx)

    if s4_candidates:
        sections['s4'] = s4_candidates[-1]   # last match is safest
    if s2_candidates:
        sections['s2'] = s2_candidates[0]    # first match

    # Fallback: first / last FUND_NAME row
    if 's2' not in sections and len(fund_name_rows) >= 1:
        sections['s2'] = fund_name_rows[0][0]
    if 's4' not in sections and len(fund_name_rows) >= 2:
        sections['s4'] = fund_name_rows[-1][0]

    return sections


def _read_section(ws, header_row_1based, col_names=None):
    """
    Read rows from header_row onwards until a blank col-A row or end of sheet.
    Returns list of dicts keyed by header name.
    Last non-blank row before the next blank section is tagged __grand_total__
    if its col-A contains 'Grand Total' or 'grand total'.
    """
    all_rows = list(ws.iter_rows(values_only=True))
    header = all_rows[header_row_1based - 1]
    cmap = col_map(header)
    if col_names:
        cmap = {k: v for k, v in cmap.items() if k in col_names}

    data = []
    for row in all_rows[header_row_1based:]:       # rows after header
        a = row[0] if row else None
        if a is None:
            break                                   # blank col-A = end of section
        record = {}
        for name, idx in cmap.items():
            record[name] = row[idx] if idx < len(row) else None
        label = str(a).strip().lower()
        if 'grand total' in label:
            record['__grand_total__'] = True
        data.append(record)
    return data


def _read_isin_column(ws, count):
    """
    Scan the sheet for a standalone 'ISIN' header cell (not in a multi-column header row),
    then read `count` ISIN values from the rows immediately below it in the same column.
    Returns a list of `count` strings (or None for missing values).
    The ISIN header row can vary per client — never hardcoded.
    """
    isin_header_row = None
    isin_col_0based = None

    for row_idx, row in enumerate(ws.iter_rows(values_only=True), start=1):
        for col_0, val in enumerate(row):
            if val == 'ISIN':
                isin_header_row = row_idx
                isin_col_0based = col_0
                break
        if isin_header_row is not None:
            break

    if isin_header_row is None:
        print("  WARNING: no standalone ISIN column found in sheet")
        return [None] * count

    all_rows = list(ws.iter_rows(values_only=True))
    isins = []
    for row in all_rows[isin_header_row: isin_header_row + count]:
        val = row[isin_col_0based] if isin_col_0based < len(row) else None
        isins.append(str(val).strip() if val else None)

    # Pad if fewer rows than expected
    while len(isins) < count:
        isins.append(None)

    print(f"  ISIN column found at row {isin_header_row}, col {isin_col_0based + 1}: {sum(1 for x in isins if x)} ISINs read")
    return isins


def _get_masterplan_sheet(wb):
    """Return the PF_MasterPlan_* sheet if present, else None."""
    for name in wb.sheetnames:
        if name.startswith('PF_MasterPlan'):
            return wb[name]
    return None


def _read_masterplan(ws):
    """
    Read the PF_MasterPlan_* sheet and return an excel_data dict compatible
    with what read_excel() produces for the old PF_Curation format.

    Sheet layout:
      Row 1 — title "Master Transition Plan"
      Row 2 — "PFV:" | <value>
      Row 3 — blank
      Row 4 — column headers (FUND_NAME, FOLIO_NUMBER, ISIN, …)
      Row 5+ — data rows (last row is Grand Total)

    Returns dict with section1/section2/section3/section4.
    ISIN is already a column in this sheet — no separate lookup needed.
    section2 and section3 are returned empty (not needed for this format).
    """
    all_rows = list(ws.iter_rows(values_only=True))

    # PFV from row 2 col B
    pfv = all_rows[1][1] if len(all_rows) > 1 and len(all_rows[1]) > 1 else 0
    pfv = pfv or 0

    # Headers at row 4 (index 3)
    header = all_rows[3]
    cmap = col_map(header)

    # Data rows from row 5 onwards; stop at first fully-blank row (col A None)
    section4 = []
    for row in all_rows[4:]:
        a = row[0] if row else None
        if a is None:
            break
        record = {}
        for name, idx in cmap.items():
            record[name] = row[idx] if idx < len(row) else None
        label = str(a).strip().lower()
        if 'grand total' in label:
            record['__grand_total__'] = True
        section4.append(record)

    # section1 — only needs Grand Total with Total Selected Value for PFV
    section1 = [{'Total Selected Value': pfv, '__grand_total__': True}]

    non_gt = [r for r in section4 if not r.get('__grand_total__')]
    print(f"  MasterPlan format detected: {len(non_gt)} data rows, PFV={pfv:,.0f}")
    return {
        'section1': section1,
        'section2': [],
        'section3': [],
        'section4': section4,
    }


def read_excel(excel_path):
    wb = openpyxl.load_workbook(excel_path, data_only=True)

    # New format: PF_MasterPlan_* sheet takes priority
    mp_ws = _get_masterplan_sheet(wb)
    if mp_ws is not None:
        return _read_masterplan(mp_ws)

    # Old format: PF_Curation_* sheet
    ws = get_curation_sheet(wb)
    sections = detect_sections(ws)

    required = ['s1', 's2', 's3', 's4']
    missing = [k for k in required if k not in sections]
    if missing:
        raise ValueError(f"Could not detect sections: {missing}")

    result = {
        'section1': _read_section(ws, sections['s1']),
        'section2': _read_section(ws, sections['s2']),
        'section3': _read_section(ws, sections['s3']),
        'section4': _read_section(ws, sections['s4']),
    }

    # Attach ISINs to section4 rows — ISIN lives in a separate standalone column
    # positioned at a variable row; read positionally (same order as section4 rows)
    isins = _read_isin_column(ws, len(result['section4']))
    for row, isin in zip(result['section4'], isins):
        row['ISIN'] = isin

    print(f"  Sections found: {list(result.keys())}")
    for k, v in result.items():
        non_gt = [r for r in v if not r.get('__grand_total__')]
        print(f"    {k}: {len(non_gt)} data rows")
    return result

# ──────────────────────────────────────────────────────────────
# DECK WRITER — writes data into a copy of the template PowerPoint
# ──────────────────────────────────────────────────────────────

# Risk-group name aliases — some client Excels use shorthand like 'S&G' instead of
# 'Gold & Silver'. Normalise at every point we read RISK_GROUP_L0.
_RG_NORMALIZE = {
    'S&G':          'Gold & Silver',
    'Gold':         'Gold & Silver',
    'Gold&Silver':  'Gold & Silver',
    'Debt like':    'Debt Like',
}

def _normalize_rg(rg):
    return _RG_NORMALIZE.get(rg, rg)


def _format_subcategory(sub):
    """Convert subcategory codes to readable text.
    Uses SUBCATEGORY_DISPLAY if available, otherwise converts
    'SOME_THING_NAME' → 'Some Thing Name'."""
    if sub in SUBCATEGORY_DISPLAY:
        return SUBCATEGORY_DISPLAY[sub]
    if not sub:
        return sub
    # Generic: replace underscores with spaces, title-case
    return sub.replace('_', ' ').title()


def _format_subcat_short(sub):
    """Short form of subcategory for label line 2.
    Uses SUBCAT_SHORT if available, otherwise falls back to _format_subcategory."""
    if sub in SUBCAT_SHORT:
        return SUBCAT_SHORT[sub]
    return _format_subcategory(sub)

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

# Template downloaded from Google Drive at runtime
_CACHE_DIR = tempfile.mkdtemp(prefix='m3_assets_')


def _get_template_path() -> str:
    """Download M3 Template from Drive (Google Slides → PPTX)."""
    local_path = os.path.join(_CACHE_DIR, 'M3_Template.pptx')
    if os.path.exists(local_path):
        return local_path

    from google_auth import get_drive_service
    from googleapiclient.http import MediaIoBaseDownload
    drive = get_drive_service()

    buf = BytesIO()
    request = drive.files().export_media(
        fileId=M3_TEMPLATE_ID,
        mimeType='application/vnd.openxmlformats-officedocument.presentationml.presentation',
    )
    downloader = MediaIoBaseDownload(buf, request)
    done = False
    while not done:
        _, done = downloader.next_chunk()
    buf.seek(0)

    with open(local_path, 'wb') as f:
        f.write(buf.read())
    print(f"  Downloaded M3 template from Google Drive")
    return local_path

SUBCATEGORY_DISPLAY = {
    'MID_CAP':                              'Mid Cap',
    'SMALL_CAP':                            'Small Cap',
    'FLEXI_CAP':                            'Flexi Cap',
    'VALUE_AND_CONTRA':                     'Value & Contra',
    'LARGE_CAP':                            'Large Cap',
    'DYNAMIC_ASSET_ALLOCATION':             'Dynamic Asset',
    'MULTI_ASSET_ALLOCATION':               'Multi Asset',
    'AGGRESSIVE_ALLOCATION':                'Aggressive Hybrid',
    'COMMODITY_GOLD':                       'Gold',
    'COMMODITY_SILVER':                     'Silver',
    'INDEX_LARGE_CAP':                      'Index Large Cap',
    'ELSS_TAX_SAVINGS':                     'ELSS',
    'LARGE_AND_MID_CAP':                    'Large & Mid',
    'FOCUSED_FUND':                         'Focused',
    'SECTORAL_THEMATIC_EXPORT_AND_SERVICES':'Sectoral Thematic',
    'INDEX_MID_CAP':                        'Index Mid Cap',
    'INDEX_SMALL_CAP':                      'Index Small Cap',
}

SUBCAT_SHORT = {
    'MID_CAP':                              'Mid',
    'SMALL_CAP':                            'Small',
    'FLEXI_CAP':                            'Flexi',
    'VALUE_AND_CONTRA':                     'Value & Contra',
    'LARGE_CAP':                            'Large Cap',
    'DYNAMIC_ASSET_ALLOCATION':             'Dynamic Asset',
    'MULTI_ASSET_ALLOCATION':               'Multi Asset',
    'COMMODITY_GOLD':                       'Gold',
    'COMMODITY_SILVER':                     'Silver',
    'INDEX_LARGE_CAP':                      'Index Large Cap',
    'ELSS_TAX_SAVINGS':                     'ELSS',
    'LARGE_AND_MID_CAP':                    'Large & Mid',
    'FOCUSED_FUND':                         'Focused',
    'AGGRESSIVE_ALLOCATION':                'Aggressive Hybrid',
    'SECTORAL_THEMATIC_EXPORT_AND_SERVICES':'Sectoral',
    'INDEX_MID_CAP':                        'Index Mid',
    'INDEX_SMALL_CAP':                      'Index Small',
}

RG_DISPLAY = {
    '1) Aggressive':   'Aggressive',
    '2) Balanced':     'Balanced',
    '3) Conservative': 'Conservative',
    'Hybrid':          'Hybrid',
    'Debt Like':       'Debt-Like',
    'Global':          'Global',
    'Gold & Silver':   'Gold',
    'Solution':        'Solution',
}

EQUITY_RGS = ['1) Aggressive', '2) Balanced', '3) Conservative']

# SIP template slide indices (0-based) per risk group
SIP_TEMPLATE = {
    '1) Aggressive':   4,
    '2) Balanced':     5,
    '3) Conservative': 6,
    'Hybrid':          7,
    'Gold & Silver':   8,
    # Global, Debt Like, Solution are corpus-only
}

# Corpus template slide indices (0-based)
CORPUS_TEMPLATE = {
    '1) Aggressive':   14,
    '2) Balanced':     15,
    '3) Conservative': 16,
    'Hybrid':          17,
    'Gold & Silver':   18,
    'Global':          19,
    'Debt Like':       20,
    'Solution':        21,
}

# Label line 1 per risk group
RG_LABEL1 = {
    '1) Aggressive':   'Equity - Aggressive',
    '2) Balanced':     'Equity - Balanced',
    '3) Conservative': 'Equity - Conservative',
    'Hybrid':          'Hybrid',
    'Gold & Silver':   'Gold',
    'Global':          'Global',
    'Debt Like':       'Debt Like',
    'Solution':        'Solution',
}

_MILESTONES = ['D0', 'D30', 'D60', 'D90', 'D120', 'D150']

_S3_SELL_COL = {
    'D0':   'Sum of Cumulative Sell Amount at D0',
    'D30':  'Sum of Cumulative Sell Amount at D30',
    'D60':  'Sum of Cumulative Sell Amount at D60',
    'D90':  'Sum of Cumulative Sell Amount at D90',
    'D120': 'Sum of Cumulative Sell Amount at D120',
    'D150': 'Sum of Cumulative Sell Amount at 150',
}

_S4_CUMM_BUY_COL = {
    'D0':   'Cummulative Buy Amount at D0',
    'D30':  'Cummulative Buy Amount at D30',
    'D60':  'Cummulative Buy Amount at D60',
    'D90':  'Cummulative Buy Amount at D90',
    'D120': 'Cummulative Buy Amount at D120',
    'D150': 'Cummulative Buy Amount at D150',
}

# New-format column names (M1-M6 milestones, direct ideal-allocation in s3)
# Section4 per-scheme ideal allocation at each milestone (new format).
# These are the actual values to SUM by risk group for the transition table.
_S4_ALLOC_COLS = {
    'D0':   'Allocation M1',
    'D30':  'Allocation M2',
    'D60':  'Allocation M3',
    'D90':  'Allocation M4',
    'D120': 'Allocation M5',
    'D150': 'Allocation M6',
}
_S4_CUMM_BUY_COL_NEW = {
    'D0':   'Cumm Buy Amount in M1',
    'D30':  'Cumm Buy Amount in M2',
    'D60':  'Cumm Buy Amount in M3',
    'D90':  'Cumm Buy Amount in M4',
    'D120': 'Cumm Buy Amount in M5',
    'D150': 'Cumm Buy Amount in M6',
}


def _get_val(row, *keys):
    """Return the first non-None value found under any of the given keys."""
    for k in keys:
        v = row.get(k)
        if v is not None:
            return v
    return None


def _detect_s4_schema(section4):
    """
    Inspect a section4 row to determine which column names are in use.
    Returns a dict with keys: corpus_pct, sip_amount, buy_value, cum_buy.
    """
    sample = next((r for r in section4 if not r.get('__grand_total__')), {})
    new = 'Cumm Buy Amount in M1' in sample
    return {
        'corpus_pct': 'Total Allocation % of PF' if new else 'Total Value as % of PF',
        'sip_amount':  'SIP Allocation Amount'    if new else 'SIP Amount',
        'buy_value':   'Buy Value Amount'          if new else 'Buy Value',
        'cum_buy': _S4_CUMM_BUY_COL_NEW if new else _S4_CUMM_BUY_COL,
    }


def _detect_s3_schema(section3):
    """
    Inspect a section3 row to determine old vs new format.
    Returns a dict with keys: current_val, new_format, sell_cols (old only).
    """
    sample = next((r for r in section3 if not r.get('__grand_total__')), {})
    if 'Sum of Current Value Amount' in sample:
        return {
            'current_val':  'Sum of Current Value Amount',
            'new_format':   True,
        }
    return {
        'current_val':  'Sum of TOTAL_VALUE',
        'new_format':   False,
        'sell_cols':    _S3_SELL_COL,
    }


# ---------------------------------------------------------------------------
# Utility helpers
# ---------------------------------------------------------------------------

def format_inr(value):
    if value is None:
        return "—"
    cr, L, K = 10_000_000, 100_000, 1_000
    if abs(value) >= cr:
        return f"{value/cr:.1f}Cr"
    elif abs(value) >= L:
        return f"{value/L:.1f}L"
    elif abs(value) >= K:
        n = round(value / K, 1)
        return f"{n:g}K"
    else:
        return f"{round(value)}"


def format_indian(n):
    """Format integer with Indian comma grouping: e.g. 150000 -> '1,50,000'"""
    n = int(round(n))
    negative = n < 0
    s = str(abs(n))
    if len(s) <= 3:
        result = s
    else:
        result = s[-3:]
        s = s[:-3]
        while s:
            result = s[-2:] + ',' + result
            s = s[:-2]
    result = result.lstrip(',')
    return ('-' if negative else '') + result


def set_cell_text(cell, text):
    tf = cell.text_frame
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = str(text)
        for r in para.runs[1:]:
            r.text = ""
    else:
        para.add_run().text = str(text)


def _set_para_text(para, text):
    if para.runs:
        para.runs[0].text = text
        for r in para.runs[1:]:
            r.text = ""
    else:
        para.add_run().text = text


def find_slide_by_text(prs, search_text):
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame and search_text in shape.text_frame.text:
                return i
    return None


def find_all_slides_by_text(prs, search_text):
    return [i for i, slide in enumerate(prs.slides)
            for shape in slide.shapes
            if shape.has_text_frame and search_text in shape.text_frame.text]


def replace_text_preserving_format(slide, old_text, new_text):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            full = "".join(r.text for r in para.runs)
            if old_text in full:
                new_full = full.replace(old_text, new_text)
                if para.runs:
                    para.runs[0].text = new_full
                    for r in para.runs[1:]:
                        r.text = ""


def delete_slide(prs, slide_idx):
    sldIdLst = prs.slides._sldIdLst
    sldId = sldIdLst[slide_idx]
    rId = sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId)
    sldIdLst.remove(sldId)


def add_table_row(table, copy_from_row_idx=1):
    """Deep-copy a row and insert it before the last row (total row)."""
    tbl = table._tbl
    rows = tbl.findall(qn('a:tr'))
    new_row = copy.deepcopy(rows[copy_from_row_idx])
    for tc in new_row.findall(qn('a:tc')):
        for r in tc.findall('.//' + qn('a:r')):
            t = r.find(qn('a:t'))
            if t is not None:
                t.text = ""
    rows[-1].addprevious(new_row)


def add_table_row_at_end(table, copy_from_row_idx=1):
    """Deep-copy a row and append it at the end."""
    tbl = table._tbl
    rows = tbl.findall(qn('a:tr'))
    new_row = copy.deepcopy(rows[copy_from_row_idx])
    for tc in new_row.findall(qn('a:tc')):
        for r in tc.findall('.//' + qn('a:r')):
            t = r.find(qn('a:t'))
            if t is not None:
                t.text = ""
    tbl.append(new_row)


def delete_table_row(table, row_idx):
    tbl = table._tbl
    rows = tbl.findall(qn('a:tr'))
    tbl.remove(rows[row_idx])


def _get_next_safe_slide_num(prs):
    """Return a slide number guaranteed not to conflict with existing slide parts."""
    import re as _re
    nums = []
    for s in prs.slides:
        m = _re.search(r'slide(\d+)', str(s.part.partname))
        if m:
            nums.append(int(m.group(1)))
    return (max(nums) + 1) if nums else 1


def duplicate_slide_after(prs, slide_idx):
    """Duplicate slide at slide_idx and insert copy immediately after it.
    Uses a safe partname to avoid conflicts when slides have been deleted.
    """
    from pptx.opc.packuri import PackURI
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    from pptx.parts.slide import SlidePart
    import lxml.etree as _etree

    src = prs.slides[slide_idx]
    safe_num = _get_next_safe_slide_num(prs)
    partname = PackURI(f'/ppt/slides/slide{safe_num}.xml')

    # Create new slide part with safe partname using the source slide's layout
    slide_layout_part = src.part.part_related_by(RT.SLIDE_LAYOUT)
    new_slide_part = SlidePart.new(partname, prs.part.package, slide_layout_part)
    new_slide = new_slide_part.slide

    # Replace spTree with deep copy of source
    sp_tree = new_slide.shapes._spTree
    for el in list(sp_tree):
        sp_tree.remove(el)
    for el in src.shapes._spTree:
        sp_tree.append(copy.deepcopy(el))

    # Add relationship from presentation part to new slide part
    rId = prs.part.relate_to(new_slide_part, RT.SLIDE)

    # Build new sldId element
    sldIdLst = prs.slides._sldIdLst
    existing_ids = [int(sld.get('id', 0)) for sld in sldIdLst]
    new_id = (max(existing_ids) + 1) if existing_ids else 256

    PML_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    R_NS = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    new_sld_id = _etree.SubElement(
        sldIdLst,
        f'{{{PML_NS}}}sldId',
    )
    new_sld_id.set('id', str(new_id))
    new_sld_id.set(f'{{{R_NS}}}id', rId)

    # Move from end of sldIdLst to slide_idx + 1
    sldIdLst.remove(new_sld_id)
    sldIdLst.insert(slide_idx + 1, new_sld_id)

    return slide_idx + 1


# ---------------------------------------------------------------------------
# Slide 1 & 2
# ---------------------------------------------------------------------------

def populate_slide1(slide, client_name):
    replace_text_preserving_format(slide, "Hari Vootori", client_name)


def populate_slide2(slide, client_name):
    first_name = client_name.split()[0]
    replace_text_preserving_format(slide, "Hari", first_name)
    today = date.today()
    today_str = f"{today.day} {today.strftime('%b')} {today.year}"
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            full = "".join(r.text for r in para.runs)
            if re.search(r"Data as of .+", full, re.IGNORECASE):
                new_full = re.sub(
                    r"(Data as of\s+).+",
                    lambda m: m.group(1) + today_str,
                    full,
                    flags=re.IGNORECASE,
                )
                if para.runs:
                    para.runs[0].text = new_full
                    for r in para.runs[1:]:
                        r.text = ""


# ---------------------------------------------------------------------------
# Slide 4 — SIP summary table (slide index 3 in template)
# ---------------------------------------------------------------------------

def _group_sip_schemes(section4):
    schema = _detect_s4_schema(section4)
    sip_amt_key = schema['sip_amount']
    groups = OrderedDict()
    for row in section4:
        if row.get('__grand_total__'):
            continue
        if not (_get_val(row, sip_amt_key, 'SIP Amount', 'SIP Allocation Amount') or 0):
            continue
        rg = row.get('RISK_GROUP_L0', '')
        sub = row.get('UPDATED_SUBCATEGORY', '')
        groups.setdefault(rg, OrderedDict()).setdefault(sub, []).append(row)
    return groups, sip_amt_key


def build_sip_rows(section4):
    groups, sip_amt_key = _group_sip_schemes(section4)
    rows = []

    def sip_amt(r):
        return _get_val(r, sip_amt_key, 'SIP Amount', 'SIP Allocation Amount') or 0

    equity_present = [rg for rg in EQUITY_RGS if rg in groups]
    if equity_present:
        eq_amount = sum(sip_amt(r) for rg in equity_present
                        for subs in groups[rg].values() for r in subs)
        eq_pct = sum(r.get('SIP Allocation %', 0) or 0
                     for rg in equity_present
                     for subs in groups[rg].values() for r in subs)
        eq_count = sum(len(subs) for rg in equity_present for subs in groups[rg].values())
        rows.append({'name': 'Equity', 'amount': eq_amount, 'pct': eq_pct,
                     'count': eq_count, 'level': 0})

        for rg in equity_present:
            rg_amount = sum(sip_amt(r) for subs in groups[rg].values() for r in subs)
            rg_pct = sum(r.get('SIP Allocation %', 0) or 0
                         for subs in groups[rg].values() for r in subs)
            rg_count = sum(len(subs) for subs in groups[rg].values())
            rows.append({'name': RG_DISPLAY.get(rg, rg),
                         'amount': rg_amount, 'pct': rg_pct,
                         'count': rg_count, 'level': 1})
            for sub, schemes in groups[rg].items():
                sub_amount = sum(sip_amt(r) for r in schemes)
                sub_pct = sum(r.get('SIP Allocation %', 0) or 0 for r in schemes)
                rows.append({'name': _format_subcategory(sub),
                             'amount': sub_amount, 'pct': sub_pct,
                             'count': len(schemes), 'level': 2})

    for rg, subs in groups.items():
        if rg in EQUITY_RGS:
            continue
        rg_amount = sum(sip_amt(r) for s in subs.values() for r in s)
        rg_pct = sum(r.get('SIP Allocation %', 0) or 0 for s in subs.values() for r in s)
        rg_count = sum(len(s) for s in subs.values())
        rows.append({'name': RG_DISPLAY.get(rg, rg),
                     'amount': rg_amount, 'pct': rg_pct,
                     'count': rg_count, 'level': 0})
        for sub, schemes in subs.items():
            sub_amount = sum(sip_amt(r) for r in schemes)
            sub_pct = sum(r.get('SIP Allocation %', 0) or 0 for r in schemes)
            rows.append({'name': _format_subcategory(sub),
                         'amount': sub_amount, 'pct': sub_pct,
                         'count': len(schemes), 'level': 1})

    return rows


def populate_slide4(prs, slide4_idx, section4):
    slide = prs.slides[slide4_idx]
    tbl_shape = None
    for shape in slide.shapes:
        if shape.has_table and shape.table.rows[0].cells[0].text_frame.text.strip() == 'Portfolio':
            tbl_shape = shape
            break
    if tbl_shape is None:
        print("  WARNING: SIP summary table not found")
        return

    main_table = tbl_shape.table
    desired = build_sip_rows(section4)
    current_data = len(main_table.rows) - 2
    target_data = len(desired)

    if target_data > current_data:
        for _ in range(target_data - current_data):
            add_table_row(main_table, copy_from_row_idx=2)
    elif target_data < current_data:
        for _ in range(current_data - target_data):
            delete_table_row(main_table, len(main_table.rows) - 2)

    # Resize rows to avoid overflow
    HEADER_H = 334375
    TOTAL_H = 100000
    DATA_H_ORIG = 218750
    TABLE_TOP = tbl_shape.top
    SAFE_BOTTOM = prs.slide_height - 300000
    available_for_data = SAFE_BOTTOM - TABLE_TOP - HEADER_H - TOTAL_H
    data_row_h = min(DATA_H_ORIG, available_for_data // max(target_data, 1))
    data_row_h = max(data_row_h, 150000)

    tbl = main_table._tbl
    tr_list = tbl.findall(qn('a:tr'))
    tr_list[0].set('h', str(HEADER_H))
    for tr in tr_list[1:-1]:
        tr.set('h', str(data_row_h))
    tr_list[-1].set('h', str(TOTAL_H))
    tbl_shape.height = HEADER_H + target_data * data_row_h + TOTAL_H

    _S4_INDENT = {0: '', 1: '   ', 2: '       '}
    for i, rd in enumerate(desired):
        row = main_table.rows[i + 1]
        indent = _S4_INDENT.get(rd.get('level', 0), '')
        pct_val = round(rd['pct'] * 100, 1)
        alloc = f"{format_inr(rd['amount'])} | {pct_val:g}%"
        set_cell_text(row.cells[0], indent + rd['name'])
        set_cell_text(row.cells[1], indent + alloc)
        set_cell_text(row.cells[2], indent + str(rd['count']))

    _sip_key = _detect_s4_schema(section4)['sip_amount']
    all_sip = [r for r in section4
               if not r.get('__grand_total__') and
               (_get_val(r, _sip_key, 'SIP Amount', 'SIP Allocation Amount') or 0) > 0]
    total_amount = sum(_get_val(r, _sip_key, 'SIP Amount', 'SIP Allocation Amount') or 0
                       for r in all_sip)
    total_pct = sum(r.get('SIP Allocation %', 0) or 0 for r in all_sip)
    total_count = len(all_sip)
    total_row = main_table.rows[len(main_table.rows) - 1]
    set_cell_text(total_row.cells[0], 'Total')
    set_cell_text(total_row.cells[1], f"{format_inr(total_amount)} | {round(total_pct * 100):g}%")
    set_cell_text(total_row.cells[2], str(total_count))

    # Update "Monthly SIP with Infinite: INR XX,XXX*" shape
    for shape in slide.shapes:
        if shape.has_text_frame and 'Monthly SIP with Infinite' in shape.text_frame.text:
            for para in shape.text_frame.paragraphs:
                full = "".join(r.text for r in para.runs)
                if 'Monthly SIP with Infinite' in full:
                    new_text = re.sub(r'INR [\d,]+', f'INR {format_indian(total_amount)}', full)
                    if para.runs:
                        para.runs[0].text = new_text
                        for r in para.runs[1:]:
                            r.text = ""
            break

    print(f"  Slide 4: SIP table {target_data} rows, monthly SIP = INR {format_indian(total_amount)}")


# ---------------------------------------------------------------------------
# Transition Plan (Slide 11 = index 10)
# ---------------------------------------------------------------------------

def _build_transition_data(section3, section4, section1=None):
    """
    Build the transition table data.  ALL data comes from section4:
      - CURRENT column  = sum of 'Current Value Amount' per risk group
      - IDEAL columns   = sum of 'Allocation M1'-'M6'  per risk group
    section1 is used only for the total corpus headline (Total Selected Value).
    section3 is kept as a fallback for older Excel formats that lack Allocation M1-M6.
    section2 is NOT used anywhere in this function.
    """
    s3_schema = _detect_s3_schema(section3)
    s4_schema = _detect_s4_schema(section4)
    cum_buy_cols = s4_schema['cum_buy']

    # --- total_pfv (denominator for %) ---
    # Priority: s1 Grand Total > sum of section4 Current Value Amount > s3 Grand Total
    gt1 = next((r for r in (section1 or []) if r.get('__grand_total__')), {})
    total_pfv = gt1.get('Total Selected Value') or 0

    # Check if section4 has Allocation M1-M6 (new format)
    sample_s4 = next((r for r in section4 if not r.get('__grand_total__')), {})
    has_alloc_cols = 'Allocation M1' in sample_s4

    if has_alloc_cols:
        # =================================================================
        # NEW FORMAT — everything from section4
        #   Current  = sum(Current Value Amount)       per RISK_GROUP_L0
        #   Ideal Dx = sum(Allocation Mx)              per RISK_GROUP_L0
        # =================================================================
        rg_alloc_s4   = {}   # {rg: {'D0': total, 'D30': total, ...}}
        rg_current_s4 = {}   # {rg: sum(Current Value Amount)}

        for row in section4:
            if row.get('__grand_total__'):
                continue
            rg = _normalize_rg(row.get('RISK_GROUP_L0', ''))
            if not rg:
                continue
            cur = row.get('Current Value Amount') or 0
            rg_current_s4[rg] = rg_current_s4.get(rg, 0.0) + cur
            for d, col in _S4_ALLOC_COLS.items():
                val = row.get(col) or 0
                rg_alloc_s4.setdefault(rg, {})
                rg_alloc_s4[rg][d] = rg_alloc_s4[rg].get(d, 0.0) + val

        if not total_pfv:
            total_pfv = sum(rg_current_s4.values())

        def current_for_rg(rg):
            return rg_current_s4.get(rg, 0)

        def ideal_at(rg, d):
            return rg_alloc_s4.get(rg, {}).get(d, 0)

        all_active = set(rg_alloc_s4.keys()) | set(rg_current_s4.keys())
        # Filter out SIP-only RGs that have no corpus allocation (they show as all-zero rows)
        all_active = {rg for rg in all_active
                      if current_for_rg(rg) > 0 or any(ideal_at(rg, d) > 0 for d in _MILESTONES)}

        print(f"  Transition (new format): {len(rg_alloc_s4)} RGs from s4 alloc, "
              f"{len(rg_current_s4)} RGs from s4 current, total_pfv={total_pfv}")

    else:
        # =================================================================
        # OLD FORMAT: s3 pivot (cur - cum_sell) + s4 cumulative buys
        # =================================================================
        if not total_pfv:
            gt3 = next((r for r in section3 if r.get('__grand_total__')), {})
            total_pfv = gt3.get(s3_schema['current_val']) or 0

        known_rgs = set(SIP_TEMPLATE.keys()) | set(CORPUS_TEMPLATE.keys())
        known_rgs.update(
            _normalize_rg(row.get('RISK_GROUP_L0', ''))
            for row in section4
            if not row.get('__grand_total__') and row.get('RISK_GROUP_L0')
        )

        s3_rg = {}
        for r in section3:
            if r.get('__grand_total__'):
                continue
            label = r.get('Row Labels', '')
            norm = _normalize_rg(label)
            if label in known_rgs or norm in known_rgs:
                s3_rg[norm] = r

        rg_cumm_buy = {}
        for row in section4:
            if row.get('__grand_total__'):
                continue
            rg = _normalize_rg(row.get('RISK_GROUP_L0', ''))
            if rg not in rg_cumm_buy:
                rg_cumm_buy[rg] = {col: 0.0 for col in cum_buy_cols.values()}
            for d, col in cum_buy_cols.items():
                rg_cumm_buy[rg][col] += (row.get(col) or 0)

        cur_col = s3_schema['current_val']

        def current_for_rg(rg):
            return (s3_rg[rg].get(cur_col) or 0) if rg in s3_rg else 0

        def ideal_at(rg, d):
            cum_buy = rg_cumm_buy.get(rg, {}).get(cum_buy_cols[d], 0)
            cur = (s3_rg[rg].get(cur_col) or 0) if rg in s3_rg else 0
            cum_sell = (s3_rg[rg].get(s3_schema['sell_cols'].get(d, ''), 0) or 0) if rg in s3_rg else 0
            return cur - cum_sell + cum_buy

        all_active = set(s3_rg.keys()) | set(rg_cumm_buy.keys())
        # Filter out SIP-only RGs that have no corpus allocation (they show as all-zero rows)
        all_active = {rg for rg in all_active
                      if current_for_rg(rg) > 0 or any(ideal_at(rg, d) > 0 for d in _MILESTONES)}

        print(f"  Transition (old format): {len(s3_rg)} RGs from s3, "
              f"{len(rg_cumm_buy)} RGs from s4 buys, total_pfv={total_pfv}")

    # ----- Build output rows (shared by both formats) -----
    def cell(val):
        pct = round(val / total_pfv * 100) if total_pfv else 0
        return (val, pct)

    equity_rgs = [rg for rg in EQUITY_RGS if rg in all_active]
    non_equity = sorted(
        [rg for rg in all_active if rg not in set(EQUITY_RGS)],
        key=lambda rg: CORPUS_TEMPLATE.get(rg, 99),
    )

    rows = []

    if equity_rgs:
        eq_cur = sum(current_for_rg(rg) for rg in equity_rgs)
        eq_row = {'name': 'Equity', 'level': 1, 'current': cell(eq_cur)}
        for d in _MILESTONES:
            eq_row[d] = cell(sum(ideal_at(rg, d) for rg in equity_rgs))
        rows.append(eq_row)
        for rg in equity_rgs:
            r = {'name': RG_DISPLAY.get(rg, rg), 'level': 2, 'current': cell(current_for_rg(rg))}
            for d in _MILESTONES:
                r[d] = cell(ideal_at(rg, d))
            rows.append(r)

    for rg in non_equity:
        r = {'name': RG_DISPLAY.get(rg, rg), 'level': 1, 'current': cell(current_for_rg(rg))}
        for d in _MILESTONES:
            r[d] = cell(ideal_at(rg, d))
        rows.append(r)

    total_cur = sum(current_for_rg(rg) for rg in all_active)
    total_row = {'name': 'Total', 'level': 0, 'current': cell(total_cur)}
    for d in _MILESTONES:
        total_row[d] = cell(sum(ideal_at(rg, d) for rg in all_active))
    rows.append(total_row)

    return rows, total_pfv


def _trim_repeated_milestones(rows_data):
    """
    Detect the last milestone where allocation values still change.
    Once ALL risk groups have the same values as the previous milestone,
    the transition is over. Returns the trimmed list of milestone keys
    and the number of months for the title.

    Milestone mapping: D0=0mo, D30=1mo, D60=2mo, D90=3mo, D120=4mo, D150=5mo.
    """
    active = _MILESTONES[:1]  # D0 is always shown
    for i in range(1, len(_MILESTONES)):
        prev_d = _MILESTONES[i - 1]
        cur_d = _MILESTONES[i]
        # Check if ANY row changed between prev and cur milestone
        changed = False
        for rd in rows_data:
            if rd[cur_d][0] != rd[prev_d][0]:
                changed = True
                break
        if changed:
            active.append(cur_d)
        else:
            break  # all subsequent milestones will also repeat
    # Months: D0=0, D30=1, D60=2, D90=3, D120=4, D150=5
    _D_TO_MONTHS = {'D0': 0, 'D30': 1, 'D60': 2, 'D90': 3, 'D120': 4, 'D150': 5}
    last_months = _D_TO_MONTHS.get(active[-1], 5)
    return active, last_months


def populate_slide10(prs, slide10_idx, section1, section3, section4):
    slide = prs.slides[slide10_idx]
    rows_data, total_pfv = _build_transition_data(section3, section4, section1)

    # Determine which milestones to show (trim repeating tail)
    active_milestones, last_month = _trim_repeated_milestones(rows_data)
    print(f"  Transition: active milestones = {active_milestones}, last_month = {last_month}")

    print("  Transition table computed values:")
    for rd in rows_data:
        vals = ' '.join(f"{d}={format_inr(rd[d][0])}/{rd[d][1]}%" for d in active_milestones)
        print(f"    {rd['name']:15s} cur={format_inr(rd['current'][0])}/{rd['current'][1]}%  {vals}")

    # Update title shape — replace corpus amount and month count
    for shape in slide.shapes:
        if shape.has_text_frame and 'transition plan for' in shape.text_frame.text.lower():
            for para in shape.text_frame.paragraphs:
                full = "".join(r.text for r in para.runs)
                if 'transition plan for' in full.lower():
                    new_text = re.sub(
                        r'(The transition plan for the ).*?(corpus)',
                        rf'\g<1>{format_inr(total_pfv)} \g<2>',
                        full,
                        flags=re.IGNORECASE,
                    )
                    # Update "ideal allocation in X months"
                    new_text = re.sub(
                        r'(ideal allocation in\s*)\d+(\s*months?)',
                        rf'\g<1>{last_month}\g<2>',
                        new_text,
                        flags=re.IGNORECASE,
                    )
                    if para.runs:
                        para.runs[0].text = new_text
                        for r in para.runs[1:]:
                            r.text = ""
            break

    # Update "ideal allocation in X months" — may be in a separate shape
    for shape in slide.shapes:
        if shape.has_text_frame and 'ideal allocation in' in shape.text_frame.text.lower():
            for para in shape.text_frame.paragraphs:
                full = "".join(r.text for r in para.runs)
                if 'ideal allocation in' in full.lower():
                    new_full = re.sub(
                        r'(ideal allocation in\s*)\d+(\s*months?)',
                        rf'\g<1>{last_month}\g<2>',
                        full,
                        flags=re.IGNORECASE,
                    )
                    if new_full != full and para.runs:
                        para.runs[0].text = new_full
                        for r in para.runs[1:]:
                            r.text = ""
            break

    # Find the table
    tbl_shape = next((s for s in slide.shapes if s.has_table), None)
    if tbl_shape is None:
        print("  WARNING: transition table not found")
        return
    tbl = tbl_shape.table

    # Remove columns for milestones that are trimmed (from right to left)
    # Table columns: 0=name, 1=Current, 2=D0, 3=D30, 4=D60, 5=D90, 6=D120, 7=D150
    cols_to_keep = 2 + len(active_milestones)  # 1 (name) + 1 (current) + active milestones
    total_cols = len(tbl.columns)
    if cols_to_keep < total_cols:
        tbl_xml = tbl._tbl
        grid_cols = tbl_xml.findall(qn('a:tblGrid') + '/' + qn('a:gridCol'))
        for col_idx in range(total_cols - 1, cols_to_keep - 1, -1):
            # Remove this column from every row
            for tr in tbl_xml.findall(qn('a:tr')):
                tcs = tr.findall(qn('a:tc'))
                if col_idx < len(tcs):
                    tr.remove(tcs[col_idx])
            # Remove from grid
            if col_idx < len(grid_cols):
                grid_cols[col_idx].getparent().remove(grid_cols[col_idx])
        print(f"  Trimmed {total_cols - cols_to_keep} repeated milestone column(s)")

    # Update header row milestone labels
    header_row = tbl.rows[0]
    header_labels = ['Portfolio Allocation', 'Current'] + [f'Ideal - {d}' for d in active_milestones]
    for j, lbl in enumerate(header_labels):
        if j < len(header_row.cells):
            set_cell_text(header_row.cells[j], lbl)

    current_data = len(tbl.rows) - 2
    target_data = len(rows_data) - 1   # exclude Total row

    if target_data > current_data:
        for _ in range(target_data - current_data):
            add_table_row(tbl, copy_from_row_idx=2)
    elif target_data < current_data:
        for _ in range(current_data - target_data):
            delete_table_row(tbl, len(tbl.rows) - 2)

    _S9_INDENT = {0: '', 1: '  ', 2: '      '}
    col_order = ['current'] + active_milestones

    for i, rd in enumerate(rows_data[:-1]):
        row = tbl.rows[i + 1]
        indent = _S9_INDENT.get(rd.get('level', 0), '')
        set_cell_text(row.cells[0], indent + rd['name'])
        for j, key in enumerate(col_order):
            val, pct = rd[key]
            set_cell_text(row.cells[j + 1], indent + f"{format_inr(val)} | {pct}%")

    total_rd = rows_data[-1]
    total_row = tbl.rows[len(tbl.rows) - 1]
    set_cell_text(total_row.cells[0], '  Total')
    for j, key in enumerate(col_order):
        val, pct = total_rd[key]
        set_cell_text(total_row.cells[j + 1], f"  {format_inr(val)} | {pct}%")

    # Tax liability — summed from section4 per-scheme rows.
    # Section4 has EXIT_LOAD, EL+STCG, Max LTCG per scheme (no grand total row).
    # Total tax impact = sum(EL+STCG) + sum(Max LTCG) across all section4 schemes.
    s4_data = [r for r in section4 if not r.get('__grand_total__')]
    el        = sum((r.get('EXIT_LOAD') or r.get('Exit Load') or 0) for r in s4_data)
    el_stcg   = sum((r.get('EL+STCG') or r.get('Total Impact (EL+STCG)') or 0) for r in s4_data)
    if not el_stcg:
        el_stcg = el   # fallback: if no EL+STCG column, use just EXIT_LOAD
    stcg      = max(0.0, el_stcg - el)
    ltcg_gain = sum((r.get('Max LTCG') or r.get('LTCG') or 0) for r in s4_data)
    tax_total = el_stcg + ltcg_gain

    def _fmt_tax(v):
        """Format values in lakhs with exactly 2 decimal places."""
        if v is None:
            return "—"
        L = 100_000
        if abs(v) >= L:
            return f"{v/L:.2f}L"
        K = 1_000
        if abs(v) >= K:
            return f"{v/K:.2f}K"
        return f"{v:.2f}"

    for shape in slide.shapes:
        if shape.has_text_frame and 'Tax Liability' in shape.text_frame.text:
            paras = shape.text_frame.paragraphs
            if len(paras) >= 1:
                _set_para_text(paras[0], f"Tax Liability: INR {_fmt_tax(tax_total)}")
            if len(paras) >= 2:
                _set_para_text(paras[1], f"EL {_fmt_tax(el)} + STCG {_fmt_tax(stcg)} + LTCG {_fmt_tax(ltcg_gain)}")
            break

    print(f"  Slide 10: transition plan updated, corpus = {format_inr(total_pfv)}")


# ---------------------------------------------------------------------------
# Scheme slides (SIP + Corpus) — slot-based approach
# ---------------------------------------------------------------------------

def _get_slot_groups(slide):
    """
    Group all table shapes on a slide into scheme-row slots.

    Uses a two-phase approach:
    1. Round top to nearest 100,000 EMU (~2.8 mm) so shapes in the same row but
       at slightly different Y positions (common in Google Slides exports) are
       grouped together.
    2. Merge any groups that ended up with fewer than 3 shapes into the nearest
       larger group — this handles the rare case where a single stray shape gets
       its own bucket.

    Returns list of slots sorted by top position.  Each slot is a list of
    shapes sorted by left (name col first, % col second, …).
    """
    tbl_shapes = [s for s in slide.shapes if s.has_table]
    if not tbl_shapes:
        return []

    # Phase 1: bucket by rounded top (100k EMU tolerance)
    slot_dict = {}
    for s in tbl_shapes:
        key = round(s.top / 100_000) * 100_000
        slot_dict.setdefault(key, []).append(s)

    # Phase 2: merge lone/small groups (< 3 shapes) into closest neighbour
    keys = sorted(slot_dict.keys())
    merged = {}
    for key in keys:
        group = slot_dict[key]
        if len(group) >= 3 or not merged:
            merged[key] = group
        else:
            # Find the nearest existing bucket
            nearest = min(merged.keys(), key=lambda k: abs(k - key))
            merged[nearest].extend(group)

    result = []
    for key in sorted(merged.keys()):
        slot = sorted(merged[key], key=lambda x: x.left)
        result.append(slot)
    return result


def _clone_slot(slide, source_slot, y_offset_emu):
    """
    Deep-copy all shapes in source_slot, shift their top by y_offset_emu,
    append to the slide spTree, and return the new slot (list of shapes).
    """
    NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    spTree = slide.shapes._spTree
    new_shapes = []
    for shape in source_slot:
        new_elem = copy.deepcopy(shape._element)
        # Update xfrm/off @y inside the table wrapper
        for off in new_elem.iter(f'{{{NS_A}}}off'):
            off.set('y', str(int(off.get('y', 0)) + y_offset_emu))
        spTree.append(new_elem)
        # Wrap in a proxy shape so the rest of the code can use .left / .top
        # We use the original shape type; append then re-fetch last shape
        new_shapes.append(slide.shapes[-1])
    return sorted(new_shapes, key=lambda s: s.left)


def _delete_slot(slot):
    """Remove all shapes in a slot from the slide's spTree."""
    for shape in slot:
        sp = shape._element
        sp.getparent().remove(sp)


def set_five_year_cell(cell, five_year, subcategory, rr_category):
    """
    Write 5Y return as two runs with different font sizes:
      Run 1: 'XX% '  at 9pt (114300 EMU)
      Run 2: '(+Y%)' at 7pt (88900 EMU)
    If data missing, write a single dash and clear the second run.
    """
    tf = cell.text_frame
    para = tf.paragraphs[0]

    if five_year is None:
        if para.runs:
            para.runs[0].text = '-'
            for r in para.runs[1:]:
                r.text = ''
        else:
            para.add_run().text = '-'
        return

    run1_text = f"{round(five_year)}% "

    cat_ret = rr_category.get(subcategory) if subcategory else None
    if cat_ret is not None:
        alpha = round(five_year - cat_ret)
        sign = '+' if alpha >= 0 else ''
        run2_text = f"({sign}{alpha}%)"
    else:
        run2_text = ''

    # Ensure exactly 2 runs (preserves existing font colour, bold, etc.)
    while len(para.runs) < 2:
        para.add_run()
    while len(para.runs) > 2:
        para.runs[-1].text = ''

    para.runs[0].text = run1_text
    para.runs[0].font.size = 114300  # 9pt

    para.runs[1].text = run2_text
    para.runs[1].font.size = 88900   # 7pt


def _fill_slot(slot, scheme, pct_key, ref_data, rr_category):
    """
    Write all 6 columns for a scheme slot (shapes sorted left-to-right):
      0: Fund name
      1: Allocation %
      2: Upside | Downside
      3: Power Rank
      4: 5Y Rolling Return
      5: AUM
    """
    _NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'

    def set_table_text(shape, text):
        t_elems = shape._element.findall(f'.//{{{_NS_A}}}t')
        if t_elems:
            t_elems[0].text = str(text)
            for t in t_elems[1:]:
                t.text = ''

    isin = scheme.get('ISIN')
    fund_name_key = (scheme.get('FUND_NAME') or '').strip()

    if isin:
        # Use ISIN directly — it is the authoritative identifier.
        # Name mismatches (typos, missing words) are common in client files; ignore them.
        isin_ref_name = ref_data.get('_isin_to_name', {}).get(isin, '')
        if isin_ref_name and _norm_name(isin_ref_name) != _norm_name(fund_name_key):
            print(f"    ISIN used (name variant): {fund_name_key!r} -> {isin_ref_name!r}")

    if not isin:
        # Resolve by fund name (exact, then normalized)
        isin = ref_data.get('_name_to_isin', {}).get(fund_name_key)
        if not isin:
            isin = ref_data.get('_norm_name_to_isin', {}).get(_norm_name(fund_name_key))
        if isin:
            print(f"    ISIN resolved via name match: {fund_name_key!r} -> {isin}")
        else:
            print(f"    WARNING: no ISIN match for {fund_name_key!r}")
    info = ref_data.get(isin, {}) if isin else {}
    subcategory = scheme.get('UPDATED_SUBCATEGORY')

    fund_name = scheme.get('FUND_NAME', '')
    pct_val = (scheme.get(pct_key, 0) or 0) * 100
    pct_str = f"{round(pct_val, 1):g}%"

    # Debug: 5Y rolling return calculation
    five_year = info.get('five_year')
    cat_ret = rr_category.get(subcategory) if subcategory else None
    alpha = round(five_year - cat_ret) if (five_year is not None and cat_ret is not None) else None
    if five_year is not None:
        sign = '+' if (alpha or 0) >= 0 else ''
        display = f"{round(five_year)}% ({sign}{alpha}%)" if alpha is not None else f"{round(five_year)}%"
    else:
        display = '-'
    print(f"    5Y [{isin}] {subcategory}: scheme={five_year}, cat_avg={cat_ret}, alpha={alpha} -> {display!r}")

    shapes_sorted = sorted(slot, key=lambda s: s.left)

    if len(shapes_sorted) >= 1:
        set_table_text(shapes_sorted[0], fund_name)
    if len(shapes_sorted) >= 2:
        set_table_text(shapes_sorted[1], pct_str)
    if len(shapes_sorted) >= 3:
        set_table_text(shapes_sorted[2],
                       fmt_upside_downside(info.get('upside'), info.get('downside')))
    if len(shapes_sorted) >= 4:
        set_table_text(shapes_sorted[3], fmt_powerrank(info.get('powerrank')))
    if len(shapes_sorted) >= 5:
        set_five_year_cell(shapes_sorted[4].table.rows[0].cells[0],
                           five_year, subcategory, rr_category)
    if len(shapes_sorted) >= 6:
        set_table_text(shapes_sorted[5], fmt_aum(info.get('aum_cr')))


def _find_label_shape(slide):
    """
    Find the category label shape: has_text_frame, contains ' | ' and '%'.
    The label shape uses vertical tab (\\x0b) as line separator.
    """
    for shape in slide.shapes:
        if shape.has_text_frame and not shape.has_table:
            txt = shape.text_frame.text
            if ' | ' in txt and '%' in txt:
                return shape
    return None


def _update_label_shape(shape, rg, schemes, pct_key):
    """Update the 2-line label shape for a risk group slide.

    The template uses two layouts:
      A) Two separate paragraphs — Aggressive/Balanced/Conservative slides.
      B) Single paragraph with <a:br/> line-break — Hybrid/Gold/etc slides.
    Both are handled here.
    """
    label1 = RG_LABEL1.get(rg, rg)
    total_pct = sum((r.get(pct_key, 0) or 0) for r in schemes) * 100
    total_pct_str = f"{round(total_pct, 1):g}"

    seen = []
    for r in schemes:
        sub = r.get('UPDATED_SUBCATEGORY', '')
        short = _format_subcat_short(sub)
        if short not in seen:
            seen.append(short)
    line2 = ' | '.join(seen)

    NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    tf = shape.text_frame
    paras = tf.paragraphs

    if len(paras) >= 2:
        # Layout A: two separate paragraphs
        _set_para_text(paras[0], f"{label1} | {total_pct_str}%")
        _set_para_text(paras[1], line2)
        # Remove any left-margin offset on paragraph 0 that causes the two lines
        # to center at different horizontal positions (making the label look crooked)
        pPr0 = paras[0]._p.find(f'{{{NS_A}}}pPr')
        if pPr0 is not None and pPr0.get('marL'):
            pPr0.attrib.pop('marL', None)
            pPr0.attrib.pop('indent', None)
    elif len(paras) == 1:
        # Layout B: single paragraph, possibly with <a:br/> separating line1/line2
        para_xml = paras[0]._p
        br_elem = para_xml.find(f'{{{NS_A}}}br')
        if br_elem is not None:
            # Collect <a:t> elements before and after the <a:br/>
            runs_before, runs_after, found_br = [], [], False
            for child in para_xml:
                tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
                if tag == 'br':
                    found_br = True
                elif tag == 'r':
                    t = child.find(f'{{{NS_A}}}t')
                    if t is not None:
                        (runs_after if found_br else runs_before).append(t)
            if runs_before:
                runs_before[0].text = f"{label1} | {total_pct_str}%"
                for t in runs_before[1:]:
                    t.text = ''
            if runs_after:
                runs_after[0].text = line2
                for t in runs_after[1:]:
                    t.text = ''
        else:
            # No line break at all — just update line 1
            _set_para_text(paras[0], f"{label1} | {total_pct_str}%")


def _find_scheme_slide_for_rg(prs, rg, is_sip):
    """
    Find the slide index for a given risk group.
    For SIP: find within slides before the appendix separator.
    For Corpus: find within slides after the appendix separator.
    Uses the RG_LABEL1 text to identify the slide.
    """
    label1 = RG_LABEL1.get(rg, rg)
    sep_idx = find_slide_by_text(prs, "List of schemes for ideal portfolio")

    if is_sip:
        # Search before separator
        limit = sep_idx if sep_idx is not None else len(prs.slides)
        for i in range(0, limit):
            s = prs.slides[i]
            # Must have at least one table shape (slot) to be a scheme slide
            if not any(sh.has_table for sh in s.shapes):
                continue
            for sh in s.shapes:
                if sh.has_text_frame and label1 in sh.text_frame.text:
                    return i
    else:
        # Search after separator
        if sep_idx is None:
            return None
        for i in range(sep_idx + 1, len(prs.slides)):
            s = prs.slides[i]
            if not any(sh.has_table for sh in s.shapes):
                continue
            for sh in s.shapes:
                if sh.has_text_frame and label1 in sh.text_frame.text:
                    return i
    return None


def _build_scheme_slides(prs, groups, template_map, pct_key, ref_data, rr_category):
    """
    Build scheme slides for SIP or corpus.
    template_map: dict of {risk_group: original_0based_slide_idx_in_template}
    pct_key: 'SIP Allocation %' or 'Total Value as % of PF'

    For each RG in template_map:
    - If no schemes: delete the template slide
    - If schemes: fill slots (up to 4 per slide), clone for overflow
    """
    is_sip = (pct_key == 'SIP Allocation %')
    rg_order = sorted(template_map.keys(), key=lambda rg: template_map[rg])

    # Process deletions last, after all fills (to keep indices stable during fills)
    rgs_to_delete = []
    rgs_to_fill = []

    for rg in rg_order:
        schemes = []
        if rg in groups:
            schemes = [r for sub in groups[rg].values() for r in sub]
        if not schemes:
            rgs_to_delete.append(rg)
        else:
            rgs_to_fill.append((rg, schemes))

    # Fill slides (process in order; cloning shifts indices so re-find each time)
    for rg, schemes in rgs_to_fill:
        slide_idx = _find_scheme_slide_for_rg(prs, rg, is_sip)
        if slide_idx is None:
            print(f"  WARNING: template slide for {rg!r} not found, skipping")
            continue

        # 4 schemes per slide; cloning handles slides that have fewer template slots
        SLOTS_PER_SLIDE = 4
        pages = [schemes[i:i + SLOTS_PER_SLIDE]
                 for i in range(0, len(schemes), SLOTS_PER_SLIDE)]

        current_slide_idx = slide_idx
        for page_num, page_schemes in enumerate(pages):
            if page_num > 0:
                current_slide_idx = duplicate_slide_after(prs, current_slide_idx)

            curr_slide = prs.slides[current_slide_idx]
            slots = _get_slot_groups(curr_slide)

            # If page has more schemes than template slots, clone extra slots.
            # This happens when the template slot count is less than what we derived
            # (e.g. due to index shift) or when a template slide has fewer slots than
            # expected. The while loop handles any number of extra schemes.
            if len(page_schemes) > len(slots) and slots:
                spacing = (slots[1][0].top - slots[0][0].top) if len(slots) >= 2 else 614062
                while len(slots) < len(page_schemes):
                    new_slot = _clone_slot(curr_slide, slots[-1], spacing)
                    slots.append(new_slot)

            # Fill used slots
            for slot_i, scheme in enumerate(page_schemes):
                if slot_i < len(slots):
                    _fill_slot(slots[slot_i], scheme, pct_key, ref_data, rr_category)

            # Delete unused slots (from bottom up)
            for slot_i in range(len(slots) - 1, len(page_schemes) - 1, -1):
                _delete_slot(slots[slot_i])

            # Update label shape on every page
            label_shape = _find_label_shape(curr_slide)
            if label_shape:
                _update_label_shape(label_shape, rg, schemes, pct_key)

            # Fix slide text for corpus slides
            if not is_sip:
                for sh in curr_slide.shapes:
                    if sh.has_text_frame and not sh.has_table:
                        tf = sh.text_frame
                        paras = tf.paragraphs
                        for pi, para in enumerate(paras):
                            for run in para.runs:
                                if 'SIP strategy' in run.text:
                                    run.text = run.text.replace('SIP strategy', 'Corpus strategy')
                                # Corpus slides: "% Monthly Allocation" or "% Monthly" → "% Allocation"
                                # Handle single-run case first to avoid double "Allocation"
                                if '% Monthly Allocation' in run.text:
                                    run.text = run.text.replace('% Monthly Allocation', '% Allocation')
                                elif '% Monthly' in run.text:
                                    run.text = run.text.replace('% Monthly', '% Allocation')
                                    # Remove the next paragraph entirely if it's just "Allocation"
                                    # (header split across two paras — deleting prevents blank line)
                                    if pi + 1 < len(paras) and paras[pi + 1].text.strip().lower() == 'allocation':
                                        orphan = paras[pi + 1]._p
                                        orphan.getparent().remove(orphan)

        print(f"  Filled scheme slide for {rg!r}: {len(schemes)} schemes, {len(pages)} page(s)")

    # Delete unused template slides in reverse index order
    slides_to_delete_idx = []
    for rg in rgs_to_delete:
        idx = _find_scheme_slide_for_rg(prs, rg, is_sip)
        if idx is not None:
            slides_to_delete_idx.append(idx)

    for idx in sorted(slides_to_delete_idx, reverse=True):
        delete_slide(prs, idx)
        print(f"  Deleted unused scheme slide at index {idx}")


def populate_sip_scheme_slides(prs, section4, ref_data, rr_category):
    schema = _detect_s4_schema(section4)
    sip_amt_key = schema['sip_amount']

    groups = OrderedDict()
    for row in section4:
        if row.get('__grand_total__'):
            continue
        if not (_get_val(row, sip_amt_key, 'SIP Amount', 'SIP Allocation Amount') or 0):
            continue
        rg = _normalize_rg(row.get('RISK_GROUP_L0', ''))
        sub = row.get('UPDATED_SUBCATEGORY', '')
        groups.setdefault(rg, OrderedDict()).setdefault(sub, []).append(row)

    # Merge multi-folio rows: same fund → one slot, summed SIP amount and allocation %
    for rg in groups:
        for sub in groups[rg]:
            groups[rg][sub] = _dedup_schemes_for_slide(
                groups[rg][sub], list(dict.fromkeys([sip_amt_key, 'SIP Allocation %', 'SIP Amount']))
            )

    _build_scheme_slides(prs, groups, SIP_TEMPLATE, 'SIP Allocation %', ref_data, rr_category)
    print("  SIP scheme slides populated")


def populate_corpus_scheme_slides(prs, section4, ref_data, rr_category):
    schema = _detect_s4_schema(section4)
    corpus_pct_key = schema['corpus_pct']

    groups = OrderedDict()
    for row in section4:
        if row.get('__grand_total__'):
            continue
        if not (_get_val(row, corpus_pct_key, 'Total Value as % of PF', 'Total Allocation % of PF') or 0):
            continue
        rg = _normalize_rg(row.get('RISK_GROUP_L0', ''))
        sub = row.get('UPDATED_SUBCATEGORY', '')
        groups.setdefault(rg, OrderedDict()).setdefault(sub, []).append(row)

    # Merge multi-folio rows: same fund (even with different actions) → one slot, summed allocation %
    for rg in groups:
        for sub in groups[rg]:
            groups[rg][sub] = _dedup_schemes_for_slide(
                groups[rg][sub],
                list(dict.fromkeys([corpus_pct_key, 'Total Value as % of PF', 'Total Allocation % of PF']))
            )

    _build_scheme_slides(prs, groups, CORPUS_TEMPLATE, corpus_pct_key, ref_data, rr_category)
    print("  Corpus scheme slides populated")


# ---------------------------------------------------------------------------
# Sell / Retain / Buy slides
# ---------------------------------------------------------------------------

MAX_ROWS_PER_SLIDE = 15



def _fill_action_table(table, schemes, name_key, value_key):
    """Populate an action table. Header is row 0; data starts at row 1."""
    current_data = len(table.rows) - 1
    target = len(schemes)

    if target > current_data:
        for _ in range(target - current_data):
            add_table_row_at_end(table, copy_from_row_idx=1)
    elif target < current_data:
        for _ in range(current_data - target):
            delete_table_row(table, len(table.rows) - 1)

    for i, scheme in enumerate(schemes):
        row = table.rows[i + 1]
        set_cell_text(row.cells[0], scheme.get(name_key, '') or '')
        # Try value_key and a common "Amount" variant (handles both old/new Excel formats)
        val = _get_val(scheme, value_key, value_key + ' Amount', value_key.replace(' Amount', '')) or 0
        set_cell_text(row.cells[1], format_inr(val))
        if len(row.cells) > 2:
            set_cell_text(row.cells[2], scheme.get('Reason', '') or '')


def _remove_section_from_slide(slide, label_text, header_text):
    """
    Remove the label shape and its associated table from the slide's spTree.

    Strategy (in order):
    1. Find the label shape by label_text; remove it.
    2. If label found  → remove the nearest table below the label (position-based,
       header_text used as a tiebreak preference only).
    3. If no label     → remove the first table whose header row contains header_text.
    """
    spTree = slide.shapes._spTree

    # Step 1: find and remove label shape
    label_shape = None
    label_bottom = None
    for shape in list(slide.shapes):
        if shape.has_text_frame and not shape.has_table:
            if label_text in shape.text_frame.text:
                label_shape = shape
                label_bottom = shape.top + shape.height
                break

    if label_shape is not None:
        spTree.remove(label_shape._element)

    # Step 2: remove the associated table
    if label_bottom is not None:
        # Position-based: pick the table whose top is closest to just below the label.
        # Prefer header_text match; fall back to closest table regardless of header.
        table_shapes = [s for s in slide.shapes if s.has_table]
        if table_shapes:
            # Try header-filtered first (only when header_text given)
            if header_text:
                preferred = [s for s in table_shapes
                             if header_text in ' '.join(c.text_frame.text
                                                        for c in s.table.rows[0].cells)]
            else:
                preferred = []
            pool = preferred if preferred else table_shapes
            best = min(pool, key=lambda s: abs(s.top - label_bottom))
            spTree.remove(best._element)
    else:
        # No label found — remove first table matching header_text
        if header_text:
            for shape in list(slide.shapes):
                if shape.has_table:
                    row0 = ' '.join(c.text_frame.text for c in shape.table.rows[0].cells)
                    if header_text in row0:
                        spTree.remove(shape._element)
                        break



def _find_table_near_label(slide, primary_label, header_text):
    """Find the table whose first-row text contains header_text,
    optionally closest to a label shape containing primary_label.

    If primary_label is given and no table matches header_text, falls back to
    the table physically nearest to the label (so mismatched template headers
    never cause data to be silently dropped).
    """
    all_tables = [s for s in slide.shapes if s.has_table]
    candidates = []
    if header_text:
        for shape in all_tables:
            row0 = ' '.join(c.text_frame.text for c in shape.table.rows[0].cells)
            if header_text in row0:
                candidates.append(shape)

    if primary_label is None:
        return candidates[0] if candidates else None

    # Locate the label shape
    lbl = next(
        (sh for sh in slide.shapes
         if sh.has_text_frame and not sh.has_table and primary_label in sh.text_frame.text),
        None
    )
    if lbl is None:
        return candidates[0] if candidates else None

    lbl_bottom = lbl.top + lbl.height
    pool = candidates if candidates else all_tables   # fall back to all tables
    if not pool:
        return None
    return min(pool, key=lambda s: abs(s.top - lbl_bottom))


def _action_is(action_str, *keywords):
    """Check if action string contains ALL given keywords (case-insensitive).
    Handles variant forms like 'Sell completely' vs 'Complete Sell' vs 'Sell partially' vs 'Partial Sell'."""
    a = (action_str or '').strip().lower()
    return all(kw in a for kw in keywords)


def _populate_single_action(prs, slide_idx, label_text, schemes, name_key, value_key,
                             strip_label=None):
    """
    Populate one action type on its own dedicated template slide.
    - If no schemes: delete the slide.
    - If schemes > MAX_ROWS_PER_SLIDE: clone the slide for each overflow page.
    - strip_label: if the template slide has an embedded secondary section (e.g. Fresh Buy
      slide still has a Buy More section below), pass its label text to remove it.
    """
    if not schemes:
        delete_slide(prs, slide_idx)
        print(f"  Deleted empty '{label_text}' slide")
        return

    pages = [schemes[i:i + MAX_ROWS_PER_SLIDE]
             for i in range(0, len(schemes), MAX_ROWS_PER_SLIDE)]

    # Strip any embedded secondary section from the original template slide first,
    # so clones are clean from the start.
    original = prs.slides[slide_idx]
    if strip_label:
        _remove_section_from_slide(original, strip_label, None)

    # Fill first page on original slide
    tbl = _find_table_near_label(original, label_text, None)
    if tbl:
        _fill_action_table(tbl.table, pages[0], name_key, value_key)
    else:
        print(f"  WARNING: table for '{label_text}' not found on slide {slide_idx}")

    # Overflow: clone the (already-stripped) original slide for each additional page
    current_idx = slide_idx
    for page in pages[1:]:
        current_idx = duplicate_slide_after(prs, current_idx)
        dup = prs.slides[current_idx]
        tbl2 = _find_table_near_label(dup, label_text, None)
        if tbl2:
            _fill_action_table(tbl2.table, page, name_key, value_key)


def _dedup_by_fund(rows, value_keys):
    """
    Merge rows that share the same FUND_NAME and Action (e.g. same fund held in two folios).
    Sums numeric value columns; keeps first row's other fields.
    Used for action slides — same fund + same action across multiple folios → one row.
    """
    seen = {}
    order = []
    for row in rows:
        key = (row.get('FUND_NAME', ''), row.get('Action', ''))
        if key not in seen:
            seen[key] = dict(row)
            order.append(key)
        else:
            for vk in value_keys:
                for col in [vk, vk + ' Amount', vk.replace(' Amount', '')]:
                    if col in row and row[col]:
                        seen[key][col] = (seen[key].get(col) or 0) + (row[col] or 0)
    return [seen[k] for k in order]


def _dedup_schemes_for_slide(rows, value_keys):
    """
    Merge rows that share the same FUND_NAME for corpus/SIP scheme slides.
    A fund may have multiple rows because of:
      (a) multiple folios (same action), OR
      (b) different actions (e.g. Retain existing + Buy More / Fresh Buy).
    On scheme slides the fund should appear exactly ONCE with summed allocation values.
    Sums all columns listed in value_keys; keeps first row's other fields (ISIN, RISK_GROUP, etc.).
    """
    seen = {}
    order = []
    for row in rows:
        key = row.get('FUND_NAME', '')
        if key not in seen:
            seen[key] = dict(row)
            order.append(key)
        else:
            for vk in value_keys:
                if vk in row and row[vk]:
                    seen[key][vk] = (seen[key].get(vk) or 0) + (row[vk] or 0)
    return [seen[k] for k in order]


def populate_action_slides(prs, section4):
    """
    Find each of the 6 action-type template slides by their label text and populate
    them independently.  ALL data comes from section4 (ideal portfolio).

    Template slide layout (0-indexed):
      23 - Complete Sell   (standalone)
      24 - Partial Sell    (standalone)
      26 - Complete Retain (standalone)
      27 - Partial Retain  (standalone)
      29 - Fresh Buy       (standalone)
      30 - Buy More        (standalone)

    Also:
    - Deletes section-header slides (02/03/04) when all actions in that section are empty.
    - Removes placeholder "link text" shapes from the transition plan slide (slide 10) for
      empty sections (shapes: "to be sold", "to be retained", "ideal portfolio").
    - Sorts all scheme lists alphabetically by fund name.
    """
    s4 = [r for r in section4 if not r.get('__grand_total__')]
    buy_val_key = _detect_s4_schema(section4)['buy_value']

    print(f"  Action slides: {len(s4)} schemes from section4")

    # All 6 action types — every one reads from section4
    action_defs = [
        ('Complete Sell',
         lambda r: _action_is(r.get('Action'), 'sell') and _action_is(r.get('Action'), 'complet'),
         'FUND_NAME', 'Redemption Value Amount', None),

        ('Partial Sell',
         lambda r: _action_is(r.get('Action'), 'sell') and _action_is(r.get('Action'), 'partial'),
         'FUND_NAME', 'Redemption Value Amount', None),

        ('Complete Retain',
         lambda r: _action_is(r.get('Action'), 'retain') and _action_is(r.get('Action'), 'complet'),
         'FUND_NAME', 'Retained Value Amount', None),

        ('Partial Retain',
         lambda r: _action_is(r.get('Action'), 'retain') and _action_is(r.get('Action'), 'partial'),
         'FUND_NAME', 'Retained Value Amount', None),

        ('Fresh Buy',
         lambda r: _action_is(r.get('Action'), 'fresh', 'buy'),
         'FUND_NAME', buy_val_key, None),

        ('Buy More',
         lambda r: (_action_is(r.get('Action'), 'buy', 'more')
                    and not _action_is(r.get('Action'), 'fresh')
                    and not _action_is(r.get('Action'), 'retain')),
         'FUND_NAME', buy_val_key, None),
    ]

    # Pre-compute scheme lists (before any slides are deleted, indices are stable for find_by_text)
    action_schemes = {}
    for label, filt, name_key, val_key, strip_lbl in action_defs:
        schemes = [r for r in s4 if filt(r)]
        schemes = _dedup_by_fund(schemes, [val_key])
        # Sort alphabetically by fund name
        schemes = sorted(schemes, key=lambda r: (r.get('FUND_NAME', '') or '').lower())
        action_schemes[label] = schemes

    # Populate (or delete) each action slide
    for label, filt, name_key, val_key, strip_lbl in action_defs:
        idx = find_slide_by_text(prs, label)
        if idx is None:
            print(f"  WARNING: slide for '{label}' not found in template")
            continue
        schemes = action_schemes[label]
        print(f"  {label}: {len(schemes)} scheme(s)")
        _populate_single_action(prs, idx, label, schemes, name_key, val_key, strip_lbl)

    # Determine which top-level sections are empty
    sell_empty   = not action_schemes['Complete Sell'] and not action_schemes['Partial Sell']
    retain_empty = not action_schemes['Complete Retain'] and not action_schemes['Partial Retain']
    buy_empty    = not action_schemes['Fresh Buy'] and not action_schemes['Buy More']

    # Delete section header slides for empty sections (find by text; indices may have shifted)
    section_headers = [
        (sell_empty,   'List of schemes from corpus to be sold'),
        (retain_empty, 'List of schemes from corpus to be retained'),
        (buy_empty,    'List of new schemes to be bought'),
    ]
    for is_empty, header_text in section_headers:
        if is_empty:
            hi = find_slide_by_text(prs, header_text)
            if hi is not None:
                delete_slide(prs, hi)
                print(f"  Deleted section header slide: '{header_text}'")

    # Remove placeholder link-text shapes from the transition plan slide
    transition_idx = find_slide_by_text(prs, 'transition plan for')
    if transition_idx is not None:
        t_slide = prs.slides[transition_idx]
        spTree = t_slide.shapes._spTree
        # Map: keyword fragment → whether to remove
        link_text_rules = [
            ('to be sold',      sell_empty),
            ('to be retained',  retain_empty),
            ('ideal portfolio', buy_empty),
        ]
        for keyword, should_remove in link_text_rules:
            if should_remove:
                for shape in list(t_slide.shapes):
                    if shape.has_text_frame and not shape.has_table:
                        if keyword in shape.text_frame.text.lower():
                            spTree.remove(shape._element)
                            print(f"  Removed link text shape containing '{keyword}' from transition slide")


# ---------------------------------------------------------------------------
# generate_deck
# ---------------------------------------------------------------------------

def generate_deck(excel_data, client_name, ref_data=None, rr_category=None):
    """
    Generate M3 transition deck.

    Args:
        excel_data: dict from read_excel() with section1-4
        client_name: Client display name
        ref_data: Reference lookup dict (loaded from Sheets if None)
        rr_category: Rolling return category dict (loaded from Sheets if None)

    Returns:
        (pptx_bytes: BytesIO, filename: str)
    """
    print(f"\n{'='*60}")
    print(f"  Generating M3 deck for: {client_name}")
    print(f"{'='*60}\n")

    # Load reference data from Google Sheets if not provided
    if ref_data is None or rr_category is None:
        print("[1/6] Loading reference data from Google Sheets...")
        ref_data, rr_category = load_reference_data()

    # Download template from Drive
    print("[2/6] Downloading template from Drive...")
    template_path = _get_template_path()
    prs = Presentation(template_path)

    # Slide 1 (idx 0) — cover
    print("[3/6] Populating slides...")
    populate_slide1(prs.slides[0], client_name)

    # Slide 2 (idx 1) — welcome (also has first name on slide idx 2)
    populate_slide2(prs.slides[1], client_name)
    # Slide 3 (idx 2) — "A quick recap, Hari" also uses first name
    first_name = client_name.split()[0]
    replace_text_preserving_format(prs.slides[2], "Hari", first_name)

    # Slide 4 (idx 3) — SIP summary
    slide4_idx = find_slide_by_text(prs, "ideal SIP strategy")
    _s4_schema = _detect_s4_schema(excel_data['section4'])
    _sip_key = _s4_schema['sip_amount']
    has_sip = any(
        (_get_val(r, _sip_key, 'SIP Amount', 'SIP Allocation Amount') or 0) > 0
        for r in excel_data['section4'] if not r.get('__grand_total__')
    )

    if has_sip and slide4_idx is not None:
        populate_slide4(prs, slide4_idx, excel_data['section4'])
    elif not has_sip and slide4_idx is not None:
        delete_slide(prs, slide4_idx)
        print("  No SIP — deleted slide 4")

    # Slide 10 (transition plan)
    print("[4/6] Transition plan...")
    slide10_idx = find_slide_by_text(prs, "transition plan for")
    if slide10_idx is not None:
        populate_slide10(prs, slide10_idx,
                         excel_data['section1'],
                         excel_data['section3'],
                         excel_data['section4'])

    # SIP scheme slides (idx 4-8)
    print("[5/6] Scheme slides...")
    _ref = ref_data or {}
    _rr = rr_category or {}
    if has_sip:
        populate_sip_scheme_slides(prs, excel_data['section4'], _ref, _rr)
    else:
        # Delete all SIP scheme slides
        for marker in ['SIP strategy - at a scheme level']:
            sep = find_slide_by_text(prs, "List of schemes for ideal portfolio")
            for idx in sorted(find_all_slides_by_text(prs, marker), reverse=True):
                if sep is None or idx < sep:
                    delete_slide(prs, idx)
                    print(f"  Deleted SIP scheme slide {idx}")

    # Corpus scheme slides (idx 14-21)
    populate_corpus_scheme_slides(prs, excel_data['section4'], _ref, _rr)

    # Action slides
    print("[6/6] Action slides...")
    populate_action_slides(prs, excel_data['section4'])

    # Save to BytesIO
    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)

    safe = re.sub(r'[^\w\s-]', '', client_name).strip().replace(' ', '_')
    filename = f'{safe}_transition_deck.pptx'

    print(f"\n{'='*60}")
    print(f"  DONE  ->  {filename}")
    print(f"{'='*60}\n")
    return buf, filename
